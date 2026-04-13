"""
概念图模板渲染器
根据 DiagramSpec 和 ContentDiagramSpec 渲染专业概念图

支持的模板：
- process_flow: 横/纵向流程图
- architecture: 层级架构图
- relationship: 因果/关系图
- framework variants: 2x2矩阵, SWOT, 金字塔, 漏斗, Venn
"""

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from models.slide_spec import DiagramSpec, DiagramNode, VisualTheme, Rect


class DiagramRenderer:
    """参数化概念图渲染器"""

    def render(self, slide, diag_spec: DiagramSpec, rect: Rect, theme: VisualTheme):
        """根据 diagram_type 分发到对应模板"""
        dtype = diag_spec.diagram_type

        if dtype == "process_flow":
            self._render_process_flow(slide, diag_spec, rect, theme)
        elif dtype == "relationship":
            self._render_relationship(slide, diag_spec, rect, theme)
        elif dtype == "architecture":
            self._render_architecture(slide, diag_spec, rect, theme)
        elif dtype == "framework":
            self._render_framework(slide, diag_spec, rect, theme)
        else:
            self._render_fallback_grid(slide, diag_spec, rect, theme)

    # ── 流程图 ──

    def _render_process_flow(self, slide, spec: DiagramSpec, rect: Rect, theme: VisualTheme):
        """横向/纵向流程图"""
        nodes = spec.nodes
        edges = spec.edges
        n = len(nodes)
        if n == 0:
            return

        is_horizontal = spec.layout_direction == "LR"
        primary = _color(theme.colors.get("primary", "#003D6E"))
        accent = _color(theme.colors.get("accent", "#FF6B35"))
        white = RGBColor(255, 255, 255)
        font = theme.fonts.get("body", "Calibri")

        if is_horizontal:
            node_w = min(int(rect.width / (n * 1.5)), 2286000)
            node_h = min(int(rect.height * 0.5), 640080)
            gap = (rect.width - n * node_w) // max(n + 1, 1)
            positions = {}
            for i, node in enumerate(nodes):
                x = rect.left + gap + i * (node_w + gap)
                y = rect.top + (rect.height - node_h) // 2
                positions[node.node_id] = (x, y)
                _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, node_w, node_h,
                           primary, node.label, white, font, 11)
            # 箭头连接
            for i in range(n - 1):
                x1 = positions[nodes[i].node_id][0] + node_w
                y1 = positions[nodes[i].node_id][1] + node_h // 2
                x2 = positions[nodes[i + 1].node_id][0]
                y2 = positions[nodes[i + 1].node_id][1] + node_h // 2
                _add_arrow(slide, x1, y1, x2, y2, accent)
                # 边标签
                edge_label = _find_edge_label(edges, nodes[i].node_id, nodes[i + 1].node_id)
                if edge_label:
                    _add_label(slide, (x1 + x2) // 2, y1 - 274320, edge_label, accent, font, 9)
        else:
            node_w = min(int(rect.width * 0.7), 3200400)
            node_h = min(int(rect.height / (n * 1.5)), 548640)
            gap = max((rect.height - n * node_h) // max(n + 1, 1), 91440)
            for i, node in enumerate(nodes):
                x = rect.left + (rect.width - node_w) // 2
                y = rect.top + gap + i * (node_h + gap)
                _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, node_w, node_h,
                           primary if i % 2 == 0 else accent,
                           node.label, white, font, 11)
                if i < n - 1:
                    arrow_x = x + node_w // 2
                    _add_arrow(slide, arrow_x, y + node_h, arrow_x,
                               y + node_h + gap, RGBColor(0x99, 0x99, 0x99))

    # ── 关系/因果图 ──

    def _render_relationship(self, slide, spec: DiagramSpec, rect: Rect, theme: VisualTheme):
        """因果/关系网络图"""
        nodes = spec.nodes
        edges = spec.edges
        n = len(nodes)
        if n == 0:
            return

        primary = _color(theme.colors.get("primary", "#003D6E"))
        accent = _color(theme.colors.get("accent", "#FF6B35"))
        white = RGBColor(255, 255, 255)
        font = theme.fonts.get("body", "Calibri")

        # 圆形布局
        node_w = min(int(rect.width * 0.22), 1828800)
        node_h = min(int(rect.height * 0.18), 548640)
        cx = rect.left + rect.width // 2
        cy = rect.top + rect.height // 2
        radius_x = rect.width // 2 - node_w // 2 - 182880
        radius_y = rect.height // 2 - node_h // 2 - 182880

        positions = {}
        for i, node in enumerate(nodes):
            import math
            angle = 2 * math.pi * i / n - math.pi / 2
            x = int(cx + radius_x * math.cos(angle)) - node_w // 2
            y = int(cy + radius_y * math.sin(angle)) - node_h // 2
            positions[node.node_id] = (x, y, node_w, node_h)

            # cause 用 primary, effect 用 accent
            color = primary
            _add_shape(slide, MSO_SHAPE.OVAL, x, y, node_w, node_h,
                       color, node.label, white, font, 10)

        # 连接边
        for edge in edges:
            if edge.from_id in positions and edge.to_id in positions:
                fx, fy, fw, fh = positions[edge.from_id]
                tx, ty, tw, th = positions[edge.to_id]
                _add_arrow(slide,
                           fx + fw // 2, fy + fh // 2,
                           tx + tw // 2, ty + th // 2,
                           RGBColor(0x99, 0x99, 0x99))
                if edge.label:
                    mx = (fx + fw // 2 + tx + tw // 2) // 2
                    my = (fy + fh // 2 + ty + th // 2) // 2
                    _add_label(slide, mx, my - 182880, edge.label, accent, font, 8)

    # ── 架构图 ──

    def _render_architecture(self, slide, spec: DiagramSpec, rect: Rect, theme: VisualTheme):
        """层级架构图"""
        nodes = spec.nodes
        if not nodes:
            return

        primary = _color(theme.colors.get("primary", "#003D6E"))
        accent = _color(theme.colors.get("accent", "#FF6B35"))
        white = RGBColor(255, 255, 255)
        font = theme.fonts.get("body", "Calibri")

        # 按 group 分层
        layers = {}
        for node in nodes:
            group = node.group or "default"
            layers.setdefault(group, []).append(node)

        layer_names = list(layers.keys())
        num_layers = len(layer_names)
        if num_layers == 0:
            return

        layer_h = rect.height // num_layers
        layer_colors = [
            RGBColor(0x1A, 0x47, 0x8A),  # dark navy
            RGBColor(0x00, 0x5B, 0x96),  # medium blue
            RGBColor(0x00, 0x7B, 0xC0),  # blue
            RGBColor(0x48, 0xA9, 0xE6),  # light blue
            RGBColor(0x7E, 0xC8, 0xE3),  # very light
        ]

        for li, layer_name in enumerate(layer_names):
            layer_nodes = layers[layer_name]
            y = rect.top + li * layer_h
            color = layer_colors[li % len(layer_colors)]

            # 层背景
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(rect.left), Emu(y),
                Emu(rect.width), Emu(layer_h - 45720)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = color
            bg.line.fill.background()

            # 节点
            n_items = len(layer_nodes)
            item_w = min(rect.width // max(n_items, 1) - 91440, 2743200)
            item_h = min(layer_h // 2, 457200)
            total_w = n_items * item_w + (n_items - 1) * 91440
            start_x = rect.left + (rect.width - total_w) // 2

            for ni, node in enumerate(layer_nodes):
                nx = start_x + ni * (item_w + 91440)
                ny = y + (layer_h - item_h) // 2
                _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                           nx, ny, item_w, item_h,
                           white, node.label, color, font, 10)

    # ── 框架图（2x2, SWOT, 金字塔, 漏斗）──

    def _render_framework(self, slide, spec: DiagramSpec, rect: Rect, theme: VisualTheme):
        """框架图：根据节点分组渲染"""
        nodes = spec.nodes
        if not nodes:
            return

        primary = _color(theme.colors.get("primary", "#003D6E"))
        accent = _color(theme.colors.get("accent", "#FF6B35"))
        white = RGBColor(255, 255, 255)
        font = theme.fonts.get("body", "Calibri")

        # 按 group 分组
        groups = {}
        for node in nodes:
            g = node.group or "default"
            groups.setdefault(g, []).append(node)

        group_keys = list(groups.keys())

        # 2x2 矩阵 / SWOT
        if len(group_keys) == 4 and set(group_keys) & {"top_left", "top_right", "bottom_left", "bottom_right"}:
            self._render_2x2(slide, groups, rect, theme, group_keys)

        # 金字塔（单列多层）
        elif len(group_keys) == 1 and len(groups[group_keys[0]]) >= 3:
            self._render_pyramid(slide, groups[group_keys[0]], rect, theme)

        # 漏斗
        elif len(group_keys) == 1:
            self._render_funnel(slide, groups[group_keys[0]], rect, theme)

        else:
            self._render_fallback_grid(slide, spec, rect, theme)

    def _render_2x2(self, slide, groups, rect: Rect, theme: VisualTheme, group_keys):
        """2x2 矩阵 / SWOT"""
        primary = _color(theme.colors.get("primary", "#003D6E"))
        white = RGBColor(255, 255, 255)
        font = theme.fonts.get("body", "Calibri")

        quadrant_colors = {
            "top_left": RGBColor(0x00, 0x3D, 0x6E),
            "top_right": RGBColor(0x00, 0x7B, 0xC0),
            "bottom_left": RGBColor(0xFF, 0x6B, 0x35),
            "bottom_right": RGBColor(0x48, 0xA9, 0xE6),
        }

        gap = 91440  # 0.1 inch
        half_w = (rect.width - gap) // 2
        half_h = (rect.height - gap) // 2

        positions = {
            "top_left": (rect.left, rect.top),
            "top_right": (rect.left + half_w + gap, rect.top),
            "bottom_left": (rect.left, rect.top + half_h + gap),
            "bottom_right": (rect.left + half_w + gap, rect.top + half_h + gap),
        }

        for key in group_keys:
            x, y = positions.get(key, (rect.left, rect.top))
            color = quadrant_colors.get(key, primary)

            # 象限背景
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(x), Emu(y), Emu(half_w), Emu(half_h)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = color
            bg.line.fill.background()

            # 标题 + 子项
            items = groups[key]
            if items:
                # 标题
                _add_label(slide, x + 91440, y + 45720, items[0].label,
                           white, font, 13, bold=True)
                # 子项
                for j, item in enumerate(items[1:6]):
                    _add_label(slide, x + 91440, y + 365760 + j * 228600,
                               f"• {item.label}", white, font, 10)

    def _render_pyramid(self, slide, nodes, rect: Rect, theme: VisualTheme):
        """金字塔图"""
        n = len(nodes)
        if n == 0:
            return

        primary = _color(theme.colors.get("primary", "#003D6E"))
        accent = _color(theme.colors.get("accent", "#FF6B35"))
        white = RGBColor(255, 255, 255)
        font = theme.fonts.get("body", "Calibri")

        colors = [primary, RGBColor(0x00, 0x5B, 0x96), accent,
                  RGBColor(0x48, 0xA9, 0xE6), RGBColor(0x7E, 0xC8, 0xE3)]

        layer_h = rect.height // n
        cx = rect.left + rect.width // 2

        for i, node in enumerate(nodes):
            # 越往上越窄
            width_ratio = (i + 1) / n
            w = int(rect.width * width_ratio * 0.85)
            x = cx - w // 2
            y = rect.top + i * layer_h

            _add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, layer_h - 45720,
                       colors[i % len(colors)], node.label, white, font, 11)

    def _render_funnel(self, slide, nodes, rect: Rect, theme: VisualTheme):
        """漏斗图"""
        n = len(nodes)
        if n == 0:
            return

        primary = _color(theme.colors.get("primary", "#003D6E"))
        white = RGBColor(255, 255, 255)
        font = theme.fonts.get("body", "Calibri")

        colors = [
            RGBColor(0x00, 0x3D, 0x6E), RGBColor(0x00, 0x5B, 0x96),
            RGBColor(0x00, 0x7B, 0xC0), RGBColor(0x48, 0xA9, 0xE6),
            RGBColor(0x7E, 0xC8, 0xE3),
        ]

        layer_h = rect.height // n
        cx = rect.left + rect.width // 2

        for i, node in enumerate(nodes):
            width_ratio = 1.0 - (i / n) * 0.6  # 递减
            w = int(rect.width * width_ratio)
            x = cx - w // 2
            y = rect.top + i * layer_h
            color = colors[i % len(colors)]

            _add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, layer_h - 45720,
                       color, node.label, white, font, 11)

    # ── 兜底网格布局 ──

    def _render_fallback_grid(self, slide, spec: DiagramSpec, rect: Rect, theme: VisualTheme):
        """简单网格布局兜底"""
        nodes = spec.nodes
        n = len(nodes)
        if n == 0:
            return

        primary = _color(theme.colors.get("primary", "#003D6E"))
        white = RGBColor(255, 255, 255)
        font = theme.fonts.get("body", "Calibri")

        cols = min(n, 4)
        rows = (n + cols - 1) // cols
        node_w = min(rect.width // cols - 91440, 2286000)
        node_h = min(rect.height // rows - 91440, 640080)

        positions = {}
        for i, node in enumerate(nodes):
            r = i // cols
            c = i % cols
            x = rect.left + c * (node_w + 91440) + 45720
            y = rect.top + r * (node_h + 91440) + 45720
            positions[node.node_id] = (x, y, node_w, node_h)
            _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, node_w, node_h,
                       primary, node.label, white, font, 11)

        for edge in spec.edges:
            if edge.from_id in positions and edge.to_id in positions:
                fx, fy, fw, fh = positions[edge.from_id]
                tx, ty, tw, th = positions[edge.to_id]
                _add_arrow(slide, fx + fw // 2, fy + fh,
                           tx + tw // 2, ty, RGBColor(0x99, 0x99, 0x99))


# ============================================================
# 辅助函数
# ============================================================

def _color(hex_str: str) -> RGBColor:
    """解析颜色"""
    if not hex_str:
        return RGBColor(0, 0, 0)
    h = hex_str.lstrip("#")
    if len(h) == 6:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    return RGBColor(0, 0, 0)


def _add_shape(slide, shape_type, x, y, w, h, fill_color, text, text_color, font_name, font_size):
    """添加一个形状 + 居中文本"""
    shape = slide.shapes.add_shape(shape_type, Emu(x), Emu(y), Emu(w), Emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = font_name
    p.alignment = PP_ALIGN.CENTER


def _add_arrow(slide, x1, y1, x2, y2, color):
    """添加连接线"""
    connector = slide.shapes.add_connector(
        1,  # straight
        Emu(x1), Emu(y1), Emu(x2), Emu(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)


def _add_label(slide, x, y, text, color, font_name, font_size, bold=False):
    """添加文字标签"""
    box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(2743200), Emu(274320))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name


def _find_edge_label(edges, from_id, to_id):
    """查找边的标签"""
    for e in edges:
        if e.from_id == from_id and e.to_id == to_id:
            return e.label
    return None
