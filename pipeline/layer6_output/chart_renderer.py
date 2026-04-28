"""
ChartRenderer — 分层图表渲染引擎

将图表类型分为两层：
- NATIVE：python-pptx 原生图表（矢量、可编辑）→ COLUMN/BAR/LINE/PIE/AREA/SCATTER
- PLOTLY_PNG：Plotly 渲染为高清 PNG 后嵌入 → WATERFALL/COMBO/SANKEY/TREEMAP/等

设计原则：
- 简单图表用 native 保证可编辑性和矢量缩放
- 复杂图表用 Plotly 保证视觉质量（python-pptx 不原生支持这些类型）
"""

import tempfile
from enum import Enum
from pathlib import Path
from typing import Optional

from models.slide_spec import ChartSpec, ChartType, Rect


class RenderMethod(Enum):
    NATIVE = "native"          # python-pptx 原生图表对象
    PLOTLY_PNG = "plotly_png"  # Plotly → PNG → 嵌入图片


# 复杂度分类：哪些类型走 Plotly
_PLOTLY_TYPES = {
    ChartType.WATERFALL,
    ChartType.COMBO,
}

# 未来可扩展的 Plotly 类型（ChartType 尚未定义，预留）
# SANKEY, TREEMAP, SUNBURST, FUNNEL, HEATMAP


class ChartRenderer:
    """
    分层图表渲染引擎。

    根据图表复杂度自动选择渲染方式：
    - 简单图表（柱/条/折线/饼/面积/散点）→ python-pptx 原生
    - 复杂图表（瀑布/组合/桑基/树状等）  → Plotly PNG
    """

    def __init__(self, dpi: int = 200):
        self.dpi = dpi
        self._plotly_available: Optional[bool] = None

    def classify(self, chart_type: ChartType) -> RenderMethod:
        """判断图表应使用哪种渲染方式"""
        if chart_type in _PLOTLY_TYPES:
            # Plotly 不可用时降级到 native
            if not self._check_plotly():
                return RenderMethod.NATIVE
            return RenderMethod.PLOTLY_PNG
        return RenderMethod.NATIVE

    def render(
        self,
        slide,
        chart_spec: ChartSpec,
        rect: Rect,
        theme,
        native_renderer=None,
    ) -> bool:
        """
        渲染图表到 slide。

        Args:
            slide: python-pptx slide 对象
            chart_spec: 图表规格
            rect: 放置区域 (EMU)
            theme: VisualTheme
            native_renderer: 原生渲染回调 (slide, chart_spec, rect, theme) -> None
                             用于 NATIVE 类型的实际渲染（由 ppt_builder 提供）

        Returns:
            True if 渲染成功
        """
        if not chart_spec.categories or not chart_spec.series:
            return False

        method = self.classify(chart_spec.chart_type)

        if method == RenderMethod.PLOTLY_PNG:
            return self._render_plotly(slide, chart_spec, rect, theme)
        else:
            # NATIVE：委托给 ppt_builder 的原有逻辑
            if native_renderer:
                try:
                    native_renderer(slide, chart_spec, rect, theme)
                    return True
                except Exception as e:
                    print(f"[ChartRenderer] native渲染失败: {e}")
                    return False
            return False

    # ================================================================
    # Plotly 渲染路径
    # ================================================================

    def _render_plotly(
        self,
        slide,
        chart_spec: ChartSpec,
        rect: Rect,
        theme,
    ) -> bool:
        """使用 Plotly 渲染图表为 PNG 并嵌入到 slide"""
        try:
            import plotly.graph_objects as go
        except ImportError:
            print("[ChartRenderer] plotly未安装，降级到native渲染")
            return False

        if theme is None:
            theme = self._default_theme()

        chart_type = chart_spec.chart_type

        try:
            if chart_type == ChartType.WATERFALL:
                fig = self._build_waterfall_fig(chart_spec, theme)
            elif chart_type == ChartType.COMBO:
                fig = self._build_combo_fig(chart_spec, theme)
            else:
                # 兜底：generic bar chart
                fig = self._build_generic_fig(chart_spec, theme)

            if fig is None:
                return False

            # 应用通用样式
            self._apply_fig_style(fig, chart_spec, theme)

            # 导出 PNG → 嵌入
            img_path = self._export_png(fig)
            if img_path is None:
                return False

            self._embed_image(slide, img_path, rect)

            # so_what 文字（如果有的话）
            if chart_spec.so_what:
                from pptx.util import Emu, Pt
                from pptx.enum.text import PP_ALIGN
                so_what_area = Rect(
                    left=rect.left,
                    top=rect.top + rect.height + 91440,
                    width=rect.width,
                    height=365760,
                )
                txBox = slide.shapes.add_textbox(
                    Emu(so_what_area.left), Emu(so_what_area.top),
                    Emu(so_what_area.width), Emu(so_what_area.height),
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = chart_spec.so_what
                p.font.size = Pt(10)
                p.font.italic = True
                text_color = self._parse_color(theme.colors.get("text_secondary", "#666666"))
                p.font.color.rgb = text_color
                p.font.name = theme.fonts.get("body", "Microsoft YaHei")
                p.alignment = PP_ALIGN.CENTER

            return True

        except Exception as e:
            print(f"[ChartRenderer] Plotly渲染失败: {e}")
            return False

    def _build_waterfall_fig(self, spec: ChartSpec, theme) -> "go.Figure":
        """构建瀑布图 Plotly figure"""
        import plotly.graph_objects as go

        values = spec.series[0].values
        categories = spec.categories

        # 计算累计值用于 text 显示
        running = 0
        text_vals = []
        for v in values:
            running += v
            text_vals.append(f"+{v}" if v >= 0 else str(v))

        # 确定 measure: 最终一个点设为 total
        measure = ["relative"] * len(values)

        fig = go.Figure(go.Waterfall(
            x=categories,
            y=[float(v) for v in values],
            measure=measure,
            text=text_vals,
            textposition="outside",
            increasing={"marker": {"color": theme.colors.get("primary", "#003D6E")}},
            decreasing={"marker": {"color": theme.colors.get("accent", "#FF6B35")}},
            connector={"line": {"color": "#D0D0D0", "width": 1}},
        ))

        return fig

    def _build_combo_fig(self, spec: ChartSpec, theme) -> "go.Figure":
        """构建组合图 Plotly figure"""
        import plotly.graph_objects as go

        palette = theme.colors.get("chart_palette", [])

        fig = go.Figure()

        for i, series in enumerate(spec.series):
            color = series.color or (palette[i] if i < len(palette) else palette[i % len(palette)])
            # 第一个系列用柱状，后续用折线
            if i == 0:
                fig.add_trace(go.Bar(
                    name=series.name,
                    x=spec.categories,
                    y=[float(v) for v in series.values],
                    marker_color=color,
                ))
            else:
                fig.add_trace(go.Scatter(
                    name=series.name,
                    x=spec.categories,
                    y=[float(v) for v in series.values],
                    mode="lines+markers",
                    line=dict(color=color, width=2),
                    marker=dict(size=6),
                    yaxis="y2" if i >= 2 else "y",
                ))

        # 双轴设置（如果有3+系列）
        if len(spec.series) >= 3:
            fig.update_layout(yaxis2=dict(
                overlaying="y",
                side="right",
                showgrid=False,
            ))

        return fig

    def _build_generic_fig(self, spec: ChartSpec, theme) -> "go.Figure":
        """通用柱状图 fallback"""
        import plotly.graph_objects as go

        palette = theme.colors.get("chart_palette", [])

        fig = go.Figure()
        for i, series in enumerate(spec.series):
            color = series.color or (palette[i] if i < len(palette) else palette[i % len(palette)])
            fig.add_trace(go.Bar(
                name=series.name,
                x=spec.categories,
                y=[float(v) for v in series.values],
                marker_color=color,
            ))
        return fig

    def _apply_fig_style(self, fig, spec: ChartSpec, theme):
        """应用通用 Plotly 样式（白底、无网格线、字体）"""
        font_family = theme.fonts.get("body", "Microsoft YaHei")
        text_color = theme.colors.get("text_primary", "#333333")

        # Hide legend when single series
        show_legend = spec.show_legend
        if show_legend and len(spec.series) <= 1:
            show_legend = False

        fig.update_layout(
            plot_bgcolor="white",
            paper_bgcolor="white",
            font=dict(family=font_family, size=12, color=text_color),
            showlegend=show_legend,
            margin=dict(l=60, r=30, t=50 if spec.title else 20, b=60),
            xaxis=dict(
                showgrid=False,
                linecolor="#D0D0D0",
                linewidth=1,
            ),
            yaxis=dict(
                showgrid=False,
                linecolor="#D0D0D0",
                linewidth=1,
            ),
        )

        # Dedup axis_title vs series_name
        if spec.title:
            for axis_key in ("xaxis", "yaxis"):
                axis_cfg = fig.layout.get(axis_key, {})
                if hasattr(axis_cfg, "title") and axis_cfg.title:
                    axis_title = axis_cfg.title.text if hasattr(axis_cfg.title, "text") else str(axis_cfg.title)
                    for s in spec.series:
                        if s.name and s.name == axis_title:
                            fig.layout[axis_key].title = None
                            break

            fig.update_layout(title=dict(
                text=spec.title,
                font=dict(size=16, color=text_color),
                x=0.5,
                xanchor="center",
            ))

    def _export_png(self, fig) -> Optional[str]:
        """导出 Plotly figure 为临时 PNG 文件"""
        try:
            tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            tmp_path = tmp.name
            tmp.close()

            fig.write_image(tmp_path, scale=self.dpi / 72, engine="kaleido")
            return tmp_path
        except Exception as e:
            print(f"[ChartRenderer] PNG导出失败: {e}")
            return None

    def _embed_image(self, slide, img_path: str, rect: Rect):
        """将 PNG 图片嵌入到 slide"""
        from pptx.util import Emu
        slide.shapes.add_picture(
            img_path,
            Emu(rect.left), Emu(rect.top),
            Emu(rect.width), Emu(rect.height),
        )

    # ================================================================
    # 工具方法
    # ================================================================

    def _check_plotly(self) -> bool:
        """检查 Plotly + Kaleido 是否可用"""
        if self._plotly_available is None:
            try:
                import plotly.graph_objects as go  # noqa: F401
                import kaleido  # noqa: F401
                self._plotly_available = True
            except ImportError:
                self._plotly_available = False
                print("[ChartRenderer] plotly或kaleido未安装，复杂图表将降级为native渲染")
        return self._plotly_available

    def render_into_pptx(
        self,
        pptx_path: str,
        placeholders: list,
        slides_data: list,
        output_path: str,
        theme=None,
    ) -> str:
        """
        Inject native charts into an existing .pptx at placeholder positions.

        Used after html2pptx.js renders text/shapes/images — this adds python-pptx
        native charts at the placeholder locations.

        Args:
            pptx_path: Path to .pptx produced by html2pptx.js
            placeholders: [{"slide_index": int, "items": [{"id": str, "x": float, "y": float, "w": float, "h": float}]}]
            slides_data: Per-slide chart specs: [{"chart_spec": ChartSpec dict, "theme": VisualTheme}]
            output_path: Where to save the final .pptx
            theme: Default VisualTheme if per-slide theme is None

        Returns:
            output_path
        """
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.chart.data import CategoryChartData

        prs = Presentation(pptx_path)

        # Map ChartType to python-pptx XL_CHART_TYPE
        _CHART_TYPE_MAP = {
            "column": 51,    # XL_CHART_TYPE.COLUMN_CLUSTERED
            "bar": 57,       # XL_CHART_TYPE.BAR_CLUSTERED
            "line": 4,       # XL_CHART_TYPE.LINE
            "pie": 5,        # XL_CHART_TYPE.PIE
            "area": 29,      # XL_CHART_TYPE.AREA
            "scatter": -4169, # XL_CHART_TYPE.XY_SCATTER
        }

        for ph_entry in placeholders:
            slide_idx = ph_entry.get("slide_index", 0)
            if slide_idx >= len(prs.slides) or slide_idx >= len(slides_data):
                continue

            slide = prs.slides[slide_idx]
            data = slides_data[slide_idx] if slide_idx < len(slides_data) else {}
            chart_spec_dict = data.get("chart_spec")
            if not chart_spec_dict:
                continue

            # Normalize: ChartSuggestion / ChartSpec / dict → ChartSpec
            if isinstance(chart_spec_dict, dict):
                chart_spec = ChartSpec.from_dict(chart_spec_dict)
            elif hasattr(chart_spec_dict, "chart_type"):
                # Likely a ChartSuggestion — convert to dict then to ChartSpec
                raw = {
                    "chart_type": getattr(chart_spec_dict, "chart_type", "column"),
                    "title": getattr(chart_spec_dict, "title", ""),
                    "categories": getattr(chart_spec_dict, "categories", []),
                    "series": getattr(chart_spec_dict, "series", []),
                    "so_what": getattr(chart_spec_dict, "so_what", ""),
                }
                chart_spec = ChartSpec.from_dict(raw)
            else:
                continue
            slide_theme = data.get("theme") or theme or self._default_theme()

            for ph_item in ph_entry.get("items", []):
                x_in = ph_item.get("x", 1)
                y_in = ph_item.get("y", 1)
                w_in = ph_item.get("w", 4)
                h_in = ph_item.get("h", 3)

                # Try direct native chart first (simpler types)
                chart_type_str = chart_spec.chart_type.value if hasattr(chart_spec.chart_type, "value") else str(chart_spec.chart_type).lower()
                xl_type = _CHART_TYPE_MAP.get(chart_type_str)

                if xl_type is not None and chart_spec.categories and chart_spec.series:
                    try:
                        chart_data = CategoryChartData()
                        chart_data.categories = list(chart_spec.categories)
                        for s in chart_spec.series:
                            chart_data.add_series(s.name, (float(v) for v in s.values))

                        chart_frame = slide.shapes.add_chart(
                            xl_type,
                            Inches(x_in), Inches(y_in),
                            Inches(w_in), Inches(h_in),
                            chart_data,
                        )
                        chart = chart_frame.chart
                        chart.has_legend = bool(chart_spec.show_legend)
                        if chart_spec.title:
                            chart.has_title = True
                            chart.chart_title.text_frame.paragraphs[0].text = chart_spec.title
                        continue
                    except Exception as e:
                        print(f"[ChartRenderer] native chart failed: {e}, trying render()")

                # Fallback to render() method (Plotly PNG etc.)
                rect = Rect(
                    left=int(x_in * 914400),
                    top=int(y_in * 914400),
                    width=int(w_in * 914400),
                    height=int(h_in * 914400),
                )
                self.render(slide, chart_spec, rect, slide_theme)

        prs.save(output_path)
        return output_path

    @staticmethod
    def _default_theme():
        """当 theme为 None 时提供默认颜色/字体，避免 Plotly 崩溃"""
        from types import SimpleNamespace
        return SimpleNamespace(
            colors={
                "primary": "#003D6E",
                "accent": "#FF6B35",
                "text_primary": "#333333",
                "text_secondary": "#666666",
                "chart_palette": ["#003D6E", "#FF6B35", "#70AD47", "#FFC000", "#5B9BD5", "#C00000"],
            },
            fonts={"body": "Microsoft YaHei", "heading": "Microsoft YaHei"},
        )

    @staticmethod
    def _parse_color(hex_color: str):
        """解析十六进制颜色为 RGBColor"""
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        if not hex_color:
            return RGBColor(0x33, 0x33, 0x33)
        hex_color = hex_color.lstrip("#")
        if len(hex_color) == 6:
            return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
        return RGBColor(0x33, 0x33, 0x33)
