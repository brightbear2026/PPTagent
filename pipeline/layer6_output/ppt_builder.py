"""
PPT Builder: 将SlideSpec转换为.pptx文件
使用python-pptx构建专业级PPT，支持图表和视觉主题
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from typing import List
from pathlib import Path

import sys
from pathlib import Path

# 添加项目根目录到Python路径
project_root = Path(__file__).parent.parent.parent
sys.path.insert(0, str(project_root))

from models import (SlideSpec, PresentationSpec, ChartSpec, DiagramSpec,
                    VisualTheme, Rect, ChartType,
                    VisualBlock, VisualBlockType, VisualBlockItem)
from pipeline.layer6_output.layout_engine import LayoutEngine
from pipeline.layer6_output import chrome
from pipeline.layer5_chart.diagram_renderer import DiagramRenderer


# ChartType → python-pptx XL_CHART_TYPE 映射
CHART_TYPE_MAP = {
    ChartType.COLUMN: XL_CHART_TYPE.COLUMN_CLUSTERED,
    ChartType.BAR: XL_CHART_TYPE.BAR_CLUSTERED,
    ChartType.LINE: XL_CHART_TYPE.LINE_MARKERS,
    ChartType.PIE: XL_CHART_TYPE.PIE,
    ChartType.AREA: XL_CHART_TYPE.AREA,
    ChartType.SCATTER: XL_CHART_TYPE.XY_SCATTER,
    ChartType.WATERFALL: XL_CHART_TYPE.COLUMN_STACKED,  # 用堆叠柱状图模拟
    ChartType.COMBO: XL_CHART_TYPE.COLUMN_CLUSTERED,
}


class PPTBuilder:
    """
    PPT构建器
    负责将SlideSpec对象转换为实际的.pptx文件
    """

    def __init__(self):
        self.prs = Presentation()
        self.layout_engine = LayoutEngine()
        self.diagram_renderer = DiagramRenderer()

        # 设置幻灯片尺寸 (16:9) - 标准16:9是13.333x7.5英寸
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    # 不要求正文内容的 slide_type（页面形态本身就是极简）
    _NO_BODY_TYPES = {"title", "agenda", "section_divider"}

    # 当前正在构建的 PresentationSpec（便于 _build_title_slide 读取 subtitle/author/date）
    _pres_spec: PresentationSpec = None

    def build(self, spec: PresentationSpec) -> str:
        """构建完整的PPT"""
        print(f"🎨 开始构建PPT: {spec.title}")
        print(f"   共 {len(spec.slides)} 页")

        self._pres_spec = spec
        for idx, slide_spec in enumerate(spec.slides, 1):
            if not slide_spec.slide_index:
                slide_spec.slide_index = idx
            self._ensure_non_empty_body(slide_spec, idx)
            print(f"   生成第 {idx} 页: {slide_spec.slide_type.value}")
            self._build_slide(slide_spec)

        # 保存文件
        output_path = f"output/{spec.title}.pptx"
        Path("output").mkdir(exist_ok=True)
        self.prs.save(output_path)

        print(f"✅ PPT生成成功: {output_path}")
        return output_path

    def _ensure_non_empty_body(self, spec: SlideSpec, page_num: int):
        """
        空白正文兜底：对 CONTENT/DATA/COMPARISON/SUMMARY 等需要正文的页面，如果
        text_blocks 和 charts/diagrams 全部为空，则用 takeaway_message 注入最低
        限度占位内容，并打印告警（上游 Layer3 不应给出空页）。
        """
        from models.slide_spec import TextBlock  # 延迟导入避免循环

        st = spec.slide_type.value if hasattr(spec.slide_type, "value") else str(spec.slide_type)
        if st in self._NO_BODY_TYPES:
            return

        has_text = bool(spec.text_blocks)
        has_chart = bool(spec.charts)
        has_diagram = bool(spec.diagrams)
        if has_text or has_chart or has_diagram:
            return

        print(f"   ⚠️  slide {page_num} ({st}): 正文为空，注入 takeaway 兜底")
        if spec.takeaway_message:
            spec.text_blocks = [
                TextBlock(content=spec.takeaway_message, level=0, is_bold=True),
                TextBlock(content="（本页内容由上游阶段填充不足，建议在检查点回退补充）", level=1),
            ]
        else:
            spec.text_blocks = [
                TextBlock(content="本页内容缺失", level=0, is_bold=True),
            ]

    def _build_slide(self, spec: SlideSpec):
        """构建单页：按 slide_type 分派到专用 renderer"""
        blank_layout = self.prs.slide_layouts[6]  # 6 = blank layout
        slide = self.prs.slides.add_slide(blank_layout)

        # 获取视觉主题
        theme = spec.visual_theme or VisualTheme()

        # 计算布局
        layout = self.layout_engine.calculate_layout(spec)

        st = spec.slide_type.value if hasattr(spec.slide_type, "value") else str(spec.slide_type)

        if st == "title":
            self._add_title_decorations(slide, theme)
            self._build_title_slide(slide, spec, layout, theme)
        elif st == "section_divider":
            self._add_section_divider_decorations(slide, theme)
            self._build_section_divider_slide(slide, spec, layout, theme)
        elif st == "agenda":
            self._add_content_decorations(slide, theme)
            self._build_agenda_slide(slide, spec, layout, theme)
        else:
            self._add_content_decorations(slide, theme)
            self._build_content_slide(slide, spec, layout, theme)
            # 页码
            self._add_page_number(slide, spec, theme)

        # 条件渲染：根据 primary_visual 只渲染对应类型
        pv = getattr(spec, 'primary_visual', '')

        # 渲染图表（仅 primary_visual=chart 或未标注时）
        if pv in ('chart', ''):
            seen_chart_ids = set()
            unique_charts = []
            for c in spec.charts:
                cid = getattr(c, "chart_id", None) or id(c)
                if cid in seen_chart_ids:
                    continue
                seen_chart_ids.add(cid)
                unique_charts.append(c)
            for i, chart_spec in enumerate(unique_charts):
                if i < len(layout.chart_areas):
                    self._add_chart(slide, chart_spec,
                                    self._clamp_rect(layout.chart_areas[i]), theme)

        # 渲染架构图（仅 primary_visual=diagram 或未标注时）
        if pv in ('diagram', ''):
            for i, diag_spec in enumerate(spec.diagrams):
                if i < len(layout.diagram_areas):
                    self._add_diagram(slide, diag_spec,
                                      self._clamp_rect(layout.diagram_areas[i]), theme)

        # 渲染原材料图片
        pictures = getattr(spec, "pictures", None) or []
        picture_areas = getattr(layout, "picture_areas", None) or []
        for i, pic_path in enumerate(pictures):
            if i < len(picture_areas):
                self._add_picture(slide, pic_path, picture_areas[i])

        # 渲染后自检：删除完全重复的 shape，告警空文本框
        self._audit_slide(slide, spec)

    def _audit_slide(self, slide, spec: SlideSpec):
        """
        渲染后质量门：
        - 删除坐标+尺寸+文本完全相同的重复 shape（保留第一个）
        - 删除空 TextBox（多见于 Layer4 计算出的空 body_areas）
        - 检测越界 / 异常超大尺寸
        - 检测正文文字框最小字号（< 10pt 提示密度过高）
        - 检测正文 shape 之间的强重叠（IoU > 0.6）
        """
        sw = self.prs.slide_width
        sh_h = self.prs.slide_height

        seen = set()
        dup_shapes = []
        empty_text_shapes = []
        offcanvas_shapes = []
        small_font_count = 0
        text_rects = []  # 用于重叠检测：(shape, l, t, r, b)

        for shape in list(slide.shapes):
            try:
                text = shape.text_frame.text if shape.has_text_frame else ""
            except Exception:
                text = ""

            fp = (shape.shape_type, shape.left, shape.top, shape.width, shape.height, text)
            if fp in seen and text:
                dup_shapes.append(shape)
                continue
            seen.add(fp)

            # 空 TextBox（17 = MSO_SHAPE_TYPE.TEXT_BOX）
            try:
                is_text_box = shape.shape_type == 17
            except Exception:
                is_text_box = False
            if is_text_box and not text.strip():
                empty_text_shapes.append(shape)
                continue

            # 越界检测
            l, t, w, h = shape.left or 0, shape.top or 0, shape.width or 0, shape.height or 0
            r, b = l + w, t + h
            if l < 0 or t < 0 or r > sw + 9144 or b > sh_h + 9144:
                offcanvas_shapes.append(shape)

            # 最小字号检测
            if shape.has_text_frame and text.strip():
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.font.size < Pt(10):
                            small_font_count += 1
                            break
                    else:
                        continue
                    break
                # 收集用于重叠检测
                if is_text_box:
                    text_rects.append((shape, l, t, r, b))

        # 删除重复
        for shape in dup_shapes:
            sp = shape._element
            sp.getparent().remove(sp)
        # 删除空 TextBox
        for shape in empty_text_shapes:
            sp = shape._element
            sp.getparent().remove(sp)

        # 重叠检测（仅文本框之间）
        overlap_pairs = []
        for i in range(len(text_rects)):
            _, l1, t1, r1, b1 = text_rects[i]
            a1 = max(0, (r1 - l1) * (b1 - t1))
            if a1 == 0:
                continue
            for j in range(i + 1, len(text_rects)):
                _, l2, t2, r2, b2 = text_rects[j]
                ix = max(0, min(r1, r2) - max(l1, l2))
                iy = max(0, min(b1, b2) - max(t1, t2))
                inter = ix * iy
                if inter == 0:
                    continue
                a2 = max(0, (r2 - l2) * (b2 - t2))
                union = a1 + a2 - inter
                iou = inter / union if union > 0 else 0
                if iou > 0.6:
                    overlap_pairs.append(iou)

        idx = spec.slide_index
        if dup_shapes:
            print(f"   ⚠️  slide {idx}: 删除了 {len(dup_shapes)} 个重复 shape")
        if empty_text_shapes:
            print(f"   ⚠️  slide {idx}: 删除了 {len(empty_text_shapes)} 个空文本框")
        if offcanvas_shapes:
            print(f"   ⚠️  slide {idx}: {len(offcanvas_shapes)} 个 shape 越出画布")
        if small_font_count:
            print(f"   ⚠️  slide {idx}: {small_font_count} 个文本框字号 < 10pt（信息密度过高）")
        if overlap_pairs:
            print(f"   ⚠️  slide {idx}: {len(overlap_pairs)} 对文本框严重重叠 (IoU>0.6)")

    def _add_title_decorations(self, slide, theme: VisualTheme):
        """封面页装饰：顶部宽色带 + 底部色带 + 强调细线"""
        primary = self._parse_color(theme.colors.get("primary", "#003D6E"))
        accent = self._parse_color(theme.colors.get("accent", "#FF6B35"))
        sw = self.prs.slide_width
        sh = self.prs.slide_height
        chrome.add_top_band(slide, sw, sh, primary, ratio=0.08)
        chrome.add_bottom_band(slide, sw, sh, primary, ratio=0.12)
        chrome.add_accent_underline(slide, sw, sh, accent, y_ratio=0.88)

    def _add_content_decorations(self, slide, theme: VisualTheme):
        """内容页装饰：顶部色条 + 左侧强调条 + 底部页脚浅色底"""
        primary = self._parse_color(theme.colors.get("primary", "#003D6E"))
        accent = self._parse_color(theme.colors.get("accent", "#FF6B35"))
        sw = self.prs.slide_width
        sh = self.prs.slide_height
        chrome.add_top_bar(slide, sw, primary)
        chrome.add_left_accent_bar(slide, sh, accent)
        chrome.add_footer_panel(slide, sw, sh, RGBColor(0xF5, 0xF6, 0xFA))

    def _build_title_slide(self, slide, spec: SlideSpec, layout, theme: VisualTheme):
        """
        封面页：主标题 + 副标题 + 分隔线 + 作者 + 日期，整体垂直居中。
        数据来源优先级：PresentationSpec.title/subtitle/author/created_at → spec.takeaway_message。
        """
        primary = self._parse_color(theme.colors.get("primary", "#003D6E"))
        accent = self._parse_color(theme.colors.get("accent", "#FF6B35"))
        text_light = self._parse_color(theme.colors.get("text_light", "#636E72"))
        title_font = theme.fonts.get("title", "Arial")
        body_font = theme.fonts.get("body", "Calibri")

        pres = self._pres_spec
        main_title = (pres.title if pres and pres.title else spec.takeaway_message) or "Presentation"
        subtitle = pres.subtitle if pres and pres.subtitle else ""
        author = pres.author if pres and pres.author else ""
        created = pres.created_at if pres and pres.created_at else ""
        if not created:
            from datetime import date
            created = date.today().strftime("%Y-%m-%d")

        sw = self.prs.slide_width
        sh = self.prs.slide_height

        # 1) 主标题：居中偏上
        title_top = int(sh * 0.32)
        title_box = slide.shapes.add_textbox(
            Emu(int(sw * 0.08)), Emu(title_top),
            Emu(int(sw * 0.84)), Emu(int(sh * 0.14))
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = main_title
        p.font.size = Pt(theme.font_sizes.get("title", 44))
        p.font.bold = True
        p.font.color.rgb = primary
        p.font.name = title_font
        p.alignment = PP_ALIGN.CENTER

        # 2) 橘色分隔短线
        bar_width = int(sw * 0.08)
        bar_left = (sw - bar_width) // 2
        bar_top = title_top + int(sh * 0.16)
        bar = slide.shapes.add_shape(
            1, bar_left, bar_top, bar_width, Emu(36000)  # ~4pt
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        # 3) 副标题（可选）
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Emu(int(sw * 0.12)), Emu(bar_top + int(sh * 0.025)),
                Emu(int(sw * 0.76)), Emu(int(sh * 0.08))
            )
            tf2 = sub_box.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(theme.font_sizes.get("subtitle", 22))
            p2.font.color.rgb = text_light
            p2.font.name = body_font
            p2.alignment = PP_ALIGN.CENTER

        # 4) 作者 / 日期：底部色带上方
        meta_parts = [x for x in [author, created] if x]
        if meta_parts:
            meta_box = slide.shapes.add_textbox(
                Emu(int(sw * 0.12)), Emu(int(sh * 0.80)),
                Emu(int(sw * 0.76)), Emu(int(sh * 0.06))
            )
            tfm = meta_box.text_frame
            pm = tfm.paragraphs[0]
            pm.text = "  ·  ".join(meta_parts)
            pm.font.size = Pt(theme.font_sizes.get("body", 14))
            pm.font.color.rgb = text_light
            pm.font.name = body_font
            pm.alignment = PP_ALIGN.CENTER

    def _add_section_divider_decorations(self, slide, theme: VisualTheme):
        """章节过渡页装饰：左侧 40% 主色块 + 右沿橘色竖条"""
        primary = self._parse_color(theme.colors.get("primary", "#003D6E"))
        accent = self._parse_color(theme.colors.get("accent", "#FF6B35"))
        sw = self.prs.slide_width
        sh = self.prs.slide_height
        chrome.add_left_panel(slide, sh, primary, width_ratio=int(sw * 0.40))
        # 右沿橘色竖条
        from pptx.util import Emu as _Emu
        bar = slide.shapes.add_shape(1, _Emu(int(sw * 0.40)), 0, _Emu(36000), _Emu(sh))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

    def _build_section_divider_slide(self, slide, spec: SlideSpec, layout, theme: VisualTheme):
        """
        章节过渡页：
        - 左侧大号章节编号 "01" / "02"（用 slide_index 或 narrative order）
        - 右侧：章节标题 + 简要描述
        """
        sw = self.prs.slide_width
        sh = self.prs.slide_height

        accent = self._parse_color(theme.colors.get("accent", "#FF6B35"))
        white = RGBColor(0xFF, 0xFF, 0xFF)
        text_light = self._parse_color(theme.colors.get("text_light", "#636E72"))
        primary = self._parse_color(theme.colors.get("primary", "#003D6E"))
        title_font = theme.fonts.get("title", "Arial")
        body_font = theme.fonts.get("body", "Calibri")

        # 1) 左侧大号编号：使用 slide_index（2 位数字）
        num = spec.slide_index or 1
        num_str = f"{num:02d}"
        num_box = slide.shapes.add_textbox(
            Emu(int(sw * 0.04)), Emu(int(sh * 0.25)),
            Emu(int(sw * 0.30)), Emu(int(sh * 0.50))
        )
        tfn = num_box.text_frame
        pn = tfn.paragraphs[0]
        pn.text = num_str
        pn.font.size = Pt(120)
        pn.font.bold = True
        pn.font.color.rgb = white
        pn.font.name = title_font
        pn.alignment = PP_ALIGN.LEFT

        # 2) 右侧章节标题
        title_text = spec.takeaway_message or "Section"
        title_box = slide.shapes.add_textbox(
            Emu(int(sw * 0.44)), Emu(int(sh * 0.38)),
            Emu(int(sw * 0.52)), Emu(int(sh * 0.18))
        )
        tft = title_box.text_frame
        tft.word_wrap = True
        pt_ = tft.paragraphs[0]
        pt_.text = title_text
        pt_.font.size = Pt(36)
        pt_.font.bold = True
        pt_.font.color.rgb = primary
        pt_.font.name = title_font
        pt_.alignment = PP_ALIGN.LEFT

        # 3) 橘色短分隔线
        bar = slide.shapes.add_shape(
            1, Emu(int(sw * 0.44)), Emu(int(sh * 0.58)),
            Emu(int(sw * 0.06)), Emu(36000)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        # 4) 章节描述（取首个 text_block 或 key_insight）
        desc = ""
        if spec.text_blocks:
            desc = spec.text_blocks[0].content or ""
        elif getattr(spec, "key_insights", None):
            desc = spec.key_insights[0] if spec.key_insights else ""
        if desc:
            desc_box = slide.shapes.add_textbox(
                Emu(int(sw * 0.44)), Emu(int(sh * 0.62)),
                Emu(int(sw * 0.52)), Emu(int(sh * 0.18))
            )
            tfd = desc_box.text_frame
            tfd.word_wrap = True
            pd = tfd.paragraphs[0]
            pd.text = desc
            pd.font.size = Pt(16)
            pd.font.color.rgb = text_light
            pd.font.name = body_font
            pd.alignment = PP_ALIGN.LEFT

    def _build_agenda_slide(self, slide, spec: SlideSpec, layout, theme: VisualTheme):
        """议程页：大号"目录" + 编号列表"""
        sw = self.prs.slide_width
        sh = self.prs.slide_height
        primary = self._parse_color(theme.colors.get("primary", "#003D6E"))
        accent = self._parse_color(theme.colors.get("accent", "#FF6B35"))
        text_dark = self._parse_color(theme.colors.get("text_dark", "#2D3436"))
        title_font = theme.fonts.get("title", "Arial")
        body_font = theme.fonts.get("body", "Calibri")

        # 页面标题
        title_text = spec.takeaway_message or "目录"
        hd_box = slide.shapes.add_textbox(
            Emu(int(sw * 0.06)), Emu(int(sh * 0.08)),
            Emu(int(sw * 0.88)), Emu(int(sh * 0.12))
        )
        tf = hd_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(theme.font_sizes.get("title", 32))
        p.font.bold = True
        p.font.color.rgb = primary
        p.font.name = title_font

        # 橘色短线
        bar = slide.shapes.add_shape(
            1, Emu(int(sw * 0.06)), Emu(int(sh * 0.21)),
            Emu(int(sw * 0.08)), Emu(36000)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        # 列表
        items = []
        for b in (spec.text_blocks or []):
            if b.content and b.level == 0:
                items.append(b.content)
        if not items:
            items = ["章节 1", "章节 2", "章节 3", "章节 4"]

        list_top = int(sh * 0.28)
        list_height = int(sh * 0.64)
        row_h = min(int(list_height / max(len(items), 1)), int(sh * 0.12))

        for i, item in enumerate(items[:6]):
            y = list_top + i * row_h
            # 编号
            num_box = slide.shapes.add_textbox(
                Emu(int(sw * 0.06)), Emu(y),
                Emu(int(sw * 0.08)), Emu(row_h)
            )
            tfn = num_box.text_frame
            pn = tfn.paragraphs[0]
            pn.text = f"{i + 1:02d}"
            pn.font.size = Pt(28)
            pn.font.bold = True
            pn.font.color.rgb = accent
            pn.font.name = title_font

            # 文本
            tx_box = slide.shapes.add_textbox(
                Emu(int(sw * 0.16)), Emu(y + int(row_h * 0.15)),
                Emu(int(sw * 0.78)), Emu(int(row_h * 0.7))
            )
            tft = tx_box.text_frame
            tft.word_wrap = True
            pt2 = tft.paragraphs[0]
            pt2.text = item
            pt2.font.size = Pt(20)
            pt2.font.color.rgb = text_dark
            pt2.font.name = body_font

    def _add_page_number(self, slide, spec: SlideSpec, theme: VisualTheme):
        """右下角页码"""
        if not spec.slide_index:
            return
        sw = self.prs.slide_width
        sh = self.prs.slide_height
        total = len(self._pres_spec.slides) if self._pres_spec else 0

        box = slide.shapes.add_textbox(
            Emu(int(sw * 0.88)), Emu(int(sh * 0.945)),
            Emu(int(sw * 0.10)), Emu(int(sh * 0.045))
        )
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{spec.slide_index} / {total}" if total else str(spec.slide_index)
        p.font.size = Pt(9)
        p.font.color.rgb = self._parse_color(theme.colors.get("text_light", "#636E72"))
        p.font.name = theme.fonts.get("footnote", "Calibri")
        p.alignment = PP_ALIGN.RIGHT

    def _build_content_slide(self, slide, spec: SlideSpec, layout, theme: VisualTheme):
        """构建内容页：标题(takeaway) + 正文区 + 来源"""
        # 1. 添加标题（即 takeaway message，带左侧强调小色块）
        if layout.title_area and spec.takeaway_message:
            self._add_takeaway_title(slide, spec.takeaway_message, layout.title_area, theme)

        # 2. 可视化块优先渲染（非bullet_list时走专属渲染器）
        vb = getattr(spec, 'visual_block', None)
        if vb and hasattr(vb, 'block_type') and vb.block_type.value != 'bullet_list' and vb.items:
            rendered = self._render_visual_block(slide, spec, layout, theme, vb)
            if rendered:
                if layout.source_area and spec.source_note:
                    self._add_source(slide, spec.source_note, layout.source_area, theme)
                return

        # 3. 添加文本内容到body_areas（传统路径）
        language = getattr(spec, 'language', 'zh') or 'zh'

        if spec.text_blocks and not layout.body_areas:
            layout.body_areas.append(self._emergency_body_rect(spec, layout))
            print(f"   ⚠️  slide {spec.slide_index}: layout 缺 body_area，已为 "
                  f"{len(spec.text_blocks)} 个 text_block 注入兜底区域")

        if layout.body_areas and spec.text_blocks:
            num_areas = len(layout.body_areas)
            if num_areas == 1:
                self._add_text_blocks(slide, spec.text_blocks, layout.body_areas[0],
                                      theme, language)
            else:
                groups = self._split_text_blocks_by_groups(spec.text_blocks, num_areas)
                for i, body_area in enumerate(layout.body_areas):
                    if i < len(groups) and groups[i]:
                        self._add_text_blocks(slide, groups[i], body_area,
                                              theme, language)

        # 4. 添加数据来源
        if layout.source_area and spec.source_note:
            self._add_source(slide, spec.source_note, layout.source_area, theme)

        # 5. 添加数据表格
        if layout.body_areas and spec.data_references:
            self._add_data_table(slide, spec.data_references, layout.body_areas[-1], theme)

    # ================================================================
    # 可视化块渲染器
    # ================================================================

    _VB_DISPATCH = {
        'kpi_cards': '_render_kpi_cards',
        'comparison_columns': '_render_comparison_columns',
        'step_cards': '_render_step_cards',
        'icon_text_grid': '_render_icon_text_grid',
        'stat_highlight': '_render_stat_highlight',
        'callout_box': '_render_callout_box',
    }

    def _render_visual_block(self, slide, spec: SlideSpec, layout, theme: VisualTheme, vb: VisualBlock) -> bool:
        """Dispatch到对应的可视化块渲染器，返回True表示渲染成功。
        优先使用 layout.visual_block_areas（逐 slot），fallback 用 body_areas[0]。
        """
        method_name = self._VB_DISPATCH.get(vb.block_type.value)
        if not method_name:
            return False
        method = getattr(self, method_name, None)
        if not method:
            return False
        try:
            vb_areas = getattr(layout, 'visual_block_areas', [])
            if vb_areas and len(vb_areas) >= len(vb.items):
                # slot 模式：每个 item 有独立坐标
                clamped = [self._clamp_rect(r) for r in vb_areas]
                method(slide, vb, clamped, theme)
            else:
                # 单区域模式：在一个大 body rect 内自行分割
                body_rect = layout.body_areas[0] if layout.body_areas else self._emergency_body_rect(spec, layout)
                method(slide, vb, self._clamp_rect(body_rect), theme)
            return True
        except Exception as e:
            print(f"   ⚠️  visual_block渲染失败({vb.block_type.value}): {e}，降级到文本模式")
            return False

    def _render_kpi_cards(self, slide, vb: VisualBlock, rect_or_rects, theme: VisualTheme):
        """KPI卡片：圆角矩形底色 + 超大数字 + 标签 + 增减趋势"""
        from pptx.enum.shapes import MSO_SHAPE
        items = vb.items[:4]
        n = len(items)
        if n == 0:
            return
        # list[Rect] slot 模式：每个 item 放到对应 slot
        if isinstance(rect_or_rects, list):
            rects = rect_or_rects
            for i, item in enumerate(items):
                if i >= len(rects):
                    break
                self._render_single_kpi_card(slide, item, rects[i], theme, i)
            return
        rect = rect_or_rects
        gap = Emu(91440)  # ~0.1 inch
        card_w = (rect.width - gap * (n - 1)) // n
        card_h = min(rect.height, Emu(2743200))  # max ~3 inch
        accent = self._parse_color(theme.colors.get("accent", "#FF6B35"))
        primary = self._parse_color(theme.colors.get("primary", "#003D6E"))
        bg_color = self._parse_color("#F0F4F8")
        green = self._parse_color("#27AE60")
        red = self._parse_color("#E74C3C")

        for i, item in enumerate(items):
            left = rect.left + i * (card_w + gap)
            # 背景卡片
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(rect.top),
                Emu(card_w), Emu(card_h)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = bg_color
            card.line.fill.background()

            # 大数字
            val_h = card_h * 45 // 100
            val_box = slide.shapes.add_textbox(
                Emu(left + card_w // 10), Emu(rect.top + card_h // 8),
                Emu(card_w * 8 // 10), Emu(val_h)
            )
            tf = val_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item.value or "—"
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = primary
            p.font.name = theme.fonts.get("title", "Arial")
            p.alignment = PP_ALIGN.CENTER

            # 标签
            label_top = rect.top + card_h // 8 + val_h
            label_box = slide.shapes.add_textbox(
                Emu(left + card_w // 10), Emu(label_top),
                Emu(card_w * 8 // 10), Emu(card_h // 5)
            )
            tf2 = label_box.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = item.title or ""
            p2.font.size = Pt(11)
            p2.font.color.rgb = self._parse_color(theme.colors.get("text_light", "#636E72"))
            p2.font.name = theme.fonts.get("body", "Calibri")
            p2.alignment = PP_ALIGN.CENTER

            # 趋势描述
            if item.description:
                desc_top = label_top + card_h // 5
                desc_box = slide.shapes.add_textbox(
                    Emu(left + card_w // 10), Emu(desc_top),
                    Emu(card_w * 8 // 10), Emu(card_h // 5)
                )
                tf3 = desc_box.text_frame
                tf3.word_wrap = True
                p3 = tf3.paragraphs[0]
                arrow = "▲ " if item.trend == "up" else ("▼ " if item.trend == "down" else "")
                p3.text = arrow + item.description
                p3.font.size = Pt(10)
                p3.font.color.rgb = green if item.trend == "up" else (red if item.trend == "down" else primary)
                p3.font.name = theme.fonts.get("body", "Calibri")
                p3.alignment = PP_ALIGN.CENTER

    def _render_single_kpi_card(self, slide, item: VisualBlockItem, rect: Rect, theme: VisualTheme, idx: int = 0):
        """渲染单个 KPI 卡片到指定 rect"""
        from pptx.enum.shapes import MSO_SHAPE
        primary = self._parse_color(theme.colors.get("primary", "#003D6E"))
        bg_color = self._parse_color("#F0F4F8")
        green = self._parse_color("#27AE60")
        red = self._parse_color("#E74C3C")

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Emu(rect.left), Emu(rect.top),
            Emu(rect.width), Emu(rect.height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.fill.background()

        val_h = rect.height * 45 // 100
        val_box = slide.shapes.add_textbox(
            Emu(rect.left + rect.width // 10), Emu(rect.top + rect.height // 8),
            Emu(rect.width * 8 // 10), Emu(val_h)
        )
        tf = val_box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item.value or "—"
        p.font.size = Pt(36); p.font.bold = True
        p.font.color.rgb = primary
        p.font.name = theme.fonts.get("title", "Arial")
        p.alignment = PP_ALIGN.CENTER

        label_top = rect.top + rect.height // 8 + val_h
        label_box = slide.shapes.add_textbox(
            Emu(rect.left + rect.width // 10), Emu(label_top),
            Emu(rect.width * 8 // 10), Emu(rect.height // 5)
        )
        tf2 = label_box.text_frame; tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = item.title or ""
        p2.font.size = Pt(11)
        p2.font.color.rgb = self._parse_color(theme.colors.get("text_light", "#636E72"))
        p2.font.name = theme.fonts.get("body", "Calibri")
        p2.alignment = PP_ALIGN.CENTER

        if item.description:
            desc_top = label_top + rect.height // 5
            desc_box = slide.shapes.add_textbox(
                Emu(rect.left + rect.width // 10), Emu(desc_top),
                Emu(rect.width * 8 // 10), Emu(rect.height // 5)
            )
            tf3 = desc_box.text_frame; tf3.word_wrap = True
            p3 = tf3.paragraphs[0]
            arrow = "▲ " if item.trend == "up" else ("▼ " if item.trend == "down" else "")
            p3.text = arrow + item.description
            p3.font.size = Pt(10)
            p3.font.color.rgb = green if item.trend == "up" else (red if item.trend == "down" else primary)
            p3.font.name = theme.fonts.get("body", "Calibri")
            p3.alignment = PP_ALIGN.CENTER

    def _render_step_cards(self, slide, vb: VisualBlock, rect_or_rects, theme: VisualTheme):
        """步骤卡片：编号圆圈 + 标题 + 描述 + 箭头连接"""
        from pptx.enum.shapes import MSO_SHAPE
        items = vb.items[:6]
        n = len(items)
        if n == 0:
            return
        # list[Rect] slot 模式：直接用提供的坐标
        if isinstance(rect_or_rects, list):
            rect = Rect(
                left=rect_or_rects[0].left,
                top=rect_or_rects[0].top,
                width=rect_or_rects[-1].left + rect_or_rects[-1].width - rect_or_rects[0].left,
                height=max(r.height for r in rect_or_rects),
            )
        else:
            rect = rect_or_rects
        gap = Emu(182880)  # ~0.2 inch (includes arrow space)
        card_w = (rect.width - gap * (n - 1)) // n
        accent = self._parse_color(theme.colors.get("accent", "#FF6B35"))
        primary = self._parse_color(theme.colors.get("primary", "#003D6E"))
        light = self._parse_color(theme.colors.get("text_light", "#636E72"))

        for i, item in enumerate(items):
            left = rect.left + i * (card_w + gap)

            # 编号圆圈
            circle_size = Emu(365760)  # ~0.4 inch
            circle_left = left + (card_w - circle_size) // 2
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Emu(circle_left), Emu(rect.top),
                Emu(circle_size), Emu(circle_size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = accent
            circle.line.fill.background()
            tf = circle.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.name = theme.fonts.get("title", "Arial")
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 标题
            title_top = rect.top + circle_size + Emu(91440)
            title_box = slide.shapes.add_textbox(
                Emu(left), Emu(title_top),
                Emu(card_w), Emu(Emu(274320))
            )
            tf2 = title_box.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = item.title or f"步骤{i+1}"
            p2.font.size = Pt(12)
            p2.font.bold = True
            p2.font.color.rgb = primary
            p2.font.name = theme.fonts.get("body", "Calibri")
            p2.alignment = PP_ALIGN.CENTER

            # 描述
            if item.description:
                desc_top = title_top + Emu(320040)
                desc_box = slide.shapes.add_textbox(
                    Emu(left), Emu(desc_top),
                    Emu(card_w), Emu(rect.height - (desc_top - rect.top))
                )
                tf3 = desc_box.text_frame
                tf3.word_wrap = True
                p3 = tf3.paragraphs[0]
                p3.text = item.description
                p3.font.size = Pt(9)
                p3.font.color.rgb = light
                p3.font.name = theme.fonts.get("body", "Calibri")
                p3.alignment = PP_ALIGN.CENTER

            # 箭头（非最后一个）
            if i < n - 1:
                arrow_left = left + card_w + Emu(27432)
                arrow_top = rect.top + circle_size // 2 - Emu(54864)
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW, Emu(arrow_left), Emu(arrow_top),
                    Emu(gap - Emu(54864)), Emu(Emu(109728))
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self._parse_color("#D0D5DD")
                arrow.line.fill.background()

    def _render_comparison_columns(self, slide, vb: VisualBlock, rect_or_rects, theme: VisualTheme):
        """对比栏：色块标题 + 下方要点列表"""
        from pptx.enum.shapes import MSO_SHAPE
        items = vb.items[:4]
        n = len(items)
        if n == 0:
            return
        if isinstance(rect_or_rects, list):
            rect = Rect(
                left=rect_or_rects[0].left,
                top=rect_or_rects[0].top,
                width=rect_or_rects[-1].left + rect_or_rects[-1].width - rect_or_rects[0].left,
                height=max(r.height for r in rect_or_rects),
            )
        else:
            rect = rect_or_rects
        gap = Emu(91440)
        col_w = (rect.width - gap * (n - 1)) // n
        primary = self._parse_color(theme.colors.get("primary", "#003D6E"))
        accent = self._parse_color(theme.colors.get("accent", "#FF6B35"))
        colors = [primary, accent,
                  self._parse_color("#00A878"),
                  self._parse_color("#6C5CE7")]

        for i, item in enumerate(items):
            left = rect.left + i * (col_w + gap)
            col_color = colors[i % len(colors)]

            # 标题色块
            header_h = Emu(365760)  # ~0.4 inch
            header = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(rect.top),
                Emu(col_w), Emu(header_h)
            )
            header.fill.solid()
            header.fill.fore_color.rgb = col_color
            header.line.fill.background()
            tf = header.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item.title or f"选项{i+1}"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.name = theme.fonts.get("body", "Calibri")
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 内容区域
            body_top = rect.top + header_h + Emu(45720)
            body_h = rect.height - header_h - Emu(45720)
            body_box = slide.shapes.add_textbox(
                Emu(left + Emu(45720)), Emu(body_top),
                Emu(col_w - Emu(91440)), Emu(body_h)
            )
            tf2 = body_box.text_frame
            tf2.word_wrap = True
            desc = item.description or ""
            lines = [l.strip() for l in desc.replace("\\n", "\n").split("\n") if l.strip()]
            for j, line in enumerate(lines):
                p2 = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
                bullet = "• " if not line.startswith("•") else ""
                p2.text = bullet + line
                p2.font.size = Pt(10)
                p2.font.color.rgb = self._parse_color(theme.colors.get("text_dark", "#2D3436"))
                p2.font.name = theme.fonts.get("body", "Calibri")
                p2.space_after = Pt(4)

    def _render_icon_text_grid(self, slide, vb: VisualBlock, rect_or_rects, theme: VisualTheme):
        """图标+文字网格：彩色小形状 + 标题 + 描述"""
        from pptx.enum.shapes import MSO_SHAPE
        items = vb.items[:6]
        n = len(items)
        if n == 0:
            return
        if isinstance(rect_or_rects, list):
            rect = Rect(
                left=rect_or_rects[0].left,
                top=rect_or_rects[0].top,
                width=rect_or_rects[-1].left + rect_or_rects[-1].width - rect_or_rects[0].left,
                height=max(r.height for r in rect_or_rects),
            )
        else:
            rect = rect_or_rects
        cols = vb.columns if vb.columns > 0 else min(3, n)
        rows = (n + cols - 1) // cols
        gap_x = Emu(91440)
        gap_y = Emu(91440)
        cell_w = (rect.width - gap_x * (cols - 1)) // cols
        cell_h = (rect.height - gap_y * (rows - 1)) // rows
        palette = [
            self._parse_color(c) for c in
            theme.colors.get("chart_palette", ["#003D6E", "#FF6B35", "#00A878", "#E17055", "#6C5CE7", "#FDCB6E"])
        ]
        primary = self._parse_color(theme.colors.get("primary", "#003D6E"))
        light = self._parse_color(theme.colors.get("text_light", "#636E72"))
        icon_shapes = [MSO_SHAPE.OVAL, MSO_SHAPE.DIAMOND, MSO_SHAPE.HEXAGON,
                       MSO_SHAPE.ROUNDED_RECTANGLE, MSO_SHAPE.PENTAGON, MSO_SHAPE.OCTAGON]

        for idx, item in enumerate(items):
            row = idx // cols
            col = idx % cols
            cell_left = rect.left + col * (cell_w + gap_x)
            cell_top = rect.top + row * (cell_h + gap_y)

            # 图标
            icon_size = Emu(274320)  # ~0.3 inch
            icon = slide.shapes.add_shape(
                icon_shapes[idx % len(icon_shapes)],
                Emu(cell_left + (cell_w - icon_size) // 2), Emu(cell_top),
                Emu(icon_size), Emu(icon_size)
            )
            icon.fill.solid()
            icon.fill.fore_color.rgb = palette[idx % len(palette)]
            icon.line.fill.background()

            # 标题
            title_top = cell_top + icon_size + Emu(45720)
            title_box = slide.shapes.add_textbox(
                Emu(cell_left), Emu(title_top),
                Emu(cell_w), Emu(Emu(228600))
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item.title or ""
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = primary
            p.font.name = theme.fonts.get("body", "Calibri")
            p.alignment = PP_ALIGN.CENTER

            # 描述
            if item.description:
                desc_top = title_top + Emu(228600)
                desc_box = slide.shapes.add_textbox(
                    Emu(cell_left), Emu(desc_top),
                    Emu(cell_w), Emu(cell_h - icon_size - Emu(274320))
                )
                tf2 = desc_box.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                p2.text = item.description
                p2.font.size = Pt(9)
                p2.font.color.rgb = light
                p2.font.name = theme.fonts.get("body", "Calibri")
                p2.alignment = PP_ALIGN.CENTER

    def _render_stat_highlight(self, slide, vb: VisualBlock, rect_or_rects, theme: VisualTheme):
        """核心数据高亮：超大字号居中 + 标签 + 上下文"""
        rect = rect_or_rects[0] if isinstance(rect_or_rects, list) else rect_or_rects
        item = vb.items[0] if vb.items else VisualBlockItem()
        primary = self._parse_color(theme.colors.get("primary", "#003D6E"))
        accent = self._parse_color(theme.colors.get("accent", "#FF6B35"))
        light = self._parse_color(theme.colors.get("text_light", "#636E72"))

        # 大数字
        val_h = rect.height * 45 // 100
        val_box = slide.shapes.add_textbox(
            Emu(rect.left), Emu(rect.top),
            Emu(rect.width), Emu(val_h)
        )
        tf = val_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item.value or "—"
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = accent
        p.font.name = theme.fonts.get("title", "Arial")
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM

        # 标签
        label_top = rect.top + val_h + Emu(91440)
        label_box = slide.shapes.add_textbox(
            Emu(rect.left), Emu(label_top),
            Emu(rect.width), Emu(Emu(365760))
        )
        tf2 = label_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = item.title or ""
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = primary
        p2.font.name = theme.fonts.get("body", "Calibri")
        p2.alignment = PP_ALIGN.CENTER

        # 上下文描述
        if item.description:
            desc_top = label_top + Emu(365760)
            desc_box = slide.shapes.add_textbox(
                Emu(rect.left + rect.width // 6), Emu(desc_top),
                Emu(rect.width * 2 // 3), Emu(rect.height - (desc_top - rect.top))
            )
            tf3 = desc_box.text_frame
            tf3.word_wrap = True
            p3 = tf3.paragraphs[0]
            p3.text = item.description
            p3.font.size = Pt(12)
            p3.font.color.rgb = light
            p3.font.name = theme.fonts.get("body", "Calibri")
            p3.alignment = PP_ALIGN.CENTER

    def _render_callout_box(self, slide, vb: VisualBlock, rect_or_rects, theme: VisualTheme):
        """引用框：左侧色条 + 大引号装饰 + 引用文字"""
        from pptx.enum.shapes import MSO_SHAPE
        rect = rect_or_rects[0] if isinstance(rect_or_rects, list) else rect_or_rects
        item = vb.items[0] if vb.items else VisualBlockItem()
        accent = self._parse_color(theme.colors.get("accent", "#FF6B35"))
        primary = self._parse_color(theme.colors.get("primary", "#003D6E"))
        light = self._parse_color(theme.colors.get("text_light", "#636E72"))

        # 左侧色条
        bar_w = Emu(54864)  # ~6pt
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(rect.left), Emu(rect.top),
            Emu(bar_w), Emu(rect.height)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        # 大引号
        quote_box = slide.shapes.add_textbox(
            Emu(rect.left + Emu(137160)), Emu(rect.top),
            Emu(Emu(365760)), Emu(Emu(365760))
        )
        tf = quote_box.text_frame
        p = tf.paragraphs[0]
        p.text = "\u201C"
        p.font.size = Pt(48)
        p.font.color.rgb = accent
        p.font.name = "Georgia"

        # 引用正文
        text_left = rect.left + Emu(182880)
        text_top = rect.top + Emu(320040)
        text_box = slide.shapes.add_textbox(
            Emu(text_left), Emu(text_top),
            Emu(rect.width - Emu(228600)), Emu(rect.height * 50 // 100)
        )
        tf2 = text_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = item.description or item.title or ""
        p2.font.size = Pt(14)
        p2.font.italic = True
        p2.font.color.rgb = primary
        p2.font.name = theme.fonts.get("body", "Calibri")

        # 出处/标题
        if item.title and item.description:
            source_top = text_top + rect.height * 50 // 100 + Emu(45720)
            source_box = slide.shapes.add_textbox(
                Emu(text_left), Emu(source_top),
                Emu(rect.width - Emu(228600)), Emu(Emu(228600))
            )
            tf3 = source_box.text_frame
            p3 = tf3.paragraphs[0]
            p3.text = "— " + item.title
            p3.font.size = Pt(10)
            p3.font.color.rgb = light
            p3.font.name = theme.fonts.get("body", "Calibri")

    def _add_title(self, slide, text: str, rect, theme: VisualTheme):
        """添加页面标题"""
        title_box = slide.shapes.add_textbox(
            Emu(rect.left), Emu(rect.top),
            Emu(rect.width), Emu(rect.height)
        )

        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(theme.font_sizes.get("title", 28))
        p.font.bold = True
        p.font.color.rgb = self._parse_color(theme.colors.get("primary", "#003D6E"))
        p.font.name = theme.fonts.get("title", "Arial")

    def _add_takeaway_title(self, slide, text: str, rect, theme: VisualTheme):
        """
        Takeaway 强调标题：左侧 4pt 橘色短竖条 + 标题文字 + 下方细分割线。
        让 takeaway 在视觉上明显区别于普通正文，符合咨询级页面"先抛结论"的范式。
        """
        primary = self._parse_color(theme.colors.get("primary", "#003D6E"))
        accent = self._parse_color(theme.colors.get("accent", "#FF6B35"))
        title_font = theme.fonts.get("title", "Arial")

        # 左侧短竖条（高度同标题）
        bar_w = Emu(54864)  # ~6pt
        side_bar = slide.shapes.add_shape(
            1, Emu(rect.left), Emu(rect.top + int(rect.height * 0.10)),
            bar_w, Emu(int(rect.height * 0.80))
        )
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = accent
        side_bar.line.fill.background()

        # 标题文本（左移 ~0.12 inch 让出竖条空间）
        text_left = rect.left + 109728
        title_box = slide.shapes.add_textbox(
            Emu(text_left), Emu(rect.top),
            Emu(rect.width - 109728), Emu(rect.height)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(theme.font_sizes.get("title", 24))
        p.font.bold = True
        p.font.color.rgb = primary
        p.font.name = title_font

        # 标题下方细分割线
        line = slide.shapes.add_shape(
            1, Emu(rect.left), Emu(rect.top + rect.height + 18288),
            Emu(rect.width), Emu(9144)  # ~1pt
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._parse_color(
            theme.colors.get("text_light", "#D0D5DD")
        )
        line.line.fill.background()

    # 各层级可见 bullet 前缀（ASCII 近似 ▪ ● ·）
    _BULLET_PREFIX = {0: "■  ", 1: "•  ", 2: "–  "}

    def _add_text_blocks(self, slide, text_blocks: List, rect, theme: VisualTheme,
                         language: str = "zh"):
        """
        添加文本块，每个 block 独立 paragraph，带 level 可见前缀 + 层级字号 +
        颜色区分 + 段前/后间距。字号由 theme.font_sizes 决定（已在 Layer4 按
        language 调整），此处不再重复 bump。
        """
        text_box = slide.shapes.add_textbox(
            Emu(rect.left), Emu(rect.top),
            Emu(rect.width), Emu(rect.height)
        )

        tf = text_box.text_frame
        tf.word_wrap = True

        body_size = theme.font_sizes.get("body", 14)
        bullet1_size = max(body_size - 1, 10)
        bullet2_size = max(body_size - 2, 9)
        text_color = self._parse_color(theme.colors.get("text_dark", "#2D3436"))
        light_color = self._parse_color(theme.colors.get("text_light", "#636E72"))
        primary_color = self._parse_color(theme.colors.get("primary", "#003D6E"))
        body_font = theme.fonts.get("body", "Calibri")

        for i, block in enumerate(text_blocks):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            level = max(0, min(int(getattr(block, "level", 0) or 0), 2))
            prefix = self._BULLET_PREFIX.get(level, "")
            p.text = prefix + (block.content or "")
            p.level = level

            # 层级字号 + 颜色 + 加粗（level-0 作为段落 heading 默认加粗）
            if level == 0:
                p.font.size = Pt(body_size)
                p.font.color.rgb = primary_color
                p.font.bold = True
                p.space_before = Pt(12)
                p.space_after = Pt(6)
            elif level == 1:
                p.font.size = Pt(bullet1_size)
                p.font.color.rgb = text_color
                p.font.bold = bool(getattr(block, "is_bold", False))
                p.space_before = Pt(6)
                p.space_after = Pt(3)
            else:
                p.font.size = Pt(bullet2_size)
                p.font.color.rgb = light_color
                p.font.bold = False
                p.space_before = Pt(3)
                p.space_after = Pt(2)

            p.font.name = body_font
            # 显式覆写（is_bold 可强制覆盖 level 默认）
            if getattr(block, "is_bold", False):
                p.font.bold = True

    def _add_source(self, slide, source_text: str, rect, theme: VisualTheme):
        """添加数据来源标注"""
        source_box = slide.shapes.add_textbox(
            Emu(rect.left), Emu(rect.top),
            Emu(rect.width), Emu(rect.height)
        )

        tf = source_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"来源：{source_text}"
        p.font.size = Pt(theme.font_sizes.get("footnote", 9))
        p.font.color.rgb = self._parse_color(theme.colors.get("text_light", "#636E72"))
        p.font.name = theme.fonts.get("footnote", "Calibri")

    def _add_so_what(self, slide, text: str, rect, theme: VisualTheme):
        """添加图表的so_what结论"""
        box = slide.shapes.add_textbox(
            Emu(rect.left), Emu(rect.top),
            Emu(rect.width), Emu(rect.height)
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(theme.font_sizes.get("chart_label", 10))
        p.font.italic = True
        p.font.color.rgb = self._parse_color(theme.colors.get("accent", "#FF6B35"))
        p.font.name = theme.fonts.get("body", "Calibri")
        p.alignment = PP_ALIGN.CENTER

    def _add_data_table(self, slide, data_refs: list, rect, theme: VisualTheme):
        """渲染数据表格"""
        if not data_refs or not rect:
            return

        # 从 data_references 中取第一个有表格数据的
        table_data = None
        for ref in data_refs:
            if hasattr(ref, 'table') and ref.table:
                table_data = ref.table
                break

        if not table_data:
            return

        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])
        if not headers or not rows:
            return

        # 限制行列数
        max_cols = min(len(headers), 8)
        max_rows = min(len(rows), 10)
        headers = headers[:max_cols]
        rows = [r[:max_cols] for r in rows[:max_rows]]

        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)

        # 表格尺寸
        table_width = min(rect.width, Inches(10))
        row_height = Emu(int(rect.height / num_rows)) if rect.height else Inches(0.35)
        col_width = Emu(int(table_width / num_cols))

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Emu(rect.left), Emu(rect.top),
            Emu(table_width), Emu(row_height * num_rows)
        )
        table = table_shape.table

        # 设置列宽
        for i in range(num_cols):
            table.columns[i].width = col_width

        primary_color = self._parse_color(theme.colors.get("primary", "#003D6E"))
        text_color = self._parse_color(theme.colors.get("text_dark", "#2D3436"))
        light_bg = RGBColor(0xF5, 0xF6, 0xFA)
        font_name = theme.fonts.get("body", "Calibri")

        # 表头
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = primary_color
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = font_name
            p.alignment = PP_ALIGN.CENTER

        # 数据行（交替色）
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(val) if val else ""
                # 交替行背景
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = light_bg
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(9)
                p.font.color.rgb = text_color
                p.font.name = font_name
                p.alignment = PP_ALIGN.CENTER

    def _add_chart(self, slide, chart_spec: ChartSpec, rect: Rect, theme: VisualTheme):
        """使用python-pptx原生图表对象渲染图表 + 智能标注"""
        if not chart_spec.categories or not chart_spec.series:
            return
        rect = self._clamp_rect(rect)

        # 瀑布图：用堆叠柱状图模拟
        if chart_spec.chart_type == ChartType.WATERFALL:
            self._add_waterfall_chart(slide, chart_spec, rect, theme)
            return

        # 组合图：柱+线
        if chart_spec.chart_type == ChartType.COMBO and len(chart_spec.series) >= 2:
            self._add_combo_chart(slide, chart_spec, rect, theme)
            return

        # 1. 准备数据（防御：LLM可能输出字符串数值）
        chart_data = CategoryChartData()
        chart_data.categories = chart_spec.categories
        for series in chart_spec.series:
            safe_vals = []
            for v in series.values:
                try:
                    safe_vals.append(float(v))
                except (ValueError, TypeError):
                    safe_vals.append(0)
            chart_data.add_series(series.name, safe_vals)

        # 2. 获取图表类型
        xl_type = CHART_TYPE_MAP.get(chart_spec.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        # 3. 添加图表形状
        chart_frame = slide.shapes.add_chart(
            xl_type,
            Emu(rect.left), Emu(rect.top),
            Emu(rect.width), Emu(rect.height),
            chart_data
        )

        chart = chart_frame.chart

        # 4. 去网格线（咨询级标准）
        # 注意：饼图 / 环形图 / 散点图等没有 value_axis / category_axis，
        # 访问会抛 ValueError，必须 try/except 包裹。
        has_axes = chart_spec.chart_type not in (ChartType.PIE,)
        if has_axes:
            try:
                value_axis = chart.value_axis
                value_axis.has_major_gridlines = False
                value_axis.has_minor_gridlines = False
                value_axis.format.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
                value_axis.format.line.width = Pt(0.5)
            except (ValueError, AttributeError):
                pass
            try:
                cat_axis = chart.category_axis
                cat_axis.format.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
                cat_axis.format.line.width = Pt(0.5)
            except (ValueError, AttributeError):
                pass

        # 5. 图表标题
        if chart_spec.title:
            chart.has_title = True
            chart.chart_title.text_frame.paragraphs[0].text = chart_spec.title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(
                theme.font_sizes.get("chart_title", 14)
            )
            chart.chart_title.text_frame.paragraphs[0].font.name = theme.fonts.get("body", "Calibri")
        else:
            chart.has_title = False

        # 6. 图例放底部
        chart.has_legend = chart_spec.show_legend
        if chart.has_legend:
            chart.legend.font.size = Pt(theme.font_sizes.get("chart_label", 9))
            chart.legend.include_in_layout = False

        # 7. 应用主题颜色
        palette = theme.colors.get("chart_palette", [])
        plot = chart.plots[0]
        # gap_width 仅适用于 bar/column；pie/scatter/line 设置会抛错
        if chart_spec.chart_type in (ChartType.BAR, ChartType.COLUMN):
            try:
                plot.gap_width = 100
            except Exception:
                pass

        # 计算最大最小值用于标注（防御：LLM可能输出字符串数值）
        all_values = []
        for s in chart_spec.series:
            for v in s.values:
                try:
                    all_values.append(float(v))
                except (ValueError, TypeError):
                    pass

        max_val = max(all_values) if all_values else None
        min_val = min(all_values) if all_values else None

        # 饼图：颜色作用于单点而非系列
        if chart_spec.chart_type == ChartType.PIE:
            try:
                pie_series = plot.series[0]
                for i, point in enumerate(pie_series.points):
                    color_str = ""
                    if i < len(palette):
                        color_str = palette[i]
                    if color_str:
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = self._parse_color(color_str)
            except Exception:
                pass
        else:
            for i, series in enumerate(plot.series):
                color_str = ""
                if i < len(chart_spec.series) and chart_spec.series[i].color:
                    color_str = chart_spec.series[i].color
                elif i < len(palette):
                    color_str = palette[i]

                if color_str:
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = self._parse_color(color_str)

        # 8. 选择性数据标签（只标最大/最小值）
        plot.has_data_labels = False  # 不全标
        if all_values and len(all_values) <= 12:
            # 少量数据时全标
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = Pt(theme.font_sizes.get("chart_label", 9))
            data_labels.font.name = theme.fonts.get("body", "Calibri")
            data_labels.number_format = '0.0'
        elif max_val is not None and min_val is not None:
            # 多数据时只标最大最小 → 用文本框标注
            self._add_min_max_labels(
                slide, chart_frame, chart_spec,
                max_val, min_val, theme
            )

        # 9. 均值基准线（单系列时，饼图除外）
        if (len(chart_spec.series) == 1 and all_values
                and chart_spec.chart_type not in (ChartType.PIE,)):
            avg = sum(all_values) / len(all_values)
            self._add_average_line(slide, chart_frame, avg, rect, theme)

        # 10. so_what 结论文字（在图表下方）
        if chart_spec.so_what:
            so_what_area = Rect(
                left=rect.left,
                top=rect.top + rect.height + 91440,  # 图表下方 0.1inch
                width=rect.width,
                height=365760  # 0.4inch
            )
            self._add_so_what(slide, chart_spec.so_what, so_what_area, theme)

    def _add_waterfall_chart(self, slide, chart_spec: ChartSpec, rect: Rect, theme: VisualTheme):
        """瀑布图：用堆叠柱状图模拟（透明底色+增量色块）"""
        if not chart_spec.series:
            return

        values = chart_spec.series[0].values
        categories = chart_spec.categories

        # 计算基底和增量
        bases = []
        increments = []
        running = 0
        for v in values:
            bases.append(running)
            increments.append(v)
            running += v

        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series("基底", bases)   # 透明
        chart_data.add_series("增量", increments)  # 彩色

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED,
            Emu(rect.left), Emu(rect.top),
            Emu(rect.width), Emu(rect.height),
            chart_data
        )
        chart = chart_frame.chart
        chart.has_title = bool(chart_spec.title)
        if chart_spec.title:
            chart.chart_title.text_frame.paragraphs[0].text = chart_spec.title

        try:
            chart.value_axis.has_major_gridlines = False
        except (ValueError, AttributeError):
            pass

        plot = chart.plots[0]
        plot.gap_width = 50

        # 底层透明
        base_series = plot.series[0]
        base_series.format.fill.background()  # 无填充

        # 增量层着色
        accent = self._parse_color(theme.colors.get("accent", "#FF6B35"))
        inc_series = plot.series[1]
        inc_series.format.fill.solid()
        inc_series.format.fill.fore_color.rgb = accent

        # 数据标签（只标增量）
        inc_series.has_data_labels = True
        inc_series.data_labels.font.size = Pt(9)
        inc_series.data_labels.font.name = theme.fonts.get("body", "Calibri")
        inc_series.data_labels.number_format = '0.0'

        # so_what
        if chart_spec.so_what:
            so_what_area = Rect(
                left=rect.left, top=rect.top + rect.height + 91440,
                width=rect.width, height=365760
            )
            self._add_so_what(slide, chart_spec.so_what, so_what_area, theme)

    def _add_combo_chart(self, slide, chart_spec: ChartSpec, rect: Rect, theme: VisualTheme):
        """
        组合图：python-pptx 不支持原生 combo，此前实现采用多 chart 叠加，导致
        PPT 中出现多个视觉重叠的图表。改为单 chart 簇状柱状图，所有系列按主题色
        分色展示；第一系列使用 primary，其余依次取 chart_palette。
        """
        chart_data = CategoryChartData()
        chart_data.categories = chart_spec.categories
        for series in chart_spec.series:
            safe_vals = [float(v) if not isinstance(v, (int, float)) else v for v in series.values]
            chart_data.add_series(series.name, safe_vals)

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Emu(rect.left), Emu(rect.top),
            Emu(rect.width), Emu(rect.height),
            chart_data
        )
        chart = chart_frame.chart
        chart.has_title = bool(chart_spec.title)
        if chart_spec.title:
            chart.chart_title.text_frame.paragraphs[0].text = chart_spec.title
        try:
            chart.value_axis.has_major_gridlines = False
        except (ValueError, AttributeError):
            pass
        chart.has_legend = len(chart_spec.series) > 1
        if chart.has_legend:
            chart.legend.include_in_layout = False

        palette = theme.colors.get("chart_palette", []) or []
        primary = self._parse_color(theme.colors.get("primary", "#003D6E"))
        accent = self._parse_color(theme.colors.get("accent", "#FF6B35"))
        fallback = [primary, accent]

        plot = chart.plots[0]
        plot.gap_width = 100
        for i, series in enumerate(plot.series):
            color = None
            if i < len(chart_spec.series) and chart_spec.series[i].color:
                color = self._parse_color(chart_spec.series[i].color)
            elif i < len(palette):
                color = self._parse_color(palette[i])
            else:
                color = fallback[i % len(fallback)]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color

        if chart_spec.so_what:
            so_what_area = Rect(
                left=rect.left, top=rect.top + rect.height + 91440,
                width=rect.width, height=365760
            )
            self._add_so_what(slide, chart_spec.so_what, so_what_area, theme)

    def _add_min_max_labels(self, slide, chart_frame, chart_spec,
                             max_val, min_val, theme):
        """在图表上标注最大/最小值"""
        accent = self._parse_color(theme.colors.get("accent", "#FF6B35"))
        gray = RGBColor(0x63, 0x6E, 0x72)

        for si, series in enumerate(chart_spec.series):
            vals = series.values
            max_idx = vals.index(max_val) if max_val in vals else -1
            min_idx = vals.index(min_val) if min_val in vals else -1

            if max_idx >= 0:
                point = chart_frame.chart.plots[0].series[si].points[max_idx]
                point.data_label.has_text_frame = True
                point.data_label.text_frame.paragraphs[0].text = str(max_val)
                point.data_label.text_frame.paragraphs[0].font.size = Pt(10)
                point.data_label.text_frame.paragraphs[0].font.bold = True
                point.data_label.text_frame.paragraphs[0].font.color.rgb = accent

            if min_idx >= 0 and min_idx != max_idx:
                point = chart_frame.chart.plots[0].series[si].points[min_idx]
                point.data_label.has_text_frame = True
                point.data_label.text_frame.paragraphs[0].text = str(min_val)
                point.data_label.text_frame.paragraphs[0].font.size = Pt(10)
                point.data_label.text_frame.paragraphs[0].font.color.rgb = gray

    def _emergency_body_rect(self, spec: SlideSpec, layout) -> Rect:
        """
        Layer4 选择的 pattern 可能未生成 body_area（如 TIMELINE/PROCESS_FLOW/
        TITLE_ONLY 等），为避免正文被静默丢弃，计算一个占用 slide 剩余空间的兜底
        body 矩形。优先避开 title_area、chart/diagram/picture 区域与 source_area。
        """
        sw = self.prs.slide_width
        sh = self.prs.slide_height
        margin = 457200  # 0.5 inch
        gap = 91440  # 0.1 inch

        top = margin
        if layout.title_area:
            top = max(top, layout.title_area.top + layout.title_area.height + gap)

        bottom = sh - margin
        if layout.source_area:
            bottom = min(bottom, layout.source_area.top - gap)

        busy_areas = (layout.chart_areas + layout.diagram_areas +
                      (getattr(layout, "picture_areas", None) or []))

        if busy_areas:
            busy_top = min(ar.top for ar in busy_areas)
            busy_bottom = max(ar.top + ar.height for ar in busy_areas)

            space_above = busy_top - gap - top
            space_below = bottom - (busy_bottom + gap)

            if space_above >= 914400:
                bottom = busy_top - gap
            elif space_below >= 914400:
                top = busy_bottom + gap
            else:
                # 无足够空间 → 用 body 上 1/3，把 busy 区域下移到 2/3
                total = bottom - top
                body_h = total // 3
                bottom = top + body_h
                new_visual_top = top + body_h + gap
                for ar in busy_areas:
                    remaining = bottom + gap + (ar.top + ar.height - busy_top)
                    ar.top = new_visual_top + (ar.top - busy_top)

        height = max(bottom - top, 914400)
        return Rect(
            left=margin,
            top=top,
            width=sw - margin * 2,
            height=height,
        )

    @staticmethod
    def _split_text_blocks_by_groups(text_blocks: List, num_groups: int) -> List[List]:
        """按 level-0 段落分组，均匀分配到 num_groups 个区域"""
        # 先按 level-0 切分成逻辑组
        groups_raw = []
        current_group = []
        for block in text_blocks:
            if block.level == 0 and current_group:
                groups_raw.append(current_group)
                current_group = []
            current_group.append(block)
        if current_group:
            groups_raw.append(current_group)

        # 将逻辑组均匀分配到 num_groups 个桶
        result = [[] for _ in range(num_groups)]
        for i, grp in enumerate(groups_raw):
            bucket = min(i * num_groups // max(len(groups_raw), 1), num_groups - 1)
            result[bucket].extend(grp)

        return result

    def _add_average_line(self, slide, chart_frame, avg, rect, theme):
        """添加均值虚线"""
        try:
            # 用一条水平线近似表示均值
            line_y = rect.top + int(rect.height * 0.5)  # 近似位置
            line = slide.shapes.add_shape(
                1,  # rectangle
                Emu(rect.left), Emu(line_y),
                Emu(rect.width), Emu(13716),  # 1.5pt 高度
            )
            line.fill.solid()
            line.fill.fore_color.rgb = self._parse_color(
                theme.colors.get("accent", "#FF6B35")
            )
            line.line.fill.background()

            # 均值标签
            label_box = slide.shapes.add_textbox(
                Emu(rect.left + rect.width - 914400),  # 右侧
                Emu(line_y - 137160),
                Emu(914400), Emu(274320)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"均值 {avg:.1f}"
            p.font.size = Pt(8)
            p.font.italic = True
            p.font.color.rgb = self._parse_color(
                theme.colors.get("text_light", "#636E72")
            )
        except Exception:
            pass

    def _add_diagram(self, slide, diag_spec: DiagramSpec, rect: Rect, theme: VisualTheme):
        """使用 DiagramRenderer 渲染概念图"""
        if not diag_spec.nodes:
            return
        rect = self._clamp_rect(rect)
        self.diagram_renderer.render(slide, diag_spec, rect, theme)

    def _add_picture(self, slide, pic_path: str, rect: Rect):
        """
        把原材料图片按 rect 等比缩放居中嵌入。失败(文件不存在/格式不支持) 时
        打印告警但不中断 build。
        """
        import os
        if not pic_path or not os.path.exists(pic_path):
            print(f"   ⚠️  picture 不存在，已跳过: {pic_path}")
            return
        try:
            # 先按 width 约束，python-pptx 会保持纵横比；若算出的高度超 rect 则
            # 再按 height 约束一次。
            pic = slide.shapes.add_picture(
                pic_path,
                Emu(rect.left), Emu(rect.top),
                width=Emu(rect.width),
            )
            if pic.height > rect.height:
                sp = pic._element
                sp.getparent().remove(sp)
                pic = slide.shapes.add_picture(
                    pic_path,
                    Emu(rect.left), Emu(rect.top),
                    height=Emu(rect.height),
                )
            # 水平/垂直居中
            dx = (rect.width - pic.width) // 2
            dy = (rect.height - pic.height) // 2
            pic.left = Emu(rect.left + max(dx, 0))
            pic.top = Emu(rect.top + max(dy, 0))
        except Exception as e:
            print(f"   ⚠️  picture 渲染失败 {pic_path}: {type(e).__name__}: {e}")

    def _clamp_rect(self, rect: Rect) -> Rect:
        """将坐标裁剪到幻灯片画布内，防止越界渲染"""
        sw = int(self.prs.slide_width)
        sh = int(self.prs.slide_height)
        min_dim = 91440  # ~0.1 inch 最小尺寸
        left = max(0, min(rect.left, sw - min_dim))
        top = max(0, min(rect.top, sh - min_dim))
        width = min(rect.width, sw - left)
        height = min(rect.height, sh - top)
        return Rect(left, top, max(width, min_dim), max(height, min_dim))

    @staticmethod
    def _parse_color(color_str: str) -> RGBColor:
        """解析颜色字符串为RGBColor"""
        if not color_str:
            return RGBColor(0, 0, 0)
        hex_str = color_str.lstrip("#")
        if len(hex_str) == 6:
            return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
        return RGBColor(0, 0, 0)


def create_hardcoded_demo() -> str:
    """
    创建一个硬编码内容的演示PPT（含图表）
    """
    from models import SlideType, NarrativeRole, TextBlock, ContentPattern, ChartSeries

    print("=" * 60)
    print("创建演示PPT（含图表）")
    print("=" * 60)

    pres_spec = PresentationSpec(
        title="四大级别PPT演示",
        subtitle="专业咨询演示样例",
        language="zh"
    )

    # 第1页：标题页
    title_slide = SlideSpec(
        slide_type=SlideType.TITLE,
        takeaway_message="数字化转型战略方案",
        narrative_arc=NarrativeRole.OPENING,
        content_pattern=ContentPattern.TITLE_ONLY,
    )
    pres_spec.slides.append(title_slide)

    # 第2页：内容页 - 含图表
    content_slide = SlideSpec(
        slide_type=SlideType.CONTENT,
        takeaway_message="各业务线收入对比分析",
        narrative_arc=NarrativeRole.EVIDENCE,
        content_pattern=ContentPattern.LEFT_CHART_RIGHT_TEXT,
        text_blocks=[
            TextBlock(content="线上业务增速显著，同比增长67%", level=0),
            TextBlock(content="传统渠道收入持续承压", level=0),
            TextBlock(content="建议加大线上投入力度", level=0),
        ],
        charts=[ChartSpec(
            chart_type=ChartType.COLUMN,
            title="各业务线季度收入（百万元）",
            categories=["线上零售", "企业服务", "广告营销", "金融科技"],
            series=[
                ChartSeries(name="2024 Q3", values=[320, 180, 95, 140]),
                ChartSeries(name="2024 Q4", values=[380, 210, 110, 165]),
            ],
            so_what="线上零售Q4环比增长18.8%，增速最快",
            show_legend=True,
            show_data_labels=True,
        )],
        source_note="公司内部财务数据，2024Q4"
    )
    pres_spec.slides.append(content_slide)

    # 第3页：纯文本页
    content_slide2 = SlideSpec(
        slide_type=SlideType.CONTENT,
        takeaway_message="当前面临三大核心挑战",
        narrative_arc=NarrativeRole.PROBLEM,
        content_pattern=ContentPattern.ARGUMENT_EVIDENCE,
        text_blocks=[
            TextBlock(content="业务增长放缓，传统渠道获客成本上升", level=0),
            TextBlock(content="2024年线上获客成本同比增加45%", level=1),
            TextBlock(content="系统架构老旧，难以支持快速迭代", level=0),
            TextBlock(content="平均需求交付周期长达3个月", level=1),
            TextBlock(content="数据孤岛严重，缺乏统一客户视图", level=0),
            TextBlock(content="客户数据分散在12个系统中", level=1),
        ],
        source_note="公司内部数据分析，2024Q1"
    )
    pres_spec.slides.append(content_slide2)

    # 构建PPT
    builder = PPTBuilder()

    # 应用Layer 4视觉设计
    from pipeline.layer4_visual import VisualDesigner
    designer = VisualDesigner()
    designer.design_slides(pres_spec.slides, language="zh")

    output_path = builder.build(pres_spec)

    print("\n" + "=" * 60)
    print(f"✨ 演示PPT生成完成！")
    print(f"📁 文件路径: {output_path}")
    print("=" * 60)

    return output_path


if __name__ == "__main__":
    create_hardcoded_demo()
