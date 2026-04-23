"""
PPTX 文件解析器
读取已有PPT文件，提取文本、表格和结构信息
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from models import RawContent, TableData, SourcePage


class PptxParser:
    """PPTX文件读取器，使用python-pptx"""

    def parse(self, file_path: str) -> RawContent:
        """
        解析PPTX文件为RawContent

        Args:
            file_path: .pptx文件路径

        Returns:
            RawContent with source_type="ppt"
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        try:
            prs = Presentation(file_path)
        except Exception as e:
            raise ValueError(f"无法解析PPTX文件，文件可能已损坏: {e}")

        all_text_parts = []
        all_tables = []
        source_pages = []

        for idx, slide in enumerate(prs.slides, 1):
            # 提取文本
            slide_text = self._extract_text_from_slide(slide)
            if slide_text.strip():
                all_text_parts.append(f"--- 第{idx}页 ---\n{slide_text}")
                title = self._extract_slide_title(slide) or f"第{idx}页"
                source_pages.append(SourcePage(
                    title=title,
                    content=slide_text,
                    page_number=idx,
                ))

            # 提取表格
            slide_tables = self._extract_tables_from_slide(slide, idx)
            all_tables.extend(slide_tables)

        raw_text = "\n\n".join(all_text_parts)
        metadata = self._extract_metadata(prs, file_path)

        # 检测语言
        from .text_parser import TextParser
        lang = TextParser()._detect_language(raw_text)

        return RawContent(
            source_type="ppt",
            raw_text=raw_text,
            tables=all_tables,
            metadata=metadata,
            detected_language=lang,
            source_pages=source_pages,
        )

    def _extract_slide_title(self, slide) -> str:
        """提取幻灯片标题：优先 title placeholder，其次首个短文本。"""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                pf = getattr(shape, "placeholder_format", None)
                if pf is not None and pf.idx == 0:
                    text = shape.text_frame.text.strip()
                    if text:
                        return text[:80]
            except Exception:
                pass
        # fallback: first non-empty paragraph shorter than 80 chars (likely a title)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and len(text) < 80:
                        return text
        return ""

    def _extract_text_from_slide(self, slide) -> str:
        """提取单页幻灯片中的所有文本"""
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        return "\n".join(texts)

    def _extract_tables_from_slide(self, slide, slide_idx: int) -> list:
        """提取幻灯片中的表格"""
        tables = []
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows_data = []

                for row in table.rows:
                    cells = []
                    for cell in row.cells:
                        cells.append(cell.text.strip())
                    rows_data.append(cells)

                if rows_data:
                    headers = rows_data[0]
                    data = rows_data[1:]
                    tables.append(TableData(
                        headers=headers,
                        rows=data,
                        source_sheet=f"第{slide_idx}页",
                        source_range=f"表格{len(tables) + 1}",
                    ))
        return tables

    def _extract_metadata(self, prs: Presentation, file_path: str) -> dict:
        """提取PPT元数据"""
        meta = {
            "file_name": Path(file_path).name,
            "file_size": os.path.getsize(file_path),
            "slide_count": len(prs.slides),
            "slide_width_inches": round(prs.slide_width / 914400, 2),
            "slide_height_inches": round(prs.slide_height / 914400, 2),
        }

        # 尝试获取第一页标题
        if prs.slides:
            first_slide = prs.slides[0]
            for shape in first_slide.shapes:
                if shape.has_text_frame:
                    title = shape.text_frame.text.strip()
                    if title:
                        meta["original_title"] = title
                        break

        return meta
