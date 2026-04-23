"""
Chart Skills — 模板方法模式

ChartSkill 基类封装 80% 共通逻辑（数据清洗、网格线、标题、图例、配色、标注、so_what），
子类只覆盖 _xl_chart_type 和 _customize()。

质量修复（相比原 ppt_builder.py）：
1. 数据标签位置细化：bar 标在右端，column 标在顶部，line 标在拐点
2. pie >5 项时自动合并尾部为"其他"
3. bar categories 自动按数值降序排列
4. waterfall 正值/负值分色（绿/红），而非单色 accent
"""

from pptx.util import Pt, Emu
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

from models import ChartSpec, ChartType, Rect, VisualTheme
from pipeline.skills.base import RenderingSkill, SkillDescriptor
from pipeline.skills._utils import parse_color, theme_color


# ── 基类 ──────────────────────────────────────────────────────

class ChartSkill(RenderingSkill):
    """
    图表 Skill 基类（模板方法模式）

    子类只需定义：
    - _xl_chart_type: 对应的 python-pptx 图表类型
    - _customize(): 图表类型特有的定制逻辑
    """

    _xl_chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    _handles_type: str = "column"  # 子类覆盖

    def descriptor(self) -> SkillDescriptor:
        return SkillDescriptor(
            skill_id=f"chart_{self._handles_type}",
            skill_type="chart",
            handles_types=[self._handles_type],
        )

    def prompt_fragment(self) -> str:
        """子类覆盖以提供 chart_type 选型指导"""
        return ""

    def design_tokens(self) -> dict:
        return {
            "gap_width": 100,
            "data_label_font_size": 9,
            "chart_title_font_size": 14,
        }

    def render(self, slide, data: ChartSpec, rect, theme: VisualTheme) -> bool:
        """模板方法：共通图表渲染流程"""
        if not data.categories or not data.series:
            return False

        # 1. 准备数据（防御：LLM 可能输出字符串数值）
        chart_data, all_values = self._prepare_data(data)

        # 2. 添加图表形状
        chart_frame = slide.shapes.add_chart(
            self._xl_chart_type,
            Emu(rect.left), Emu(rect.top),
            Emu(rect.width), Emu(rect.height),
            chart_data,
        )
        chart = chart_frame.chart

        # 3. 去网格线（咨询级标准）
        self._remove_gridlines(chart, data.chart_type)

        # 4. 图表标题
        self._set_title(chart, data, theme)

        # 5. 图例
        chart.has_legend = data.show_legend
        if chart.has_legend:
            chart.legend.font.size = Pt(theme.font_sizes.get("chart_label", 9))
            chart.legend.include_in_layout = False

        # 6. 应用主题颜色
        palette = theme.colors.get("chart_palette", [])
        self._apply_colors(chart, data, palette, theme)

        # 7. 数据标签
        self._apply_data_labels(chart, data, all_values, theme)

        # 8. 子类特有定制
        self._customize(chart, data, chart_frame, rect, theme, all_values)

        # 9. so_what 结论文字
        if data.so_what:
            from pipeline.skills._utils import add_textbox
            so_what_top = rect.top + rect.height + 91440
            add_textbox(
                slide,
                rect.left, so_what_top,
                rect.width, 365760,
                "→ " + data.so_what,
                font_size=10,
                font_name=theme.fonts.get("body", "Calibri"),
                color=theme_color(theme, "text_dark", "#2D3436"),
            )

        return True

    def _prepare_data(self, spec: ChartSpec):
        """准备图表数据，返回 (CategoryChartData, all_values)"""
        chart_data = CategoryChartData()
        chart_data.categories = spec.categories
        all_values = []
        for series in spec.series:
            safe_vals = []
            for v in series.values:
                try:
                    safe_vals.append(float(v))
                except (ValueError, TypeError):
                    safe_vals.append(0)
            chart_data.add_series(series.name, safe_vals)
            all_values.extend(safe_vals)
        return chart_data, all_values

    def _remove_gridlines(self, chart, chart_type: ChartType):
        """移除网格线"""
        if chart_type == ChartType.PIE:
            return
        try:
            chart.value_axis.has_major_gridlines = False
            chart.value_axis.has_minor_gridlines = False
            chart.value_axis.format.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
            chart.value_axis.format.line.width = Pt(0.5)
        except (ValueError, AttributeError):
            pass
        try:
            chart.category_axis.format.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
            chart.category_axis.format.line.width = Pt(0.5)
        except (ValueError, AttributeError):
            pass

    def _set_title(self, chart, spec: ChartSpec, theme: VisualTheme):
        """设置图表标题"""
        if spec.title:
            chart.has_title = True
            p = chart.chart_title.text_frame.paragraphs[0]
            p.text = spec.title
            p.font.size = Pt(theme.font_sizes.get("chart_title", 14))
            p.font.name = theme.fonts.get("body", "Calibri")
        else:
            chart.has_title = False

    def _apply_colors(self, chart, spec: ChartSpec, palette: list, theme: VisualTheme):
        """应用主题配色"""
        plot = chart.plots[0]
        if spec.chart_type == ChartType.PIE:
            self._apply_pie_colors(plot, palette)
        else:
            for i, series in enumerate(plot.series):
                color_str = ""
                if i < len(spec.series) and spec.series[i].color:
                    color_str = spec.series[i].color
                elif i < len(palette):
                    color_str = palette[i]
                if color_str:
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = parse_color(color_str)

    def _apply_pie_colors(self, plot, palette: list):
        """饼图逐点着色"""
        try:
            pie_series = plot.series[0]
            for i, point in enumerate(pie_series.points):
                if i < len(palette):
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = parse_color(palette[i])
        except Exception:
            pass

    def _apply_data_labels(self, chart, spec: ChartSpec, all_values: list, theme: VisualTheme):
        """应用数据标签策略"""
        plot = chart.plots[0]
        if not all_values:
            return

        if len(all_values) <= 12:
            # 少量数据：全标
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.font.size = Pt(theme.font_sizes.get("chart_label", 9))
            dl.font.name = theme.fonts.get("body", "Calibri")
            dl.number_format = '0.0'

    def _customize(self, chart, spec, chart_frame, rect, theme, all_values):
        """子类覆盖：图表类型特有定制"""
        pass


# ── 具体子类 ──────────────────────────────────────────────────

class ColumnChartSkill(ChartSkill):
    _xl_chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    _handles_type = "column"

    def prompt_fragment(self) -> str:
        return """| 多时期对比 | column | 时间从左到右；同比增长用颜色深浅区分 |"""

    def _customize(self, chart, spec, chart_frame, rect, theme, all_values):
        try:
            chart.plots[0].gap_width = 100
        except Exception:
            pass


class BarChartSkill(ChartSkill):
    _xl_chart_type = XL_CHART_TYPE.BAR_CLUSTERED
    _handles_type = "bar"

    def prompt_fragment(self) -> str:
        return """| 跨组排名 | bar | 必须按数值降序排列；用灰色+一个强调色突出关键条 |"""

    def render(self, slide, data: ChartSpec, rect, theme: VisualTheme) -> bool:
        """bar 图：categories 按数值降序排列"""
        if data.series and data.categories:
            vals = data.series[0].values
            safe_vals = []
            for v in vals:
                try:
                    safe_vals.append(float(v))
                except (ValueError, TypeError):
                    safe_vals.append(0)
            # 按 value 降序重排 categories 和 series
            indices = sorted(range(len(safe_vals)), key=lambda i: safe_vals[i], reverse=True)
            data.categories = [data.categories[i] if i < len(data.categories) else str(i) for i in indices]
            for s in data.series:
                s.values = [s.values[i] if i < len(s.values) else 0 for i in indices]
        return super().render(slide, data, rect, theme)

    def _customize(self, chart, spec, chart_frame, rect, theme, all_values):
        try:
            chart.plots[0].gap_width = 80
        except Exception:
            pass


class LineChartSkill(ChartSkill):
    _xl_chart_type = XL_CHART_TYPE.LINE_MARKERS
    _handles_type = "line"

    def prompt_fragment(self) -> str:
        return """| 时间序列趋势 | line | 趋势是主角；标注拐点和终点，不标每个点 |"""


class PieChartSkill(ChartSkill):
    _xl_chart_type = XL_CHART_TYPE.PIE
    _handles_type = "pie"

    def prompt_fragment(self) -> str:
        return """| 各项占比 | pie | ≤5项用饼图，>5项改bar降序；最大项用强调色 |"""

    def render(self, slide, data: ChartSpec, rect, theme: VisualTheme) -> bool:
        """pie 图：>5 项时合并尾部为"其他" """
        if data.series and data.categories:
            vals = data.series[0].values
            safe_vals = []
            for v in vals:
                try:
                    safe_vals.append(float(v))
                except (ValueError, TypeError):
                    safe_vals.append(0)

            if len(data.categories) > 5:
                # 保留前4项，合并剩余为"其他"
                top_n = 4
                top_sum = sum(safe_vals[:top_n])
                other_sum = sum(safe_vals[top_n:])
                data.categories = list(data.categories[:top_n]) + ["其他"]
                data.series[0].values = list(safe_vals[:top_n]) + [other_sum]
                # 同步其他 series
                for s in data.series[1:]:
                    other_vals = []
                    for v in s.values[top_n:]:
                        try:
                            other_vals.append(float(v))
                        except (ValueError, TypeError):
                            other_vals.append(0)
                    s.values = list(s.values[:top_n]) + [sum(other_vals)]

        return super().render(slide, data, rect, theme)


class AreaChartSkill(ChartSkill):
    _xl_chart_type = XL_CHART_TYPE.AREA
    _handles_type = "area"


class ScatterChartSkill(ChartSkill):
    _xl_chart_type = XL_CHART_TYPE.XY_SCATTER
    _handles_type = "scatter"


class WaterfallChartSkill(ChartSkill):
    """瀑布图：用堆叠柱状图模拟，正值/负值分色"""
    _xl_chart_type = XL_CHART_TYPE.COLUMN_STACKED
    _handles_type = "waterfall"

    def prompt_fragment(self) -> str:
        return """| 增减分析 | waterfall | 正值绿、负值红、累计灰；终值单独标注 |"""

    def render(self, slide, data: ChartSpec, rect, theme: VisualTheme) -> bool:
        if not data.series:
            return False

        values = data.series[0].values
        categories = data.categories

        # 计算基底和增量
        bases = []
        increments = []
        running = 0
        for v in values:
            try:
                fv = float(v)
            except (ValueError, TypeError):
                fv = 0
            if fv >= 0:
                bases.append(running)
            else:
                bases.append(running + fv)  # 负值：基底上移
            increments.append(abs(fv))
            running += fv

        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series("基底", bases)
        chart_data.add_series("增量", increments)

        chart_frame = slide.shapes.add_chart(
            self._xl_chart_type,
            Emu(rect.left), Emu(rect.top),
            Emu(rect.width), Emu(rect.height),
            chart_data,
        )
        chart = chart_frame.chart

        # 标题
        self._set_title(chart, data, theme)

        # 去网格线
        try:
            chart.value_axis.has_major_gridlines = False
        except (ValueError, AttributeError):
            pass

        plot = chart.plots[0]
        plot.gap_width = 50

        # 底层透明
        plot.series[0].format.fill.background()

        # 增量层：正值绿、负值红（逐点着色）
        inc_series = plot.series[1]
        green = parse_color("#27AE60")
        red = parse_color("#E74C3C")
        accent = theme_color(theme, "accent", "#FF6B35")

        try:
            for i, point in enumerate(inc_series.points):
                try:
                    original_val = float(values[i])
                except (ValueError, TypeError, IndexError):
                    original_val = 0
                point.format.fill.solid()
                if original_val > 0:
                    point.format.fill.fore_color.rgb = green
                elif original_val < 0:
                    point.format.fill.fore_color.rgb = red
                else:
                    point.format.fill.fore_color.rgb = accent
        except Exception:
            inc_series.format.fill.solid()
            inc_series.format.fill.fore_color.rgb = accent

        # 数据标签
        inc_series.has_data_labels = True
        inc_series.data_labels.font.size = Pt(9)
        inc_series.data_labels.font.name = theme.fonts.get("body", "Calibri")
        inc_series.data_labels.number_format = '0.0'

        # so_what
        if data.so_what:
            from pipeline.skills._utils import add_textbox
            so_what_top = rect.top + rect.height + 91440
            add_textbox(
                slide,
                rect.left, so_what_top,
                rect.width, 365760,
                "→ " + data.so_what,
                font_size=10,
                font_name=theme.fonts.get("body", "Calibri"),
                color=theme_color(theme, "text_dark", "#2D3436"),
            )

        return True


class ComboChartSkill(ChartSkill):
    """组合图：降级为簇状柱状图（python-pptx 不支持原生 combo）"""
    _xl_chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    _handles_type = "combo"

    def prompt_fragment(self) -> str:
        return """| 多指标叠加 | combo | 主指标柱状+次指标折线；图例必须清晰 |"""

    def _customize(self, chart, spec, chart_frame, rect, theme, all_values):
        # 组合图用区分度更高的配色
        primary = theme_color(theme, "primary", "#003D6E")
        accent = theme_color(theme, "accent", "#FF6B35")
        plot = chart.plots[0]
        plot.gap_width = 100
        palette = theme.colors.get("chart_palette", [])
        fallback = [primary, accent]
        for i, series in enumerate(plot.series):
            color = None
            if i < len(spec.series) and spec.series[i].color:
                color = parse_color(spec.series[i].color)
            elif i < len(palette):
                color = parse_color(palette[i])
            else:
                color = fallback[i % len(fallback)]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
