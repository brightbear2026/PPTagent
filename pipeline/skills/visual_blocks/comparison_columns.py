"""
Comparison Columns Skill

A vs B 对比栏：色块标题 + 逐条要点列表。

质量修复（相比原 ppt_builder.py）：
1. body 区域从 \n split 改为逐 item 结构化渲染
2. 增加行间距（space_after Pt(6)）
3. 标题色块增加圆角视觉效果
"""

from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from models import VisualBlock, VisualBlockItem, Rect
from pipeline.skills.base import RenderingSkill, SkillDescriptor
from pipeline.skills._utils import parse_color, theme_color, add_textbox


class ComparisonColumnsSkill(RenderingSkill):

    def descriptor(self) -> SkillDescriptor:
        return SkillDescriptor(
            skill_id="vb_comparison_columns",
            skill_type="visual_block",
            handles_types=["comparison_columns"],
            content_pattern="two_column",
        )

    def prompt_fragment(self) -> str:
        return """**comparison_columns**（A vs B 对比）
  items字段: title, description
  设计理念: 对比要鲜明——对比色区分两列，标题栏突出差异核心
  质量要求: 每列title是一个立场/选项名，description分点列出关键差异"""

    def design_tokens(self) -> dict:
        return {
            "max_items": 4,
            "header_height_inches": 0.4,
            "column_gap_inches": 0.1,
            "body_font_size": 11,
            "header_font_size": 13,
        }

    def render(self, slide, data: VisualBlock, rect, theme) -> bool:
        items = data.items[:4]
        n = len(items)
        if n == 0:
            return False

        tokens = self.design_tokens()

        # 合并 rect（list[Rect] slot 模式）
        if isinstance(rect, list):
            from models import Rect as R
            rect = R(
                left=rect[0].left,
                top=rect[0].top,
                width=rect[-1].left + rect[-1].width - rect[0].left,
                height=max(r.height for r in rect),
            )

        gap = Emu(int(tokens["column_gap_inches"] * 914400))
        col_w = (rect.width - gap * (n - 1)) // n
        header_h = Emu(int(tokens["header_height_inches"] * 914400))

        primary = theme_color(theme, "primary", "#003D6E")
        accent = theme_color(theme, "accent", "#FF6B35")
        text_dark = theme_color(theme, "text_dark", "#2D3436")
        body_font = theme.fonts.get("body", "Calibri")

        # 对比色方案（从 theme 扩展）
        col_colors = [primary, accent,
                      parse_color("#00A878"),
                      parse_color("#6C5CE7")]

        for i, item in enumerate(items):
            left = rect.left + i * (col_w + gap)
            col_color = col_colors[i % len(col_colors)]

            # 标题色块（圆角矩形）
            header = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Emu(left), Emu(rect.top),
                Emu(col_w), Emu(header_h),
            )
            header.fill.solid()
            header.fill.fore_color.rgb = col_color
            header.line.fill.background()
            tf = header.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item.title or f"选项{i+1}"
            p.font.size = Pt(tokens["header_font_size"])
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.name = body_font
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 内容区域 — 逐行渲染而非 \n split
            body_top = rect.top + header_h + Emu(45720)
            body_h = rect.height - header_h - Emu(45720)
            padding = Emu(45720)

            # 如果 description 中有换行，逐条渲染为子弹点
            desc = item.description or ""
            lines = [l.strip() for l in desc.replace("\\n", "\n").split("\n") if l.strip()]

            if lines:
                body_box = slide.shapes.add_textbox(
                    Emu(left + padding), Emu(body_top),
                    Emu(col_w - padding * 2), Emu(body_h),
                )
                tf2 = body_box.text_frame
                tf2.word_wrap = True

                for j, line in enumerate(lines):
                    p2 = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
                    bullet = "" if line.startswith("•") else "• "
                    p2.text = bullet + line
                    p2.font.size = Pt(tokens["body_font_size"])
                    p2.font.color.rgb = text_dark
                    p2.font.name = body_font
                    p2.space_after = Pt(6)  # 增加行间距
            else:
                # 无分条内容，整体渲染
                add_textbox(
                    slide,
                    left + padding, body_top,
                    col_w - padding * 2, body_h,
                    desc,
                    font_size=tokens["body_font_size"],
                    font_name=body_font,
                    color=text_dark,
                )

        return True
