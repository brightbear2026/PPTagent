"""
Icon Text Grid Skill

4-6个并列要点：圆形色块图标（首字缩写）+ 标题 + 描述。

质量修复（相比原 ppt_builder.py）：
1. 图标从随机几何体统一为圆形色块 + 首字缩写（更专业）
2. 描述字号 Pt(9) → Pt(10)（提升可读性）
3. 增加卡片底色区分
"""

from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from models import VisualBlock, VisualBlockItem, Rect
from pipeline.skills.base import RenderingSkill, SkillDescriptor
from pipeline.skills._utils import (
    parse_color, theme_color, add_textbox, add_centered_text,
)


class IconTextGridSkill(RenderingSkill):

    def descriptor(self) -> SkillDescriptor:
        return SkillDescriptor(
            skill_id="vb_icon_text_grid",
            skill_type="visual_block",
            handles_types=["icon_text_grid"],
            content_pattern="icon_grid",
        )

    def prompt_fragment(self) -> str:
        return """**icon_text_grid**（4-6个并列要点）
  items字段: title, description
  设计理念: 信息平铺——每格一个核心观点，避免长段落
  质量要求: 所有title等长（字数相近），description控制在15字以内"""

    def design_tokens(self) -> dict:
        return {
            "max_items": 6,
            "icon_size_inches": 0.35,
            "desc_font_size": 10,
            "title_font_size": 11,
            "card_bg_key": "surface",
        }

    def render(self, slide, data: VisualBlock, rect, theme) -> bool:
        items = data.items[:6]
        n = len(items)
        if n == 0:
            return False

        tokens = self.design_tokens()

        # 合并 rect（list[Rect] slot 模式）
        if isinstance(rect, list):
            from models import Rect as R
            rect = R(
                left=rect[0].left,
                top=rect[0].top,
                width=rect[-1].left + rect[-1].width - rect[0].left,
                height=max(r.height for r in rect),
            )

        cols = data.columns if data.columns > 0 else min(3, n)
        rows = (n + cols - 1) // cols
        gap = Emu(91440)  # 0.1 inch
        cell_w = (rect.width - gap * (cols - 1)) // cols
        cell_h = (rect.height - gap * (rows - 1)) // rows

        icon_size = Emu(int(tokens["icon_size_inches"] * 914400))
        primary = theme_color(theme, "primary", "#003D6E")
        text_light = theme_color(theme, "text_light", "#636E72")
        card_bg = theme_color(theme, tokens["card_bg_key"], "#F0F4F8")
        body_font = theme.fonts.get("body", "Calibri")

        palette = [
            parse_color(c) for c in
            theme.colors.get(
                "chart_palette",
                ["#003D6E", "#FF6B35", "#00A878", "#E17055", "#6C5CE7", "#FDCB6E"],
            )
        ]

        for idx, item in enumerate(items):
            row = idx // cols
            col = idx % cols
            cell_left = rect.left + col * (cell_w + gap)
            cell_top = rect.top + row * (cell_h + gap)

            # 卡片背景（增加底色区分）
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Emu(cell_left), Emu(cell_top),
                Emu(cell_w), Emu(cell_h),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = card_bg
            card.line.fill.background()

            # 图标 — 统一为圆形色块 + 首字缩写
            icon_left = cell_left + (cell_w - icon_size) // 2
            icon = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Emu(icon_left), Emu(cell_top + Emu(91440)),
                Emu(icon_size), Emu(icon_size),
            )
            icon.fill.solid()
            icon.fill.fore_color.rgb = palette[idx % len(palette)]
            icon.line.fill.background()

            # 首字缩写
            tf = icon.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            title_text = item.title or ""
            p.text = title_text[0] if title_text else str(idx + 1)
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.name = body_font
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 标题
            title_top = cell_top + icon_size + Emu(137160)
            add_centered_text(
                slide,
                cell_left, title_top,
                cell_w, Emu(228600),
                title_text,
                font_size=tokens["title_font_size"],
                font_name=body_font,
                color=primary,
                bold=True,
            )

            # 描述
            if item.description:
                desc_top = int(title_top) + Emu(228600)
                add_centered_text(
                    slide,
                    cell_left, desc_top,
                    cell_w, cell_h - icon_size - Emu(365760),
                    item.description,
                    font_size=tokens["desc_font_size"],
                    font_name=body_font,
                    color=text_light,
                )

        return True
