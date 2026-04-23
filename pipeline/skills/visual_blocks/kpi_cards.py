"""
KPI Cards Skill

2-4个关键指标卡片，圆角矩形底色 + 自适应字号数字 + 标签 + 趋势箭头。

质量修复（相比原 ppt_builder.py）：
1. 背景色从硬编码 #F0F4F8 改为 theme.colors["surface"]
2. value 字号根据字符数动态调整（不再固定 Pt(36)）
3. 卡片间距根据卡片数量自适应
"""

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from models import VisualBlock, VisualBlockItem, Rect
from pipeline.skills.base import RenderingSkill, SkillDescriptor
from pipeline.skills._utils import (
    parse_color, theme_color, fit_font_size,
    add_centered_text, TREND_ARROWS, TREND_COLORS,
)


class KpiCardsSkill(RenderingSkill):

    def descriptor(self) -> SkillDescriptor:
        return SkillDescriptor(
            skill_id="vb_kpi_cards",
            skill_type="visual_block",
            handles_types=["kpi_cards"],
            content_pattern="kpi_highlight",
        )

    def prompt_fragment(self) -> str:
        return """**kpi_cards**（2-4个关键指标）
  items字段: title, value, description, trend(up/down/flat)
  设计理念: 少即是多——大字号突出数值，小字号标注含义
  质量要求: value必须简短（"32%"、"1.56亿"），不写成句子；4个KPI之间要有逻辑关联，不是随机拼凑"""

    def design_tokens(self) -> dict:
        return {
            "max_items": 4,
            "value_font_range": (18, 44),   # min/max pt，根据字符数自适应
            "bg_color_key": "surface",       # 从 theme 读取，不硬编码
            "card_max_height_inches": 3.0,
            "card_gap_inches": 0.1,
        }

    def render(self, slide, data: VisualBlock, rect, theme) -> bool:
        """
        渲染 KPI 卡片

        支持两种 rect 模式：
        - list[Rect]: 每个 item 有独立 slot
        - Rect: 单区域，内部分割为 n 个卡片
        """
        items = data.items[:4]
        n = len(items)
        if n == 0:
            return False

        tokens = self.design_tokens()

        # 颜色：从 theme 读取，不再硬编码
        primary = theme_color(theme, "primary", "#003D6E")
        bg_color = theme_color(theme, tokens["bg_color_key"], "#F0F4F8")
        text_light = theme_color(theme, "text_light", "#636E72")
        title_font = theme.fonts.get("title", "Arial")
        body_font = theme.fonts.get("body", "Calibri")

        if isinstance(rect, list):
            # slot 模式
            for i, item in enumerate(items):
                if i >= len(rect):
                    break
                self._render_single(slide, item, rect[i], theme,
                                    tokens, primary, bg_color, text_light,
                                    title_font, body_font)
        else:
            # 单区域分割模式
            gap = Emu(int(tokens["card_gap_inches"] * 914400))
            card_w = (rect.width - gap * (n - 1)) // n
            card_h = min(rect.height, Emu(int(tokens["card_max_height_inches"] * 914400)))

            for i, item in enumerate(items):
                left = rect.left + i * (card_w + gap)
                from models import Rect as R
                card_rect = R(left, rect.top, card_w, card_h)
                self._render_single(slide, item, card_rect, theme,
                                    tokens, primary, bg_color, text_light,
                                    title_font, body_font)
        return True

    def _render_single(self, slide, item: VisualBlockItem, rect: Rect,
                       theme, tokens, primary, bg_color, text_light,
                       title_font, body_font):
        """渲染单个 KPI 卡片"""
        # 背景卡片
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(rect.left), Emu(rect.top),
            Emu(rect.width), Emu(rect.height),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.fill.background()

        # 顶部装饰条（主题色，2px 高）
        accent = theme_color(theme, "accent", "#FF6B35")
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(rect.left), Emu(rect.top),
            Emu(rect.width), Emu(18000),  # ~2px
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        # 大数字 — 字号根据字符数自适应
        value_text = item.value or "—"
        min_pt, max_pt = tokens["value_font_range"]
        font_size = fit_font_size(value_text, max_pt, min_pt)

        val_h = rect.height * 42 // 100
        add_centered_text(
            slide,
            rect.left + rect.width // 10,
            rect.top + rect.height // 8,
            rect.width * 8 // 10,
            val_h,
            value_text,
            font_size=font_size,
            font_name=title_font,
            color=primary,
            bold=True,
        )

        # 标签
        label_top = rect.top + rect.height // 8 + val_h
        add_centered_text(
            slide,
            rect.left + rect.width // 10,
            label_top,
            rect.width * 8 // 10,
            rect.height // 5,
            item.title or "",
            font_size=11,
            font_name=body_font,
            color=text_light,
        )

        # 趋势描述
        if item.description:
            desc_top = label_top + rect.height // 5
            trend = item.trend or "flat"
            arrow = TREND_ARROWS.get(trend, "")
            trend_color = parse_color(TREND_COLORS.get(trend, "#003D6E"))

            add_centered_text(
                slide,
                rect.left + rect.width // 10,
                desc_top,
                rect.width * 8 // 10,
                rect.height // 5,
                arrow + item.description,
                font_size=10,
                font_name=body_font,
                color=trend_color,
            )
