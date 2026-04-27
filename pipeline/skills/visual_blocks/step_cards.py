"""
Step Cards Skill

3-6个步骤卡片：编号圆圈 + 标题 + 描述 + 箭头连接。

质量修复（相比原 ppt_builder.py）：
1. 描述字号 Pt(9) → Pt(11)（原来太小不可读）
2. 箭头颜色从硬编码 #D0D5DD 改为跟随主题
3. 编号圆圈增加阴影效果（通过半透明偏移圆模拟）
"""

from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from models import VisualBlock, VisualBlockItem, Rect
from pipeline.skills.base import RenderingSkill, SkillDescriptor
from pipeline.skills._utils import theme_color, add_centered_text, parse_color


class StepCardsSkill(RenderingSkill):

    def descriptor(self) -> SkillDescriptor:
        return SkillDescriptor(
            skill_id="vb_step_cards",
            skill_type="visual_block",
            handles_types=["step_cards"],
            content_pattern="step_flow",
        )

    def prompt_fragment(self) -> str:
        return """**step_cards**（3-6个步骤）
  items字段: title, description
  设计理念: 步骤间视觉间隔清晰，编号+箭头引导视线
  质量要求: title ≤6字，description ≤20字，步骤之间有因果或时序关系
  反模式: 不要用于无顺序关系的并列要点（用icon_text_grid）；不要用于仅2个步骤（用comparison_columns）"""

    def design_tokens(self) -> dict:
        return {
            "max_items": 6,
            "circle_size_inches": 0.4,
            "desc_font_size": 11,        # 从 Pt(9) 提升到 Pt(11)
            "title_font_size": 12,
            "step_gap_inches": 0.2,      # 步骤间距（含箭头空间）
            "arrow_color_key": "border",  # 从 theme 读取，不硬编码
        }

    def render(self, slide, data: VisualBlock, rect, theme) -> bool:
        items = data.items[:6]
        n = len(items)
        if n == 0:
            return False

        tokens = self.design_tokens()
        gap = Emu(int(tokens["step_gap_inches"] * 914400))
        circle_size = Emu(int(tokens["circle_size_inches"] * 914400))

        # 合并 rect（如果是 list[Rect] slot 模式）
        if isinstance(rect, list):
            from models import Rect as R
            rect = R(
                left=rect[0].left,
                top=rect[0].top,
                width=rect[-1].left + rect[-1].width - rect[0].left,
                height=max(r.height for r in rect),
            )

        card_w = (rect.width - gap * (n - 1)) // n
        accent = theme_color(theme, "accent", "#FF6B35")
        primary = theme_color(theme, "primary", "#003D6E")
        text_light = theme_color(theme, "text_light", "#636E72")
        arrow_color = theme_color(theme, tokens["arrow_color_key"], "#D0D5DD")
        title_font = theme.fonts.get("body", "Calibri")

        for i, item in enumerate(items):
            left = rect.left + i * (card_w + gap)

            # 编号圆圈阴影（偏移 2px 的半透明圆）
            shadow_offset = 18000  # ~2px
            shadow = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Emu(left + (card_w - circle_size) // 2 + shadow_offset),
                Emu(rect.top + shadow_offset),
                Emu(circle_size), Emu(circle_size),
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            shadow.line.fill.background()

            # 编号圆圈
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Emu(left + (card_w - circle_size) // 2),
                Emu(rect.top),
                Emu(circle_size), Emu(circle_size),
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = accent
            circle.line.fill.background()

            tf = circle.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.name = title_font
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 标题
            title_top = rect.top + circle_size + Emu(91440)
            add_centered_text(
                slide,
                left, title_top,
                card_w, Emu(274320),
                item.title or f"步骤{i+1}",
                font_size=tokens["title_font_size"],
                font_name=title_font,
                color=primary,
                bold=True,
            )

            # 描述 — Pt(11) 替代原来的 Pt(9)
            if item.description:
                desc_top = int(title_top) + Emu(320040)
                add_centered_text(
                    slide,
                    left, desc_top,
                    card_w, rect.height - (desc_top - rect.top),
                    item.description,
                    font_size=tokens["desc_font_size"],
                    font_name=title_font,
                    color=text_light,
                )

            # 箭头（非最后一个）— 颜色跟随主题
            if i < n - 1:
                arrow_left = left + card_w + Emu(27432)
                arrow_top = rect.top + circle_size // 2 - Emu(54864)
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Emu(arrow_left), Emu(arrow_top),
                    Emu(gap - Emu(54864)), Emu(109728),
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = arrow_color
                arrow.line.fill.background()

        return True
