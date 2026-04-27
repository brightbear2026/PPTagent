"""
Callout Box Skill

关键洞察/金句引用框：左侧色条 + 大引号装饰 + 引用文字。

质量修复（相比原 ppt_builder.py）：
1. 引号字体跟随主题（中文场景 Georgia 引号不好看 → 用主题字体）
2. 色条增加圆角顶部装饰效果
3. 引用文字增加左右 padding
"""

from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE

from models import VisualBlock, VisualBlockItem, Rect
from pipeline.skills.base import RenderingSkill, SkillDescriptor
from pipeline.skills._utils import theme_color, parse_color, add_textbox


class CalloutBoxSkill(RenderingSkill):

    def descriptor(self) -> SkillDescriptor:
        return SkillDescriptor(
            skill_id="vb_callout_box",
            skill_type="visual_block",
            handles_types=["callout_box"],
            content_pattern="stat_callout",
        )

    def prompt_fragment(self) -> str:
        return """**callout_box**（关键洞察/金句）
  items字段: title, description
  设计理念: 引用的力量——大引号+斜体+色条，让洞察成为视觉焦点
  质量要求: description是一个完整的洞察句，不是标题短语
  反模式: 不要用于列举多个要点；不要用于数据展示（用kpi_cards或chart）；仅用于单条洞察/金句"""

    def design_tokens(self) -> dict:
        return {
            "bar_width_emu": 54864,       # ~6pt 色条宽度
            "quote_font_size": 48,
            "text_font_size": 14,
            "source_font_size": 10,
            "padding_left_emu": 182880,   # 引用文字左边距
            "padding_right_emu": 228600,  # 引用文字右边距
        }

    def render(self, slide, data: VisualBlock, rect, theme) -> bool:
        if isinstance(rect, list):
            rect = rect[0] if rect else None
        if rect is None:
            return False

        item = data.items[0] if data.items else VisualBlockItem()
        tokens = self.design_tokens()

        accent = theme_color(theme, "accent", "#FF6B35")
        primary = theme_color(theme, "primary", "#003D6E")
        text_light = theme_color(theme, "text_light", "#636E72")
        body_font = theme.fonts.get("body", "Calibri")

        # 左侧色条
        bar_w = tokens["bar_width_emu"]
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(rect.left), Emu(rect.top),
            Emu(bar_w), Emu(rect.height),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        # 顶部装饰圆角块（色条顶部的视觉加强）
        cap_w = bar_w * 3
        cap_h = Emu(27432)  # ~3px
        cap = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(rect.left), Emu(rect.top),
            Emu(cap_w), Emu(cap_h),
        )
        cap.fill.solid()
        cap.fill.fore_color.rgb = accent
        cap.line.fill.background()

        # 大引号 — 使用主题字体（不再硬编码 Georgia）
        quote_box = slide.shapes.add_textbox(
            Emu(rect.left + Emu(137160)), Emu(rect.top),
            Emu(365760), Emu(365760),
        )
        tf = quote_box.text_frame
        p = tf.paragraphs[0]
        p.text = "\u201C"
        p.font.size = Pt(tokens["quote_font_size"])
        p.font.color.rgb = accent
        p.font.name = body_font  # 跟随主题字体

        # 引用正文
        text_left = rect.left + tokens["padding_left_emu"]
        text_top = rect.top + Emu(320040)
        text_w = rect.width - tokens["padding_right_emu"] - tokens["padding_left_emu"]
        add_textbox(
            slide,
            text_left, text_top,
            text_w, rect.height * 50 // 100,
            item.description or item.title or "",
            font_size=tokens["text_font_size"],
            font_name=body_font,
            color=primary,
        )

        # 正文设为斜体（需要修改刚创建的段落）
        try:
            slide.shapes[-1].text_frame.paragraphs[0].font.italic = True
        except (IndexError, AttributeError):
            pass

        # 出处/标题
        if item.title and item.description:
            source_top = int(text_top) + rect.height * 50 // 100 + Emu(45720)
            add_textbox(
                slide,
                text_left, source_top,
                text_w, Emu(228600),
                "— " + item.title,
                font_size=tokens["source_font_size"],
                font_name=body_font,
                color=text_light,
            )

        return True
