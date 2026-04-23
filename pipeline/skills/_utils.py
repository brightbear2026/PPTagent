"""
Skill 共享渲染工具函数

从 ppt_builder.py 提取的通用渲染操作，
供各 Skill 的 render() 方法复用。
"""

from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ── 常量 ──────────────────────────────────────────────────────

SLIDE_WIDTH = 12192000   # 13.333 inches in EMU (16:9)
SLIDE_HEIGHT = 6858000   # 7.5 inches in EMU (16:9)
MIN_DIM = 91440          # ~0.1 inch 最小尺寸


# ── 颜色工具 ──────────────────────────────────────────────────

def parse_color(color_str: str) -> RGBColor:
    """解析颜色字符串为 RGBColor"""
    if not color_str:
        return RGBColor(0, 0, 0)
    hex_str = color_str.lstrip("#")
    if len(hex_str) == 6:
        return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
    return RGBColor(0, 0, 0)


def theme_color(theme, key: str, default: str = "#000000") -> RGBColor:
    """从 theme.colors 读取颜色，带默认值"""
    return parse_color(theme.colors.get(key, default))


# ── 字号自适应 ────────────────────────────────────────────────

def fit_font_size(text: str, max_pt: int = 44, min_pt: int = 18) -> int:
    """
    根据文本长度自适应字号。

    短文本（≤4字符）用大字号，长文本逐步缩小。
    适用于 KPI value、stat_highlight 等大字号场景。
    """
    length = len(text)
    if length <= 4:
        return max_pt
    if length <= 8:
        return int(max_pt * 0.85)
    if length <= 12:
        return int(max_pt * 0.7)
    if length <= 16:
        return int(max_pt * 0.6)
    return min_pt


# ── 文本框快捷创建 ───────────────────────────────────────────

def add_textbox(slide, left, top, width, height, text: str,
                font_size: int = 12, font_name: str = "Calibri",
                color: RGBColor = None, bold: bool = False,
                alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                word_wrap: bool = True):
    """
    快速创建带样式的文本框

    Returns:
        TextFrame 对象（可进一步自定义段落）
    """
    color = color or RGBColor(0, 0, 0)
    txBox = slide.shapes.add_textbox(
        Emu(left), Emu(top), Emu(width), Emu(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment

    # 设置垂直锚点（需通过 XML 操作）
    try:
        from pptx.oxml.ns import qn
        txBox.text_frame._txBody.bodyPr.set("anchor", {
            MSO_ANCHOR.TOP: "t",
            MSO_ANCHOR.MIDDLE: "ctr",
            MSO_ANCHOR.BOTTOM: "b",
        }.get(anchor, "t"))
    except Exception:
        pass

    return tf


def add_centered_text(slide, left, top, width, height, text: str,
                      font_size: int = 12, font_name: str = "Calibri",
                      color: RGBColor = None, bold: bool = False):
    """居中文本框的快捷方法"""
    return add_textbox(
        slide, left, top, width, height, text,
        font_size=font_size, font_name=font_name,
        color=color, bold=bold, alignment=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


# ── 矩形裁剪 ─────────────────────────────────────────────────

def clamp_rect(rect, slide_width=SLIDE_WIDTH, slide_height=SLIDE_HEIGHT):
    """将坐标裁剪到幻灯片画布内"""
    left = max(0, min(rect.left, slide_width - MIN_DIM))
    top = max(0, min(rect.top, slide_height - MIN_DIM))
    width = min(rect.width, slide_width - left)
    height = min(rect.height, slide_height - top)
    from models import Rect
    return Rect(left, top, max(width, MIN_DIM), max(height, MIN_DIM))


# ── 趋势箭头 ─────────────────────────────────────────────────

TREND_ARROWS = {
    "up": "\u25B2 ",     # ▲
    "down": "\u25BC ",   # ▼
    "flat": "",          # 无箭头
}

TREND_COLORS = {
    "up": "#27AE60",     # 绿色
    "down": "#E74C3C",   # 红色
}
