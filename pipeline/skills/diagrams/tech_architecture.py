"""
Tech Architecture Diagram Skill

技术分层架构图：IT 语义色板 + 协议标签。
与通用 architecture skill 的区别：层内 items 含具体技术组件名（Nginx/Kafka/MySQL 等），
使用 IT 语义色板而非主题色轮转。
"""

from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from models import DiagramSpec, Rect, VisualTheme
from pipeline.skills.base import RenderingSkill, SkillDescriptor
from pipeline.skills._utils import theme_color


# IT 语义色板：按层级类型分配颜色
_IT_LAYER_COLORS = {
    "接入": "#5B9BD5",
    "网关": "#5B9BD5",
    "前端": "#5B9BD5",
    "展示": "#5B9BD5",
    "业务": "#70AD47",
    "应用": "#70AD47",
    "服务": "#70AD47",
    "逻辑": "#70AD47",
    "数据": "#FFC000",
    "存储": "#FFC000",
    "缓存": "#FFC000",
    "基础设施": "#A5A5A5",
    "infra": "#A5A5A5",
    "运维": "#A5A5A5",
    "外部": "#7030A0",
    "第三方": "#7030A0",
    "安全": "#C00000",
    "认证": "#C00000",
}
_DEFAULT_IT_COLOR = "#48A9E6"


class TechArchitectureSkill(RenderingSkill):

    def descriptor(self) -> SkillDescriptor:
        return SkillDescriptor(
            skill_id="diagram_tech_architecture",
            skill_type="diagram",
            handles_types=["tech_architecture"],
        )

    def prompt_fragment(self) -> str:
        return """**tech_architecture**（技术分层架构图）:
  约束: 3-7层，每层items ≤6个；层按从底到顶排列；protocols标注层间通信协议
  判别: 如果层内items含具体技术组件名（Nginx/Kafka/MySQL/Kubernetes/Docker/Redis等），用 tech_architecture；如果层内是业务概念/组织架构/抽象分类，用 architecture
  反模式: 不要用于微服务调用关系（用component_topology）；不要用于流程步骤（用process_flow）"""

    def design_tokens(self) -> dict:
        return {
            "separator_height_emu": 18288,
            "separator_color": "#FFFFFF",
            "label_bg_width_emu": 1828800,
            "label_bg_height_emu": 228600,
        }

    def render(self, slide, data: DiagramSpec, rect, theme: VisualTheme) -> bool:
        nodes = data.nodes
        if not nodes:
            return False

        tokens = self.design_tokens()
        white = RGBColor(255, 255, 255)
        font = theme.fonts.get("body", "Calibri")

        # 按 group 分层
        layers = {}
        layer_colors_map = {}
        for node in nodes:
            group = node.group or "default"
            layers.setdefault(group, []).append(node)
            # 尝试匹配 IT 语义色板
            if group not in layer_colors_map:
                layer_colors_map[group] = self._match_layer_color(group)

        layer_names = list(layers.keys())
        num_layers = len(layer_names)
        if num_layers == 0:
            return False

        layer_h = rect.height // num_layers

        for li, layer_name in enumerate(layer_names):
            layer_nodes = layers[layer_name]
            y = rect.top + li * layer_h
            hex_color = layer_colors_map[layer_name]
            color = self._hex_to_rgb(hex_color)

            # 层背景
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(rect.left), Emu(y),
                Emu(rect.width), Emu(layer_h - 45720),
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = color
            bg.line.fill.background()

            # 层标签
            label_w = tokens["label_bg_width_emu"]
            label_h = tokens["label_bg_height_emu"]
            label_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Emu(rect.left + 45720), Emu(y + 45720),
                Emu(label_w), Emu(label_h),
            )
            label_bg.fill.solid()
            label_bg.fill.fore_color.rgb = RGBColor(
                max(0, color[0] - 40), max(0, color[1] - 40), max(0, color[2] - 40),
            )
            label_bg.line.fill.background()
            tf_label = label_bg.text_frame
            tf_label.word_wrap = True
            p_label = tf_label.paragraphs[0]
            p_label.text = layer_name
            p_label.font.size = Pt(9)
            p_label.font.color.rgb = white
            p_label.font.bold = True
            p_label.font.name = font
            p_label.alignment = PP_ALIGN.CENTER

            # 分隔线
            if li > 0:
                sep_h = tokens["separator_height_emu"]
                sep = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Emu(rect.left), Emu(y - sep_h // 2),
                    Emu(rect.width), Emu(sep_h),
                )
                sep.fill.solid()
                sep.fill.fore_color.rgb = white
                sep.line.fill.background()

            # 节点
            n_items = len(layer_nodes)
            item_w = min(rect.width // max(n_items, 1) - 91440, 2743200)
            item_h = min(layer_h // 2, 457200)
            total_w = n_items * item_w + (n_items - 1) * 91440
            start_x = rect.left + (rect.width - total_w) // 2

            for ni, node in enumerate(layer_nodes):
                nx = start_x + ni * (item_w + 91440)
                ny = y + (layer_h - item_h) // 2 + label_h // 2
                self._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                                nx, ny, item_w, item_h,
                                white, node.label, color, font, 10)

        return True

    @staticmethod
    def _match_layer_color(group_name: str) -> str:
        """根据层名匹配 IT 语义色板"""
        g = group_name.lower()
        for keyword, color in _IT_LAYER_COLORS.items():
            if keyword in g:
                return color
        return _DEFAULT_IT_COLOR

    @staticmethod
    def _hex_to_rgb(hex_color: str):
        hex_color = hex_color.lstrip("#")
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

    @staticmethod
    def _add_shape(slide, shape_type, x, y, w, h, fill_color, text, text_color, font_name, font_size):
        shape = slide.shapes.add_shape(shape_type, Emu(x), Emu(y), Emu(w), Emu(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.bold = True
        p.font.name = font_name
        p.alignment = PP_ALIGN.CENTER
