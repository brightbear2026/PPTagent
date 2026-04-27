"""
Data Flow Diagram Skill

数据流管线图：水平 pipeline，stage 按类型区分形状。
source→transform→store→consume 线性流转。
"""

from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from models import DiagramSpec, Rect, VisualTheme
from pipeline.skills.base import RenderingSkill, SkillDescriptor
from pipeline.skills._utils import theme_color

# 按 stage type 分配颜色
_STAGE_COLORS = {
    "source": "#5B9BD5",
    "transform": "#FF6B35",
    "store": "#FFC000",
    "consume": "#70AD47",
}
_DEFAULT_STAGE_COLOR = "#48A9E6"


class DataFlowSkill(RenderingSkill):

    def descriptor(self) -> SkillDescriptor:
        return SkillDescriptor(
            skill_id="diagram_data_flow",
            skill_type="diagram",
            handles_types=["data_flow"],
        )

    def prompt_fragment(self) -> str:
        return """**data_flow**（数据流管线图）:
  约束: 3-8个stage，每个stage的type必须是 source/transform/store/consume 之一；flows标注数据量/格式
  反模式: 不要用于业务流程（用process_flow）；不要用于组件调用关系（用component_topology）"""

    def design_tokens(self) -> dict:
        return {
            "node_font_size": 10,
            "flow_label_font_size": 8,
            "arrow_spacing_emu": 137160,
        }

    def render(self, slide, data: DiagramSpec, rect, theme: VisualTheme) -> bool:
        nodes = data.nodes
        if not nodes or len(nodes) < 2:
            return False

        tokens = self.design_tokens()
        primary = theme_color(theme, "primary", "#003D6E")
        accent = theme_color(theme, "accent", "#FF6B35")
        white = RGBColor(255, 255, 255)
        muted = RGBColor(0x66, 0x66, 0x66)
        font = theme.fonts.get("body", "Calibri")

        n = min(len(nodes), 8)
        # 水平均匀分布
        total_w = rect.width - 182880 * 2
        spacing = total_w // n
        node_w = min(spacing - 91440, 1600200)
        node_h = 640080

        # 按 sublabel 或 label 判断 stage type（简化：用 node.group 或默认）
        positions = {}
        for i, node in enumerate(nodes[:n]):
            x = rect.left + 91440 + i * spacing + (spacing - node_w) // 2
            y = rect.top + (rect.height - node_h) // 2

            stage_type = (node.sublabel or "").lower()
            color_hex = _STAGE_COLORS.get(stage_type, _DEFAULT_STAGE_COLOR)
            color = RGBColor(int(color_hex[1:3], 16), int(color_hex[3:5], 16), int(color_hex[5:7], 16))

            # 根据类型选形状
            shape = MSO_SHAPE.ROUNDED_RECTANGLE
            if stage_type == "source":
                shape = MSO_SHAPE.PARALLELOGRAM
            elif stage_type == "store":
                shape = MSO_SHAPE.ROUNDED_RECTANGLE  # 圆角矩形近似存储

            box = slide.shapes.add_shape(shape, Emu(x), Emu(y), Emu(node_w), Emu(node_h))
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.fill.background()
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = node.label
            p.font.size = Pt(tokens["node_font_size"])
            p.font.color.rgb = white
            p.font.bold = True
            p.font.name = font
            p.alignment = PP_ALIGN.CENTER

            # stage type 小标签
            if stage_type:
                type_box = slide.shapes.add_textbox(
                    Emu(x), Emu(y + node_h + 45720),
                    Emu(node_w), Emu(182880),
                )
                tt = type_box.text_frame
                tp = tt.paragraphs[0]
                tp.text = stage_type.upper()
                tp.font.size = Pt(7)
                tp.font.color.rgb = muted
                tp.font.name = font
                tp.alignment = PP_ALIGN.CENTER

            positions[node.node_id] = (x, y, node_w, node_h, color)

        # 连线
        for i in range(n - 1):
            node_a = nodes[i]
            node_b = nodes[i + 1]
            if node_a.node_id in positions and node_b.node_id in positions:
                ax, ay, aw, _, _ = positions[node_a.node_id]
                bx, by, bw, _, _ = positions[node_b.node_id]

                connector = slide.shapes.add_connector(
                    1,
                    Emu(ax + aw), Emu(ay + node_h // 2),
                    Emu(bx), Emu(by + node_h // 2),
                )
                connector.line.color.rgb = muted
                connector.line.width = Pt(1.5)

                # flow label（查找匹配的 edge）
                edge_label = ""
                for edge in data.edges:
                    if edge.from_id == node_a.node_id and edge.to_id == node_b.node_id:
                        edge_label = edge.label
                        break
                if edge_label:
                    mx = (ax + aw + bx) // 2
                    lb = slide.shapes.add_textbox(
                        Emu(mx - 685800), Emu(ay + node_h // 2 - 365760),
                        Emu(1371600), Emu(228600),
                    )
                    lt = lb.text_frame
                    lp = lt.paragraphs[0]
                    lp.text = edge_label
                    lp.font.size = Pt(tokens["flow_label_font_size"])
                    lp.font.color.rgb = accent
                    lp.font.name = font
                    lp.alignment = PP_ALIGN.CENTER

        return True
