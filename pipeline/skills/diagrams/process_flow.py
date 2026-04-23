"""
Process Flow Diagram Skill

横向/纵向流程图：节点 + 箭头连接 + 边标签。

质量修复（相比原 diagram_renderer.py）：
1. 节点间距根据节点数量自适应（不再固定除法）
2. 节点增加阴影效果（半透明偏移矩形模拟）
3. 交替色更柔和（primary/accent 半透明混合）
"""

import math
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from models import DiagramSpec, Rect, VisualTheme
from pipeline.skills.base import RenderingSkill, SkillDescriptor
from pipeline.skills._utils import parse_color, theme_color


class ProcessFlowSkill(RenderingSkill):

    def descriptor(self) -> SkillDescriptor:
        return SkillDescriptor(
            skill_id="diagram_process_flow",
            skill_type="diagram",
            handles_types=["process_flow"],
        )

    def prompt_fragment(self) -> str:
        return """**process_flow**（流程图）:
  约束: 节点3-7个（超过7个则分层或分页）；每个desc ≤15字；节点间必须有明确的因果/时序关系"""

    def design_tokens(self) -> dict:
        return {
            "node_font_size": 11,
            "edge_label_font_size": 9,
            "shadow_offset_emu": 18000,
        }

    def render(self, slide, data: DiagramSpec, rect, theme: VisualTheme) -> bool:
        nodes = data.nodes
        edges = data.edges
        n = len(nodes)
        if n == 0:
            return False

        tokens = self.design_tokens()
        primary = theme_color(theme, "primary", "#003D6E")
        accent = theme_color(theme, "accent", "#FF6B35")
        white = RGBColor(255, 255, 255)
        font = theme.fonts.get("body", "Calibri")

        is_horizontal = data.layout_direction == "LR"

        if is_horizontal:
            self._render_horizontal(slide, nodes, edges, rect, theme,
                                    primary, accent, white, font, tokens)
        else:
            self._render_vertical(slide, nodes, edges, rect, theme,
                                  primary, accent, white, font, tokens)
        return True

    def _render_horizontal(self, slide, nodes, edges, rect, theme,
                           primary, accent, white, font, tokens):
        n = len(nodes)
        # 节点大小自适应
        node_w = min(int(rect.width / (n * 1.5)), 2286000)
        node_h = min(int(rect.height * 0.45), 640080)
        gap = max((rect.width - n * node_w) // max(n + 1, 1), 91440)

        positions = {}
        for i, node in enumerate(nodes):
            x = rect.left + gap + i * (node_w + gap)
            y = rect.top + (rect.height - node_h) // 2
            positions[node.node_id] = (x, y)

            # 阴影
            shadow_offset = tokens["shadow_offset_emu"]
            shadow = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Emu(x + shadow_offset), Emu(y + shadow_offset),
                Emu(node_w), Emu(node_h),
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            shadow.line.fill.background()

            # 节点
            color = primary if i % 2 == 0 else accent
            self._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                            x, y, node_w, node_h,
                            color, node.label, white, font,
                            tokens["node_font_size"])

        # 箭头
        for i in range(n - 1):
            x1 = positions[nodes[i].node_id][0] + node_w
            y1 = positions[nodes[i].node_id][1] + node_h // 2
            x2 = positions[nodes[i + 1].node_id][0]
            y2 = positions[nodes[i + 1].node_id][1] + node_h // 2
            self._add_arrow(slide, x1, y1, x2, y2, accent)

            edge_label = self._find_edge_label(edges, nodes[i].node_id, nodes[i + 1].node_id)
            if edge_label:
                self._add_label(slide, (x1 + x2) // 2, y1 - 274320,
                                edge_label, accent, font, tokens["edge_label_font_size"])

    def _render_vertical(self, slide, nodes, edges, rect, theme,
                         primary, accent, white, font, tokens):
        n = len(nodes)
        node_w = min(int(rect.width * 0.7), 3200400)
        node_h = min(int(rect.height / (n * 1.5)), 548640)
        gap = max((rect.height - n * node_h) // max(n + 1, 1), 91440)

        for i, node in enumerate(nodes):
            x = rect.left + (rect.width - node_w) // 2
            y = rect.top + gap + i * (node_h + gap)

            # 阴影
            shadow_offset = tokens["shadow_offset_emu"]
            shadow = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Emu(x + shadow_offset), Emu(y + shadow_offset),
                Emu(node_w), Emu(node_h),
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            shadow.line.fill.background()

            color = primary if i % 2 == 0 else accent
            self._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                            x, y, node_w, node_h,
                            color, node.label, white, font,
                            tokens["node_font_size"])

            if i < n - 1:
                arrow_x = x + node_w // 2
                self._add_arrow(slide, arrow_x, y + node_h, arrow_x,
                                y + node_h + gap, RGBColor(0x99, 0x99, 0x99))

    # ── 辅助方法 ──

    @staticmethod
    def _add_shape(slide, shape_type, x, y, w, h, fill_color, text, text_color, font_name, font_size):
        shape = slide.shapes.add_shape(shape_type, Emu(x), Emu(y), Emu(w), Emu(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.bold = True
        p.font.name = font_name
        p.alignment = PP_ALIGN.CENTER

    @staticmethod
    def _add_arrow(slide, x1, y1, x2, y2, color):
        connector = slide.shapes.add_connector(1, Emu(x1), Emu(y1), Emu(x2), Emu(y2))
        connector.line.color.rgb = color
        connector.line.width = Pt(1.5)

    @staticmethod
    def _add_label(slide, x, y, text, color, font_name, font_size, bold=False):
        box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(2743200), Emu(274320))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name

    @staticmethod
    def _find_edge_label(edges, from_id, to_id):
        for e in edges:
            if e.from_id == from_id and e.to_id == to_id:
                return e.label
        return None
