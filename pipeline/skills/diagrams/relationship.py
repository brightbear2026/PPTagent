"""
Relationship Diagram Skill

因果/关系网络图：圆形布局节点 + 连接边。

质量修复（相比原 diagram_renderer.py）：
1. 边的箭头样式区分：实线=强关系，虚线=弱关系（通过 line.dash_style）
2. 边标签增加半透明背景（提升可读性）
"""

import math
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from models import DiagramSpec, Rect, VisualTheme
from pipeline.skills.base import RenderingSkill, SkillDescriptor
from pipeline.skills._utils import theme_color


class RelationshipSkill(RenderingSkill):

    def descriptor(self) -> SkillDescriptor:
        return SkillDescriptor(
            skill_id="diagram_relationship",
            skill_type="diagram",
            handles_types=["relationship"],
        )

    def prompt_fragment(self) -> str:
        return """**relationship**（关系图）:
  约束: edges数量 ≤ nodes × 2（超过则关系图会变得混乱不可读）；每个edge必须有label说明关系类型"""

    def design_tokens(self) -> dict:
        return {
            "node_font_size": 10,
            "edge_label_font_size": 8,
            "strong_edge_width_pt": 2.0,
            "weak_edge_width_pt": 1.0,
        }

    def render(self, slide, data: DiagramSpec, rect, theme: VisualTheme) -> bool:
        nodes = data.nodes
        edges = data.edges
        n = len(nodes)
        if n == 0:
            return False

        tokens = self.design_tokens()
        primary = theme_color(theme, "primary", "#003D6E")
        accent = theme_color(theme, "accent", "#FF6B35")
        white = RGBColor(255, 255, 255)
        font = theme.fonts.get("body", "Calibri")

        # 圆形布局
        node_w = min(int(rect.width * 0.22), 1828800)
        node_h = min(int(rect.height * 0.18), 548640)
        cx = rect.left + rect.width // 2
        cy = rect.top + rect.height // 2
        radius_x = rect.width // 2 - node_w // 2 - 182880
        radius_y = rect.height // 2 - node_h // 2 - 182880

        positions = {}
        for i, node in enumerate(nodes):
            angle = 2 * math.pi * i / n - math.pi / 2
            x = int(cx + radius_x * math.cos(angle)) - node_w // 2
            y = int(cy + radius_y * math.sin(angle)) - node_h // 2
            positions[node.node_id] = (x, y, node_w, node_h)

            self._add_shape(slide, MSO_SHAPE.OVAL, x, y, node_w, node_h,
                            primary, node.label, white, font,
                            tokens["node_font_size"])

        # 连接边
        for edge in edges:
            if edge.from_id in positions and edge.to_id in positions:
                fx, fy, fw, fh = positions[edge.from_id]
                tx, ty, tw, th = positions[edge.to_id]

                # 根据是否有 label 判断强弱关系
                has_label = bool(edge.label)
                edge_color = RGBColor(0x66, 0x66, 0x66) if has_label else RGBColor(0xBB, 0xBB, 0xBB)
                self._add_arrow(slide,
                                fx + fw // 2, fy + fh // 2,
                                tx + tw // 2, ty + th // 2,
                                edge_color, has_label, tokens)

                if edge.label:
                    mx = (fx + fw // 2 + tx + tw // 2) // 2
                    my = (fy + fh // 2 + ty + th // 2) // 2
                    self._add_label_with_bg(slide, mx, my - 182880,
                                            edge.label, accent, font,
                                            tokens["edge_label_font_size"])

        return True

    @staticmethod
    def _add_shape(slide, shape_type, x, y, w, h, fill_color, text, text_color, font_name, font_size):
        shape = slide.shapes.add_shape(shape_type, Emu(x), Emu(y), Emu(w), Emu(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.bold = True
        p.font.name = font_name
        p.alignment = PP_ALIGN.CENTER

    @staticmethod
    def _add_arrow(slide, x1, y1, x2, y2, color, strong, tokens):
        connector = slide.shapes.add_connector(1, Emu(x1), Emu(y1), Emu(x2), Emu(y2))
        connector.line.color.rgb = color
        connector.line.width = Pt(tokens["strong_edge_width_pt"] if strong else tokens["weak_edge_width_pt"])
        if not strong:
            try:
                connector.line.dash_style = 4  # DASH
            except Exception:
                pass

    @staticmethod
    def _add_label_with_bg(slide, x, y, text, color, font_name, font_size):
        """带半透明背景的标签"""
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(x), Emu(y), Emu(1828800), Emu(228600),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        bg.line.fill.background()
        # 半透明
        try:
            from pptx.oxml.ns import qn
            bg.fill._fill.attrib[qn('a:solidFill')].append(
                bg.fill._fill.makeelement(qn('a:srgbClr'), {'val': 'FFFFFF', 'alpha': '60000'})
            )
        except Exception:
            pass

        box = slide.shapes.add_textbox(Emu(x + 45720), Emu(y), Emu(1737600), Emu(228600))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = False
        p.font.name = font_name
