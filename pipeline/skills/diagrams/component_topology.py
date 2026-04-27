"""
Component Topology Diagram Skill

微服务/组件拓扑图：分组框 + 组间简化连线。
第一版限制：≤6 节点，单层线性调用。超出截断显示 +N more。
"""

from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from models import DiagramSpec, Rect, VisualTheme
from pipeline.skills.base import RenderingSkill, SkillDescriptor
from pipeline.skills._utils import theme_color

_MAX_NODES = 6
_MAX_CONNECTIONS_RATIO = 2  # connections ≤ groups × ratio


class ComponentTopologySkill(RenderingSkill):

    def descriptor(self) -> SkillDescriptor:
        return SkillDescriptor(
            skill_id="diagram_component_topology",
            skill_type="diagram",
            handles_types=["component_topology"],
        )

    def prompt_fragment(self) -> str:
        return """**component_topology**（组件/微服务拓扑图）:
  约束: ≤6个组件节点，groups描述服务分组，connections描述调用关系；每条connection必须有label说明协议/接口
  反模式: 不要用于多层堆叠架构（用tech_architecture或architecture）；不要用于因果关系（用relationship）"""

    def design_tokens(self) -> dict:
        return {
            "group_border_pt": 1.5,
            "node_font_size": 10,
            "connection_font_size": 8,
        }

    def render(self, slide, data: DiagramSpec, rect, theme: VisualTheme) -> bool:
        nodes = data.nodes
        if not nodes:
            return False

        tokens = self.design_tokens()
        primary = theme_color(theme, "primary", "#003D6E")
        accent = theme_color(theme, "accent", "#FF6B35")
        white = RGBColor(255, 255, 255)
        muted = RGBColor(0x99, 0x99, 0x99)
        font = theme.fonts.get("body", "Calibri")

        # 按 group 分组
        groups = {}
        for node in nodes:
            group = node.group or "default"
            groups.setdefault(group, []).append(node)

        # 视觉降级：超过 _MAX_NODES 截断
        total_nodes = len(nodes)
        truncated = total_nodes > _MAX_NODES

        group_names = list(groups.keys())
        num_groups = len(group_names)
        if num_groups == 0:
            return False

        group_h = (rect.height - (num_groups - 1) * 91440) // num_groups

        group_positions = {}
        node_positions = {}

        for gi, group_name in enumerate(group_names):
            group_nodes = groups[group_name]
            # 如果是最后一个 group 且有截断，显示 +N more
            if truncated and gi == num_groups - 1:
                remaining = total_nodes - sum(len(groups[gn]) for gn in group_names[:gi])
                if remaining > _MAX_NODES:
                    group_nodes = group_nodes[:_MAX_NODES]

            gy = rect.top + gi * (group_h + 91440)

            # 分组框
            group_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(rect.left), Emu(gy),
                Emu(rect.width), Emu(group_h),
            )
            group_box.fill.solid()
            group_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
            group_box.line.color.rgb = primary
            group_box.line.width = Pt(tokens["group_border_pt"])

            # 分组标题
            title_box = slide.shapes.add_textbox(
                Emu(rect.left + 91440), Emu(gy + 45720),
                Emu(rect.width - 182880), Emu(228600),
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = group_name
            p.font.size = Pt(11)
            p.font.color.rgb = primary
            p.font.bold = True
            p.font.name = font

            # 组内节点
            n = len(group_nodes)
            item_w = min((rect.width - 182880) // max(n, 1) - 45720, 2286000)
            item_h = 365760
            total_w = n * item_w + (n - 1) * 45720
            start_x = rect.left + (rect.width - total_w) // 2

            for ni, node in enumerate(group_nodes):
                nx = start_x + ni * (item_w + 45720)
                ny = gy + group_h // 2 + 91440

                node_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Emu(nx), Emu(ny), Emu(item_w), Emu(item_h),
                )
                node_box.fill.solid()
                node_box.fill.fore_color.rgb = white
                node_box.line.color.rgb = accent
                node_box.line.width = Pt(1)
                ntf = node_box.text_frame
                ntf.word_wrap = True
                np_ = ntf.paragraphs[0]
                np_.text = node.label
                np_.font.size = Pt(tokens["node_font_size"])
                np_.font.color.rgb = primary
                np_.font.bold = True
                np_.font.name = font
                np_.alignment = PP_ALIGN.CENTER

                node_positions[node.node_id] = (nx, ny, item_w, item_h)

            group_positions[group_name] = (rect.left, gy, rect.width, group_h)

            # +N more 标签
            if truncated and gi == num_groups - 1:
                more_box = slide.shapes.add_textbox(
                    Emu(rect.left + rect.width // 2 - 914400),
                    Emu(gy + group_h - 274320),
                    Emu(1828800), Emu(228600),
                )
                mt = more_box.text_frame
                mp = mt.paragraphs[0]
                mp.text = f"+{total_nodes - _MAX_NODES} more"
                mp.font.size = Pt(9)
                mp.font.color.rgb = muted
                mp.font.italic = True
                mp.font.name = font
                mp.alignment = PP_ALIGN.CENTER

        # 组间连线（简化：只画组间第一条连线）
        edges = data.edges
        drawn_between = set()
        for edge in edges:
            if edge.from_id in node_positions and edge.to_id in node_positions:
                from_group = None
                to_group = None
                for gn, gnodes in groups.items():
                    if any(n.node_id == edge.from_id for n in gnodes):
                        from_group = gn
                    if any(n.node_id == edge.to_id for n in gnodes):
                        to_group = gn
                if from_group and to_group and from_group != to_group:
                    pair = (from_group, to_group)
                    if pair in drawn_between:
                        continue
                    drawn_between.add(pair)
                    # 画组间箭头（简化：从上组底部到下组顶部）
                    if from_group in group_positions and to_group in group_positions:
                        _, fy, _, fh = group_positions[from_group]
                        _, ty, _, _ = group_positions[to_group]
                        mid_x = rect.left + rect.width // 2
                        connector = slide.shapes.add_connector(
                            1,
                            Emu(mid_x), Emu(fy + fh),
                            Emu(mid_x), Emu(ty),
                        )
                        connector.line.color.rgb = muted
                        connector.line.width = Pt(1)
                        # 连线标签
                        if edge.label:
                            lb = slide.shapes.add_textbox(
                                Emu(mid_x + 91440), Emu((fy + fh + ty) // 2 - 91440),
                                Emu(1371600), Emu(182880),
                            )
                            lt = lb.text_frame
                            lp = lt.paragraphs[0]
                            lp.text = edge.label
                            lp.font.size = Pt(tokens["connection_font_size"])
                            lp.font.color.rgb = accent
                            lp.font.name = font

        return True
