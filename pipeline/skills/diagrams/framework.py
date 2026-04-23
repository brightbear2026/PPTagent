"""
Framework Diagram Skill

框架图变体：2x2 矩阵、SWOT、金字塔、漏斗。

质量修复（相比原 diagram_renderer.py）：
1. 2x2 矩阵增加象限标签（角标注含义）
2. 金字塔层间增加颜色渐变（顶深底浅）
3. 漏斗每层宽度递减更自然
"""

from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from models import DiagramSpec, Rect, VisualTheme
from pipeline.skills.base import RenderingSkill, SkillDescriptor
from pipeline.skills._utils import theme_color


class FrameworkSkill(RenderingSkill):

    def descriptor(self) -> SkillDescriptor:
        return SkillDescriptor(
            skill_id="diagram_framework",
            skill_type="diagram",
            handles_types=["framework"],
        )

    def prompt_fragment(self) -> str:
        return """**framework**（框架图）:
  约束: matrix_2x2必须4象限全填；pyramid ≤5层；funnel ≤6层且每层宽度要有实际含义"""

    def design_tokens(self) -> dict:
        return {
            "title_font_size": 13,
            "item_font_size": 10,
            "pyramid_font_size": 11,
        }

    def render(self, slide, data: DiagramSpec, rect, theme: VisualTheme) -> bool:
        nodes = data.nodes
        if not nodes:
            return False

        # 按 group 分组
        groups = {}
        for node in nodes:
            g = node.group or "default"
            groups.setdefault(g, []).append(node)
        group_keys = list(groups.keys())

        # 2x2 矩阵 / SWOT
        if len(group_keys) == 4 and set(group_keys) & {"top_left", "top_right", "bottom_left", "bottom_right"}:
            self._render_2x2(slide, groups, rect, theme, group_keys)
        # 金字塔
        elif len(group_keys) == 1 and len(groups[group_keys[0]]) >= 3:
            self._render_pyramid(slide, groups[group_keys[0]], rect, theme)
        # 漏斗
        elif len(group_keys) == 1:
            self._render_funnel(slide, groups[group_keys[0]], rect, theme)
        else:
            return False  # 交给 fallback

        return True

    def _render_2x2(self, slide, groups, rect, theme, group_keys):
        """2x2 矩阵 + 象限标签"""
        tokens = self.design_tokens()
        primary = theme_color(theme, "primary", "#003D6E")
        white = RGBColor(255, 255, 255)
        font = theme.fonts.get("body", "Calibri")

        quadrant_colors = {
            "top_left": RGBColor(0x00, 0x3D, 0x6E),
            "top_right": RGBColor(0x00, 0x7B, 0xC0),
            "bottom_left": RGBColor(0xFF, 0x6B, 0x35),
            "bottom_right": RGBColor(0x48, 0xA9, 0xE6),
        }

        # 象限中文标签（角标注含义）
        axis_labels = {
            "top_left": "高",
            "top_right": "高",
            "bottom_left": "低",
            "bottom_right": "低",
        }

        gap = 91440
        half_w = (rect.width - gap) // 2
        half_h = (rect.height - gap) // 2

        positions = {
            "top_left": (rect.left, rect.top),
            "top_right": (rect.left + half_w + gap, rect.top),
            "bottom_left": (rect.left, rect.top + half_h + gap),
            "bottom_right": (rect.left + half_w + gap, rect.top + half_h + gap),
        }

        for key in group_keys:
            x, y = positions.get(key, (rect.left, rect.top))
            color = quadrant_colors.get(key, primary)

            # 象限背景
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(x), Emu(y), Emu(half_w), Emu(half_h),
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = color
            bg.line.fill.background()

            # 象限标签（角标注）
            label = axis_labels.get(key, "")
            if label:
                self._add_label(slide, x + 45720, y + 45720, label,
                                white, font, 8, bold=False)

            # 标题 + 子项
            items = groups[key]
            if items:
                self._add_label(slide, x + 91440, y + 91440, items[0].label,
                                white, font, tokens["title_font_size"], bold=True)
                for j, item in enumerate(items[1:6]):
                    self._add_label(slide, x + 91440, y + 365760 + j * 228600,
                                    f"• {item.label}", white, font,
                                    tokens["item_font_size"])

    def _render_pyramid(self, slide, nodes, rect, theme):
        """金字塔图 — 顶深底浅渐变"""
        tokens = self.design_tokens()
        n = len(nodes)
        if n == 0:
            return

        primary = theme_color(theme, "primary", "#003D6E")
        accent = theme_color(theme, "accent", "#FF6B35")
        white = RGBColor(255, 255, 255)
        font = theme.fonts.get("body", "Calibri")

        # 渐变色：从深到浅
        base_colors = [
            RGBColor(0x00, 0x1A, 0x3A),
            RGBColor(0x00, 0x3D, 0x6E),
            RGBColor(0x00, 0x5B, 0x96),
            RGBColor(0x00, 0x7B, 0xC0),
            RGBColor(0x48, 0xA9, 0xE6),
        ]

        layer_h = rect.height // n
        cx = rect.left + rect.width // 2

        for i, node in enumerate(nodes):
            width_ratio = (i + 1) / n
            w = int(rect.width * width_ratio * 0.85)
            x = cx - w // 2
            y = rect.top + i * layer_h

            color = base_colors[min(i, len(base_colors) - 1)]
            self._add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, layer_h - 45720,
                            color, node.label, white, font,
                            tokens["pyramid_font_size"])

    def _render_funnel(self, slide, nodes, rect, theme):
        """漏斗图"""
        tokens = self.design_tokens()
        n = len(nodes)
        if n == 0:
            return

        white = RGBColor(255, 255, 255)
        font = theme.fonts.get("body", "Calibri")

        colors = [
            RGBColor(0x00, 0x3D, 0x6E), RGBColor(0x00, 0x5B, 0x96),
            RGBColor(0x00, 0x7B, 0xC0), RGBColor(0x48, 0xA9, 0xE6),
            RGBColor(0x7E, 0xC8, 0xE3), RGBColor(0xB8, 0xE4, 0xF0),
        ]

        layer_h = rect.height // n
        cx = rect.left + rect.width // 2

        for i, node in enumerate(nodes):
            width_ratio = 1.0 - (i / n) * 0.6
            w = int(rect.width * width_ratio)
            x = cx - w // 2
            y = rect.top + i * layer_h
            color = colors[i % len(colors)]

            self._add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, layer_h - 45720,
                            color, node.label, white, font,
                            tokens["pyramid_font_size"])

    # ── 辅助方法 ──

    @staticmethod
    def _add_shape(slide, shape_type, x, y, w, h, fill_color, text, text_color, font_name, font_size):
        shape = slide.shapes.add_shape(shape_type, Emu(x), Emu(y), Emu(w), Emu(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.bold = True
        p.font.name = font_name
        p.alignment = PP_ALIGN.CENTER

    @staticmethod
    def _add_label(slide, x, y, text, color, font_name, font_size, bold=False):
        box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(2743200), Emu(274320))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
