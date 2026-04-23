"""
Architecture Diagram Skill

层级架构图：按 group 分层渲染。

质量修复（相比原 diagram_renderer.py）：
1. 层与层之间增加分隔线
2. 层标签增加背景色条
"""

from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from models import DiagramSpec, Rect, VisualTheme
from pipeline.skills.base import RenderingSkill, SkillDescriptor
from pipeline.skills._utils import theme_color


class ArchitectureSkill(RenderingSkill):

    def descriptor(self) -> SkillDescriptor:
        return SkillDescriptor(
            skill_id="diagram_architecture",
            skill_type="diagram",
            handles_types=["architecture"],
        )

    def prompt_fragment(self) -> str:
        return """**architecture**（架构图）:
  约束: 最多4层，每层items ≤5个；层与层之间要有逻辑依赖关系"""

    def design_tokens(self) -> dict:
        return {
            "separator_height_emu": 18288,  # ~2px 分隔线
            "separator_color": "#FFFFFF",
            "label_bg_width_emu": 1828800,
            "label_bg_height_emu": 228600,
        }

    def render(self, slide, data: DiagramSpec, rect, theme: VisualTheme) -> bool:
        nodes = data.nodes
        if not nodes:
            return False

        tokens = self.design_tokens()
        white = RGBColor(255, 255, 255)
        font = theme.fonts.get("body", "Calibri")

        # 按 group 分层
        layers = {}
        for node in nodes:
            group = node.group or "default"
            layers.setdefault(group, []).append(node)

        layer_names = list(layers.keys())
        num_layers = len(layer_names)
        if num_layers == 0:
            return False

        layer_h = rect.height // num_layers
        layer_colors = [
            RGBColor(0x1A, 0x47, 0x8A),
            RGBColor(0x00, 0x5B, 0x96),
            RGBColor(0x00, 0x7B, 0xC0),
            RGBColor(0x48, 0xA9, 0xE6),
        ]

        for li, layer_name in enumerate(layer_names):
            layer_nodes = layers[layer_name]
            y = rect.top + li * layer_h
            color = layer_colors[li % len(layer_colors)]

            # 层背景
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(rect.left), Emu(y),
                Emu(rect.width), Emu(layer_h - 45720),
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = color
            bg.line.fill.background()

            # 层标签（左侧色条 + 标签名）
            label_w = tokens["label_bg_width_emu"]
            label_h = tokens["label_bg_height_emu"]
            label_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Emu(rect.left + 45720), Emu(y + 45720),
                Emu(label_w), Emu(label_h),
            )
            label_bg.fill.solid()
            label_bg.fill.fore_color.rgb = RGBColor(
                max(0, color[0] - 40), max(0, color[1] - 40), max(0, color[2] - 40),
            )
            label_bg.line.fill.background()
            tf_label = label_bg.text_frame
            tf_label.word_wrap = True
            p_label = tf_label.paragraphs[0]
            p_label.text = layer_name
            p_label.font.size = Pt(9)
            p_label.font.color.rgb = white
            p_label.font.bold = True
            p_label.font.name = font
            p_label.alignment = PP_ALIGN.CENTER

            # 分隔线（层与层之间）
            if li > 0:
                sep_h = tokens["separator_height_emu"]
                sep = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Emu(rect.left), Emu(y - sep_h // 2),
                    Emu(rect.width), Emu(sep_h),
                )
                sep.fill.solid()
                sep.fill.fore_color.rgb = white
                sep.line.fill.background()

            # 节点
            n_items = len(layer_nodes)
            item_w = min(rect.width // max(n_items, 1) - 91440, 2743200)
            item_h = min(layer_h // 2, 457200)
            total_w = n_items * item_w + (n_items - 1) * 91440
            start_x = rect.left + (rect.width - total_w) // 2

            for ni, node in enumerate(layer_nodes):
                nx = start_x + ni * (item_w + 91440)
                ny = y + (layer_h - item_h) // 2 + label_h // 2
                self._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                                nx, ny, item_w, item_h,
                                white, node.label, color, font, 10)

        return True

    @staticmethod
    def _add_shape(slide, shape_type, x, y, w, h, fill_color, text, text_color, font_name, font_size):
        shape = slide.shapes.add_shape(shape_type, Emu(x), Emu(y), Emu(w), Emu(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.bold = True
        p.font.name = font_name
        p.alignment = PP_ALIGN.CENTER
