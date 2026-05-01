"""
R1 regression test: chart_renderer must skip chart injection on structural slides
(section_divider, agenda, title) per REQUIREMENTS.md H3.
"""
import pytest
from unittest.mock import MagicMock, patch, call


def _make_renderer():
    from pipeline.layer6_output.chart_renderer import ChartRenderer
    return ChartRenderer()


def _chart_spec_dict():
    return {
        "chart_type": "column",
        "title": "Test Chart",
        "categories": ["A", "B", "C"],
        "series": [{"name": "S1", "values": [10, 20, 30]}],
        "so_what": "峰值 30",
    }


def _placeholder_entry(slide_index=0):
    return {
        "slide_index": slide_index,
        "items": [{"id": "chart_1", "x": 1.0, "y": 1.0, "w": 4.0, "h": 3.0}],
    }


def _slide_data(slide_type="content"):
    return {"slide_type": slide_type, "chart_spec": _chart_spec_dict(), "theme": None}


@pytest.mark.parametrize("slide_type", ["section_divider", "agenda", "title"])
def test_chart_skipped_on_structural_slide(slide_type, tmp_path):
    """Chart must NOT be injected into section_divider / agenda / title slides."""
    import json

    renderer = _make_renderer()

    # Create a minimal real pptx for Presentation() to open
    from pptx import Presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    input_path = str(tmp_path / "input.pptx")
    output_path = str(tmp_path / "output.pptx")
    prs.save(input_path)

    placeholders = [_placeholder_entry(0)]
    slides_data = [_slide_data(slide_type)]

    renderer.render_into_pptx(input_path, placeholders, slides_data, output_path)

    # Verify no chart shapes were added (only the blank slide should exist)
    result_prs = Presentation(output_path)
    slide = result_prs.slides[0]
    chart_shapes = [s for s in slide.shapes if s.has_chart]
    assert len(chart_shapes) == 0, (
        f"Expected 0 charts on {slide_type} slide, got {len(chart_shapes)}"
    )


def test_chart_injected_on_content_slide(tmp_path):
    """Chart MUST be injected on normal content slides."""
    renderer = _make_renderer()

    from pptx import Presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    input_path = str(tmp_path / "input.pptx")
    output_path = str(tmp_path / "output.pptx")
    prs.save(input_path)

    placeholders = [_placeholder_entry(0)]
    slides_data = [_slide_data("content")]

    renderer.render_into_pptx(input_path, placeholders, slides_data, output_path)

    result_prs = Presentation(output_path)
    slide = result_prs.slides[0]
    chart_shapes = [s for s in slide.shapes if s.has_chart]
    assert len(chart_shapes) >= 1, "Expected chart on content slide"


def test_mixed_slides_only_inject_content(tmp_path):
    """When multiple slides, only content slides get charts."""
    from pptx import Presentation

    renderer = _make_renderer()
    prs = Presentation()
    # 3 slides: section_divider, content, section_divider
    for _ in range(3):
        prs.slides.add_slide(prs.slide_layouts[6])
    input_path = str(tmp_path / "input.pptx")
    output_path = str(tmp_path / "output.pptx")
    prs.save(input_path)

    placeholders = [
        _placeholder_entry(0),
        _placeholder_entry(1),
        _placeholder_entry(2),
    ]
    slides_data = [
        _slide_data("section_divider"),
        _slide_data("content"),
        _slide_data("section_divider"),
    ]

    renderer.render_into_pptx(input_path, placeholders, slides_data, output_path)

    result_prs = Presentation(output_path)
    # slide 0: no chart (section_divider)
    assert len([s for s in result_prs.slides[0].shapes if s.has_chart]) == 0
    # slide 1: has chart (content)
    assert len([s for s in result_prs.slides[1].shapes if s.has_chart]) >= 1
    # slide 2: no chart (section_divider)
    assert len([s for s in result_prs.slides[2].shapes if s.has_chart]) == 0
