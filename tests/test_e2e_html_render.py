"""
End-to-end test: synthetic data through HTMLDesignAgent fallback path.

Tests the full render pipeline without LLM:
  outline+content → HTMLDesignAgent(fallback HTML) → html2pptx.js → ChartRenderer
"""

import json
import os
import sys
import tempfile

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pipeline.agents.html_design_agent import HTMLDesignAgent
from pipeline.layer6_output.node_bridge import is_node_available


def test_e2e_synthetic():
    if not is_node_available():
        print("SKIP: Node.js not available")
        return

    # Simulate 5 slides of outline+content
    outline_items = [
        {"page_number": 1, "slide_type": "title", "takeaway_message": "AI Agent 安全风险分析报告", "title": "封面"},
        {"page_number": 2, "slide_type": "content", "takeaway_message": "AI Agent 面临三类核心安全威胁", "title": "威胁概览"},
        {"page_number": 3, "slide_type": "content", "takeaway_message": "提示注入攻击是最普遍的攻击向量", "title": "提示注入",
         "chart_suggestion": {"chart_type": "column", "title": "攻击类型分布", "categories": ["Prompt Injection", "Data Leakage", "Tool Abuse", "Privilege Escalation"], "series": [{"name": "Incidents", "values": [45, 28, 18, 9]}]}},
        {"page_number": 4, "slide_type": "content", "takeaway_message": "防御需要分层架构", "title": "防御框架"},
        {"page_number": 5, "slide_type": "content", "takeaway_message": "未来趋势：可验证 Agent 与形式化验证", "title": "展望"},
    ]

    content_slides = [
        {"page_number": 1, "slide_type": "title", "takeaway_message": "AI Agent 安全风险分析报告", "text_blocks": [], "chart_suggestion": None, "diagram_spec": None, "visual_block": None},
        {"page_number": 2, "slide_type": "content", "takeaway_message": "AI Agent 面临三类核心安全威胁", "text_blocks": [
            {"content": "提示注入攻击：通过恶意输入操控 Agent 行为", "is_bold": True},
            {"content": "数据泄露：Agent 在推理过程中暴露敏感信息", "is_bold": False},
            {"content": "权限滥用：Agent 越权执行未授权操作", "is_bold": False},
            {"content": "2024 年已报告超过 200 起相关安全事件", "is_bold": False},
        ], "chart_suggestion": None, "diagram_spec": None, "visual_block": None},
        {"page_number": 3, "slide_type": "content", "takeaway_message": "提示注入攻击是最普遍的攻击向量", "text_blocks": [
            {"content": "数据显示提示注入占所有攻击的 45%", "is_bold": True},
            {"content": "数据泄露占 28%，工具滥用占 18%", "is_bold": False},
            {"content": "权限提升类攻击较少但危害最大", "is_bold": False},
        ], "chart_suggestion": {"chart_type": "column", "title": "攻击类型分布", "categories": ["Prompt Injection", "Data Leakage", "Tool Abuse", "Privilege Escalation"], "series": [{"name": "Incidents", "values": [45, 28, 18, 9]}]},
        "diagram_spec": None, "visual_block": None},
        {"page_number": 4, "slide_type": "content", "takeaway_message": "防御需要分层架构", "text_blocks": [
            {"content": "输入验证层：严格校验所有用户输入", "is_bold": True},
            {"content": "权限控制层：最小权限原则 + RBAC", "is_bold": False},
            {"content": "审计监控层：全链路日志 + 异常检测", "is_bold": False},
            {"content": "输出过滤层：敏感信息脱敏", "is_bold": False},
        ], "chart_suggestion": None,
        "diagram_spec": {"diagram_type": "process_flow", "title": "防御架构", "nodes": [
            {"id": "n1", "label": "输入验证"}, {"id": "n2", "label": "权限控制"},
            {"id": "n3", "label": "审计监控"}, {"id": "n4", "label": "输出过滤"},
        ]}, "visual_block": None},
        {"page_number": 5, "slide_type": "content", "takeaway_message": "未来趋势：可验证 Agent 与形式化验证", "text_blocks": [
            {"content": "可验证 Agent：通过形式化方法证明安全属性", "is_bold": True},
            {"content": "零知识证明应用于 Agent 隐私保护", "is_bold": False},
            {"content": "多 Agent 协作中的博弈论安全机制", "is_bold": False},
        ], "chart_suggestion": None, "diagram_spec": None, "visual_block": None},
    ]

    context = {
        "task": {"task_id": "test_e2e", "source_type": "text"},
        "task_id": "test_e2e",
        "outline": {"items": outline_items},
        "content": {"slides": content_slides},
        "analysis": {},
        "report_progress": lambda pct, msg: print(f"  [{pct}%] {msg}"),
    }

    agent = HTMLDesignAgent(llm_client=None)  # No LLM, use fallback
    result = agent.run(context)

    print(f"\n=== E2E Result ===")
    print(f"Output: {result.get('output_file')}")
    print(f"Slides: {result.get('slide_count')}")
    print(f"Charts: {result.get('chart_count')}")
    print(f"Diagrams: {result.get('diagram_count')}")
    print(f"Render errors: {len(result.get('render_errors', []))}")

    # Verify output file
    output_path = result.get("output_file", "")
    if output_path and os.path.exists(output_path):
        from pptx import Presentation
        prs = Presentation(output_path)
        total_shapes = sum(len(s.shapes) for s in prs.slides)
        has_charts = any(shape.has_chart for s in prs.slides for shape in s.shapes)
        has_text = any(shape.has_text_frame for s in prs.slides for shape in s.shapes)

        print(f"\n=== PPTX Verification ===")
        print(f"Slides: {len(prs.slides)}")
        print(f"Total shapes: {total_shapes}")
        print(f"Has charts: {has_charts}")
        print(f"Has text: {has_text}")

        # Check text content
        all_text = []
        for s in prs.slides:
            for shape in s.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            all_text.append(para.text.strip())

        print(f"Text elements: {len(all_text)}")
        print(f"Sample text: {all_text[:5]}")

        assert len(prs.slides) >= 4, f"Expected >= 4 slides, got {len(prs.slides)}"
        assert has_text, "No text found in PPTX"
        print("\n=== E2E TEST PASSED ===")
    else:
        print(f"\nERROR: Output file not found at {output_path}")
        sys.exit(1)


if __name__ == "__main__":
    test_e2e_synthetic()
