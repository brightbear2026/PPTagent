"""
PPTBuilder 端到端回归测试

覆盖 P0/P1/P2 关键修复点：
- 封面 / 章节过渡 / 议程 / 内容 / 数据 / 对比 五种 slide_type 都能成功生成
- 不出现重复 shape、不出现空 TextBox、不出现越界 shape
- 中英文字号策略生效
- chart_id 去重
- 空白正文页兜底
"""
from __future__ import annotations

import os
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

from models import (
    PresentationSpec, SlideSpec, SlideType, NarrativeRole, TextBlock,
    ChartSpec, ChartSeries, ChartType,
)
from pipeline.layer4_visual import VisualDesigner
from pipeline.layer4_visual.visual_designer import LANGUAGE_FONT_STRATEGY
from pipeline.layer6_output.ppt_builder import PPTBuilder


OUTPUT_DIR = Path("output")


def _build(pres: PresentationSpec, language: str = "zh") -> Presentation:
    VisualDesigner().design_slides(pres.slides, language=language)
    path = PPTBuilder().build(pres)
    return Presentation(path)


def _shape_texts(slide):
    return [
        sh.text_frame.text.strip()
        for sh in slide.shapes
        if sh.has_text_frame and sh.text_frame.text.strip()
    ]


def _shapes_have_overlap(slide, iou_threshold: float = 0.6) -> bool:
    rects = []
    for sh in slide.shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        l, t = sh.left or 0, sh.top or 0
        rects.append((l, t, l + (sh.width or 0), t + (sh.height or 0)))
    for i in range(len(rects)):
        l1, t1, r1, b1 = rects[i]
        a1 = max(0, (r1 - l1) * (b1 - t1))
        if a1 == 0:
            continue
        for j in range(i + 1, len(rects)):
            l2, t2, r2, b2 = rects[j]
            ix = max(0, min(r1, r2) - max(l1, l2))
            iy = max(0, min(b1, b2) - max(t1, t2))
            inter = ix * iy
            if inter == 0:
                continue
            a2 = max(0, (r2 - l2) * (b2 - t2))
            union = a1 + a2 - inter
            if (inter / union) > iou_threshold:
                return True
    return False


def _has_empty_textbox(slide) -> bool:
    for sh in slide.shapes:
        if sh.has_text_frame and sh.shape_type == 17:
            if not sh.text_frame.text.strip():
                return True
    return False


def _has_offcanvas_shape(slide, sw: int, sh_h: int) -> bool:
    tol = 9144  # 0.01 inch
    for sh in slide.shapes:
        l, t = sh.left or 0, sh.top or 0
        r, b = l + (sh.width or 0), t + (sh.height or 0)
        if l < -tol or t < -tol or r > sw + tol or b > sh_h + tol:
            return True
    return False


class TestPPTBuilderRegression(unittest.TestCase):

    def setUp(self):
        OUTPUT_DIR.mkdir(exist_ok=True)

    def _make_full_deck(self, title: str, language: str) -> PresentationSpec:
        pres = PresentationSpec(
            title=title,
            subtitle="自动化回归测试",
            author="pytest",
            created_at="2026-04-11",
            language=language,
        )
        # 1 cover
        pres.slides.append(SlideSpec(
            slide_type=SlideType.TITLE,
            takeaway_message="自动化回归测试",
            language=language,
        ))
        # 2 agenda
        pres.slides.append(SlideSpec(
            slide_type=SlideType.AGENDA,
            takeaway_message="目录",
            text_blocks=[
                TextBlock(content="背景与挑战", level=0),
                TextBlock(content="解决方案", level=0),
                TextBlock(content="预期收益", level=0),
            ],
            language=language,
        ))
        # 3 section divider
        pres.slides.append(SlideSpec(
            slide_type=SlideType.SECTION_DIVIDER,
            takeaway_message="背景与挑战",
            text_blocks=[TextBlock(content="行业宏观分析", level=0)],
            language=language,
        ))
        # 4 content
        pres.slides.append(SlideSpec(
            slide_type=SlideType.CONTENT,
            takeaway_message="线上业务同比增长67%，成为第一增长引擎",
            narrative_arc=NarrativeRole.EVIDENCE,
            text_blocks=[
                TextBlock(content="线上零售收入超传统渠道", level=0),
                TextBlock(content="Q4单季 3.8 亿元", level=1),
                TextBlock(content="企业服务环比+18%", level=0),
            ],
            source_note="2024Q4 内部数据",
            language=language,
        ))
        # 5 comparison
        pres.slides.append(SlideSpec(
            slide_type=SlideType.COMPARISON,
            takeaway_message="新方案在效率和扩展性上明显优于传统方案",
            narrative_arc=NarrativeRole.COMPARISON,
            text_blocks=[
                TextBlock(content="传统方案", level=0),
                TextBlock(content="周期 3 个月", level=1),
                TextBlock(content="维护成本高", level=1),
                TextBlock(content="新方案", level=0),
                TextBlock(content="周期 1 个月", level=1),
                TextBlock(content="弹性扩展", level=1),
            ],
            language=language,
        ))
        # 6 data
        pres.slides.append(SlideSpec(
            slide_type=SlideType.DATA,
            takeaway_message="各业务线收入对比：线上零售领跑",
            charts=[ChartSpec(
                chart_type=ChartType.COLUMN,
                title="2024 Q3 vs Q4 收入(百万)",
                categories=["线上零售", "企业服务", "广告", "金融"],
                series=[
                    ChartSeries(name="Q3", values=[320, 180, 95, 140]),
                    ChartSeries(name="Q4", values=[380, 210, 110, 165]),
                ],
                so_what="线上零售Q4环比+18.8%",
                show_legend=True,
            )],
            language=language,
        ))
        return pres

    def test_full_deck_zh(self):
        pres = self._make_full_deck("回归测试_zh", "zh")
        prs = _build(pres, language="zh")
        self.assertEqual(len(prs.slides), 6)

        for i, slide in enumerate(prs.slides, 1):
            with self.subTest(slide=i):
                self.assertFalse(_has_empty_textbox(slide),
                                 f"slide {i} 含空 TextBox")
                self.assertFalse(_has_offcanvas_shape(slide, prs.slide_width,
                                                      prs.slide_height),
                                 f"slide {i} 有 shape 越界")
                self.assertFalse(_shapes_have_overlap(slide),
                                 f"slide {i} 文本框严重重叠")

    def test_cover_has_title_subtitle_author(self):
        pres = self._make_full_deck("封面字段验证", "zh")
        prs = _build(pres, language="zh")
        cover = prs.slides[0]
        texts = " ".join(_shape_texts(cover))
        self.assertIn("封面字段验证", texts)
        self.assertIn("自动化回归测试", texts)  # subtitle
        self.assertIn("pytest", texts)         # author
        self.assertIn("2026-04-11", texts)     # date

    def test_section_divider_has_number_and_title(self):
        pres = self._make_full_deck("章节过渡验证", "zh")
        prs = _build(pres, language="zh")
        divider = prs.slides[2]
        texts = _shape_texts(divider)
        # 应有 "03" 编号和章节标题
        self.assertTrue(any("03" in t for t in texts), f"texts={texts}")
        self.assertTrue(any("背景与挑战" in t for t in texts), f"texts={texts}")

    def test_agenda_has_numbered_items(self):
        pres = self._make_full_deck("议程验证", "zh")
        prs = _build(pres, language="zh")
        agenda = prs.slides[1]
        texts = _shape_texts(agenda)
        for label in ("01", "02", "03"):
            self.assertTrue(any(t == label for t in texts),
                            f"未找到编号 {label} in {texts}")

    def test_content_has_takeaway_and_bullets(self):
        pres = self._make_full_deck("内容验证", "zh")
        prs = _build(pres, language="zh")
        content = prs.slides[3]
        texts = _shape_texts(content)
        self.assertTrue(any("线上业务同比增长67%" in t for t in texts))
        self.assertTrue(any("■" in t for t in texts), f"未找到 L0 bullet 前缀")

    def test_chart_dedup_by_id(self):
        from models import ChartSpec
        same_id_chart = ChartSpec(
            chart_id="dup-1",
            chart_type=ChartType.COLUMN,
            categories=["A", "B"],
            series=[ChartSeries(name="x", values=[1, 2])],
        )
        pres = PresentationSpec(title="去重测试", language="zh")
        pres.slides.append(SlideSpec(
            slide_type=SlideType.DATA,
            takeaway_message="同一 chart_id 的两个 ChartSpec 不应渲染两次",
            charts=[same_id_chart, same_id_chart],
            language="zh",
        ))
        prs = _build(pres, language="zh")
        chart_count = sum(1 for sh in prs.slides[0].shapes if sh.has_chart)
        self.assertEqual(chart_count, 1)

    def test_empty_body_fallback(self):
        pres = PresentationSpec(title="空白页兜底", language="zh")
        pres.slides.append(SlideSpec(
            slide_type=SlideType.CONTENT,
            takeaway_message="本页正文为空，应被兜底注入",
            language="zh",
        ))
        prs = _build(pres, language="zh")
        slide = prs.slides[0]
        texts = " ".join(_shape_texts(slide))
        self.assertIn("本页正文为空", texts)

    def test_language_font_strategy_zh_vs_en(self):
        zh_strat = LANGUAGE_FONT_STRATEGY["zh"]
        en_strat = LANGUAGE_FONT_STRATEGY["en"]
        # 中文 body 应不小于英文 body
        self.assertGreaterEqual(zh_strat["font_sizes"]["body"],
                                en_strat["font_sizes"]["body"])
        # 中文 title 不小于英文 title
        self.assertGreaterEqual(zh_strat["font_sizes"]["title"],
                                en_strat["font_sizes"]["title"])
        # 中文应使用 YaHei
        self.assertEqual(zh_strat["fonts"]["body"], "Microsoft YaHei")

    def test_pie_chart_does_not_crash_on_value_axis(self):
        """Regression: PIE 图无 value_axis，旧代码会抛 ValueError。"""
        pres = PresentationSpec(title="饼图回归", language="zh")
        pres.slides.append(SlideSpec(
            slide_type=SlideType.DATA,
            takeaway_message="市场份额分布",
            charts=[ChartSpec(
                chart_type=ChartType.PIE,
                title="份额",
                categories=["A", "B", "C"],
                series=[ChartSeries(name="share", values=[40, 35, 25])],
            )],
            language="zh",
        ))
        # 不应抛 ValueError("chart has no value axis")
        prs = _build(pres, language="zh")
        self.assertEqual(len(prs.slides), 1)
        chart_count = sum(1 for sh in prs.slides[0].shapes if sh.has_chart)
        self.assertEqual(chart_count, 1)

    def test_all_chart_types_render(self):
        """烟测：每种 ChartType 都能不报错地渲染（COMBO 需 ≥2 系列）"""
        for ct in [ChartType.BAR, ChartType.COLUMN, ChartType.LINE,
                   ChartType.PIE, ChartType.AREA, ChartType.WATERFALL]:
            with self.subTest(chart_type=ct):
                pres = PresentationSpec(title=f"chart_{ct.value}", language="zh")
                pres.slides.append(SlideSpec(
                    slide_type=SlideType.DATA,
                    takeaway_message=f"{ct.value} 测试",
                    charts=[ChartSpec(
                        chart_type=ct,
                        title="x",
                        categories=["A", "B", "C"],
                        series=[ChartSeries(name="v", values=[10, 20, 15])],
                    )],
                    language="zh",
                ))
                prs = _build(pres, language="zh")
                self.assertEqual(len(prs.slides), 1)

    def test_docx_with_image_does_not_crash(self):
        """Regression: 含嵌入图片的 docx 必须能成功解析"""
        import struct, zlib, tempfile
        from docx import Document
        from docx.shared import Inches
        from pipeline.layer1_input.docx_parser import DocxParser

        def _png():
            sig = b"\x89PNG\r\n\x1a\n"
            def chunk(t, d):
                return struct.pack(">I", len(d)) + t + d + struct.pack(
                    ">I", zlib.crc32(t + d) & 0xFFFFFFFF
                )
            ihdr = struct.pack(">IIBBBBB", 10, 10, 8, 2, 0, 0, 0)
            raw = b"".join(b"\x00" + b"\xff\x00\x00" * 10 for _ in range(10))
            return (sig + chunk(b"IHDR", ihdr)
                    + chunk(b"IDAT", zlib.compress(raw))
                    + chunk(b"IEND", b""))

        with tempfile.TemporaryDirectory() as td:
            img_path = Path(td) / "i.png"
            img_path.write_bytes(_png())
            doc = Document()
            doc.add_heading("测试", 0)
            doc.add_paragraph("段落1")
            doc.add_picture(str(img_path), width=Inches(1))
            doc.add_paragraph("段落2")
            doc_path = Path(td) / "t.docx"
            doc.save(str(doc_path))

            rc = DocxParser().parse(str(doc_path))
            self.assertEqual(rc.source_type, "doc")
            self.assertIn("段落1", rc.raw_text)
            self.assertGreaterEqual(len(rc.images), 1)

    def test_source_pictures_render_on_content_slide(self):
        """Regression: SlideSpec.pictures 中的原材料图必须渲染到 slide 上。"""
        import struct, zlib, tempfile
        def _png():
            sig = b"\x89PNG\r\n\x1a\n"
            def chunk(t, d):
                return struct.pack(">I", len(d)) + t + d + struct.pack(
                    ">I", zlib.crc32(t + d) & 0xFFFFFFFF)
            ihdr = struct.pack(">IIBBBBB", 20, 20, 8, 2, 0, 0, 0)
            raw = b"".join(b"\x00" + b"\x33\x88\xEE" * 20 for _ in range(20))
            return (sig + chunk(b"IHDR", ihdr)
                    + chunk(b"IDAT", zlib.compress(raw))
                    + chunk(b"IEND", b""))

        with tempfile.TemporaryDirectory() as td:
            pic = Path(td) / "pic.png"
            pic.write_bytes(_png())
            pres = PresentationSpec(title="图片渲染回归", language="zh")
            pres.slides.append(SlideSpec(
                slide_type=SlideType.CONTENT,
                takeaway_message="右侧附图左侧正文",
                narrative_arc=NarrativeRole.EVIDENCE,
                text_blocks=[
                    TextBlock(content="主论据", level=0),
                    TextBlock(content="证据一", level=1),
                    TextBlock(content="证据二", level=1),
                ],
                pictures=[str(pic)],
                language="zh",
            ))
            prs = _build(pres, language="zh")
            slide = prs.slides[0]
            pic_shapes = [sh for sh in slide.shapes if sh.shape_type == 13]
            self.assertEqual(len(pic_shapes), 1,
                             f"应该渲染 1 张原材料图片，实际 {len(pic_shapes)}")
            # 图片应位于右侧 (left 超过 slide_width 的一半)
            self.assertGreater(pic_shapes[0].left, prs.slide_width // 2)

    def test_content_slide_never_drops_body_when_layout_has_no_body(self):
        """Regression: 若 layout 未提供 body_area 但 spec 有 text_blocks，应兜底。"""
        from pipeline.layer6_output.ppt_builder import PPTBuilder
        from models import LayoutCoordinates, Rect
        b = PPTBuilder.__new__(PPTBuilder)
        b.prs = Presentation()  # 需要 slide_width/height
        layout = LayoutCoordinates(
            title_area=Rect(left=457200, top=457200, width=11277600, height=457200),
            source_area=Rect(left=457200, top=6400000, width=11277600, height=182880),
        )
        spec = SlideSpec(
            slide_type=SlideType.CONTENT,
            takeaway_message="t",
            text_blocks=[TextBlock(content="x", level=0)],
        )
        rect = b._emergency_body_rect(spec, layout)
        self.assertGreater(rect.height, 0)
        self.assertGreater(rect.width, 0)
        # 应位于 title 下方、source 上方
        self.assertGreaterEqual(rect.top, layout.title_area.top + layout.title_area.height)
        self.assertLessEqual(rect.top + rect.height, layout.source_area.top)

    def test_outline_extract_json_handles_llm_quirks(self):
        """Regression: LLM 围栏/尾随逗号/行注释/缺失围栏都要能解析"""
        import re

        def _parse_json_from_text(text):
            patterns = [
                r'```json\s*(\[[\s\S]*?\])\s*```',
                r'```\s*(\[[\s\S]*?\])\s*```',
                r'(\[[\s\S]*"slide_type"[\s\S]*?\])',
            ]
            def _clean(s):
                s = re.sub(r'//[^\n]*', '', s)
                s = re.sub(r',(\s*[\]\}])', r'\1', s)
                return s
            for pattern in patterns:
                for match in re.finditer(pattern, text):
                    raw = match.group(1)
                    for candidate in (raw, _clean(raw)):
                        try:
                            data = json.loads(candidate)
                            if isinstance(data, list) and data and isinstance(data[0], dict):
                                if "slide_type" in data[0] or "page_number" in data[0]:
                                    return data
                        except Exception:
                            continue
            return []

        slide = '{"page_number":1,"slide_type":"title","takeaway":"封面"}'

        cases = {
            "fenced_array": (
                f'```json\n[{slide}]\n```'
            ),
            "fenced_trailing_comma": (
                f'```json\n[{slide},]\n```'
            ),
            "missing_closing_fence": (
                f'```json\n[{slide}]'
            ),
            "inline_comment": (
                f'```json\n[\n  // 注释\n  {slide}\n]\n```'
            ),
            "with_prose": (
                f'这是大纲：\n```json\n[{slide}]\n```\n请审核'
            ),
        }
        for name, raw in cases.items():
            with self.subTest(case=name):
                parsed = _parse_json_from_text(raw)
                self.assertIsNotNone(parsed, f"{name} 解析失败")
                self.assertGreater(len(parsed), 0, f"{name} 结果为空")
                self.assertIn("page_number", parsed[0])

    def test_section_divider_routes_to_title_template(self):
        from pipeline.layer4_visual.visual_designer import (
            SLIDE_TYPE_LAYOUT_OVERRIDE,
        )
        self.assertEqual(SLIDE_TYPE_LAYOUT_OVERRIDE[SlideType.SECTION_DIVIDER],
                         "title_center")
        self.assertEqual(SLIDE_TYPE_LAYOUT_OVERRIDE[SlideType.TITLE],
                         "title_center")
        self.assertEqual(SLIDE_TYPE_LAYOUT_OVERRIDE[SlideType.AGENDA],
                         "agenda_vertical")


if __name__ == "__main__":
    unittest.main()
