"""Tests for R29 — RawContent headings + structured_blocks."""
import pytest
import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))


class TestHeadingDataclass:
    def test_heading_creation(self):
        from models.slide_spec import Heading
        h = Heading(level=1, text="第一章 概述")
        assert h.level == 1
        assert h.text == "第一章 概述"
        assert h.page_idx == 0
        assert h.char_offset == 0

    def test_heading_defaults(self):
        from models.slide_spec import Heading
        h = Heading(level=2, text="1.1 子节", page_idx=3)
        assert h.page_idx == 3


class TestStructuredBlockDataclass:
    def test_block_paragraph(self):
        from models.slide_spec import StructuredBlock
        b = StructuredBlock(type="paragraph", text="这是一段正文")
        assert b.type == "paragraph"
        assert b.table_idx == -1
        assert b.image_idx == -1
        assert b.heading_path == []

    def test_block_table(self):
        from models.slide_spec import StructuredBlock
        b = StructuredBlock(type="table", text="表格: A | B", table_idx=2)
        assert b.table_idx == 2

    def test_block_heading_path(self):
        from models.slide_spec import StructuredBlock
        b = StructuredBlock(type="paragraph", text="text", heading_path=["第一章", "1.2 小节"])
        assert b.heading_path == ["第一章", "1.2 小节"]


class TestRawContentNewFields:
    def test_raw_content_has_headings(self):
        from models.slide_spec import RawContent
        rc = RawContent(source_type="doc")
        assert rc.headings == []
        assert rc.structured_blocks == []

    def test_raw_content_with_headings(self):
        from models.slide_spec import RawContent, Heading, StructuredBlock
        h = [Heading(level=1, text="Ch1"), Heading(level=2, text="1.1")]
        b = [StructuredBlock(type="heading", text="Ch1")]
        rc = RawContent(source_type="doc", headings=h, structured_blocks=b)
        assert len(rc.headings) == 2
        assert rc.headings[0].level == 1
        assert len(rc.structured_blocks) == 1


class TestDocxParserFlatStructure:
    """Test that docx_parser extracts headings and structured_blocks."""

    def test_docx_with_headings(self, tmp_path):
        """Create a simple .docx with Heading 1/2 styles and verify extraction."""
        pytest.importorskip("docx")
        from docx import Document
        from docx.shared import Pt
        from pipeline.layer1_input.docx_parser import DocxParser

        doc = Document()
        doc.add_heading("第一章 概述", level=1)
        doc.add_paragraph("这是第一章的正文内容，包含一些描述性文字。")
        doc.add_heading("1.1 子节", level=2)
        doc.add_paragraph("子节的正文内容。")
        doc.add_heading("第二章 分析", level=1)
        doc.add_paragraph("第二章正文。")

        path = tmp_path / "test.docx"
        doc.save(str(path))

        parser = DocxParser()
        rc = parser.parse(str(path))

        assert len(rc.headings) >= 3
        heading_texts = [h.text for h in rc.headings]
        assert any("概述" in t for t in heading_texts)
        assert any("分析" in t for t in heading_texts)

        # Check heading levels
        h1_count = sum(1 for h in rc.headings if h.level == 1)
        assert h1_count >= 2

    def test_docx_structured_blocks(self, tmp_path):
        """Verify structured_blocks contain heading + paragraph types."""
        pytest.importorskip("docx")
        from docx import Document
        from pipeline.layer1_input.docx_parser import DocxParser

        doc = Document()
        doc.add_heading("第一章", level=1)
        doc.add_paragraph("段落A。")
        doc.add_paragraph("段落B。")

        path = tmp_path / "test_blocks.docx"
        doc.save(str(path))

        parser = DocxParser()
        rc = parser.parse(str(path))

        block_types = [b.type for b in rc.structured_blocks]
        assert "heading" in block_types
        assert "paragraph" in block_types

    def test_docx_heading_path_propagation(self, tmp_path):
        """Verify heading_path is set on paragraph blocks."""
        pytest.importorskip("docx")
        from docx import Document
        from pipeline.layer1_input.docx_parser import DocxParser

        doc = Document()
        doc.add_heading("第一章", level=1)
        doc.add_paragraph("正文段落。")

        path = tmp_path / "test_path.docx"
        doc.save(str(path))

        parser = DocxParser()
        rc = parser.parse(str(path))

        para_blocks = [b for b in rc.structured_blocks if b.type == "paragraph"]
        assert len(para_blocks) >= 1
        assert any("第一章" in b.heading_path for b in para_blocks)


class TestMarkdownParserFlatStructure:
    def test_markdown_headings(self, tmp_path):
        from pipeline.layer1_input.markdown_parser import MarkdownParser

        md_content = "# 第一章 概述\n\n一些正文。\n\n## 1.1 子节\n\n子节内容。\n\n# 第二章 分析\n\n第二章正文。"
        path = tmp_path / "test.md"
        path.write_text(md_content, encoding="utf-8")

        parser = MarkdownParser()
        rc = parser.parse(str(path))

        assert len(rc.headings) == 3
        assert rc.headings[0].level == 1
        assert rc.headings[1].level == 2
        assert rc.headings[2].level == 1

    def test_markdown_structured_blocks(self, tmp_path):
        from pipeline.layer1_input.markdown_parser import MarkdownParser

        md_content = "# 标题\n\n段落文字。\n\n- 列表项1\n- 列表项2"
        path = tmp_path / "test_blocks.md"
        path.write_text(md_content, encoding="utf-8")

        parser = MarkdownParser()
        rc = parser.parse(str(path))

        block_types = [b.type for b in rc.structured_blocks]
        assert "heading" in block_types
        assert "paragraph" in block_types
        assert "list" in block_types


class TestPptxParserFlatStructure:
    def test_pptx_headings(self, tmp_path):
        pytest.importorskip("pptx")
        from pptx import Presentation
        from pptx.util import Inches
        from pipeline.layer1_input.pptx_parser import PptxParser

        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "第一章 概述"

        slide2 = prs.slides.add_slide(slide_layout)
        slide2.shapes.title.text = "第二章 分析"

        path = tmp_path / "test.pptx"
        prs.save(str(path))

        parser = PptxParser()
        rc = parser.parse(str(path))

        assert len(rc.headings) >= 2
        heading_texts = [h.text for h in rc.headings]
        assert "第一章 概述" in heading_texts


class TestParseAgentSerialization:
    def test_parse_agent_includes_new_fields(self, tmp_path):
        """Verify ParseAgent serializes headings and structured_blocks."""
        pytest.importorskip("docx")
        from docx import Document
        from pipeline.agents.parse_agent import ParseAgent

        doc = Document()
        doc.add_heading("测试章节", level=1)
        doc.add_paragraph("测试内容。")

        path = tmp_path / "test.docx"
        doc.save(str(path))

        agent = ParseAgent()
        result = agent.run({"task": {"file_path": str(path)}})

        assert "headings" in result
        assert "structured_blocks" in result
        assert isinstance(result["headings"], list)
        assert isinstance(result["structured_blocks"], list)
        if result["headings"]:
            assert "level" in result["headings"][0]
            assert "text" in result["headings"][0]
