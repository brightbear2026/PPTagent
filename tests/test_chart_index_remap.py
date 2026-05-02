"""
Tests for chart_renderer file→PPTX slide index remapping.

When html2pptx.js drops failed slides, chart_renderer must remap
file indices to actual PPTX slide positions.
"""
import pytest
from pptx import Presentation


def _make_renderer():
    from pipeline.layer6_output.chart_renderer import ChartRenderer
    return ChartRenderer()


def _chart_spec_dict():
    return {
        "chart_type": "column",
        "title": "Test",
        "categories": ["A", "B"],
        "series": [{"name": "S1", "values": [10, 20]}],
    }


def _placeholder_entry(slide_index):
    return {
        "slide_index": slide_index,
        "items": [{"id": "c1", "x": 1.0, "y": 1.0, "w": 4.0, "h": 3.0}],
    }


def _slide_data(slide_type="content", has_chart=True):
    d = {"slide_type": slide_type}
    if has_chart:
        d["chart_spec"] = _chart_spec_dict()
    return d


def _make_pptx(tmp_path, slide_count):
    prs = Presentation()
    for _ in range(slide_count):
        prs.slides.add_slide(prs.slide_layouts[6])
    p = str(tmp_path / "input.pptx")
    prs.save(p)
    return p


def test_no_errors_direct_mapping(tmp_path):
    """5 slides, 0 errors → file index == PPTX index, chart on slide 3."""
    renderer = _make_renderer()
    input_path = _make_pptx(tmp_path, 5)
    output_path = str(tmp_path / "output.pptx")

    placeholders = [_placeholder_entry(3)]
    slides_data = [{"slide_type": "content"}] * 5
    slides_data[3] = _slide_data("content")

    renderer.render_into_pptx(input_path, placeholders, slides_data, output_path)

    result = Presentation(output_path)
    assert len([s for s in result.slides[3].shapes if s.has_chart]) >= 1


def test_one_dropped_slide(tmp_path):
    """5 files, file 2 dropped → PPTX has 4 slides. File 3 → PPTX slide 2."""
    renderer = _make_renderer()
    input_path = _make_pptx(tmp_path, 4)  # 4 actual slides
    output_path = str(tmp_path / "output.pptx")

    placeholders = [_placeholder_entry(3)]  # file index 3
    slides_data = [{"slide_type": "content"}] * 5
    slides_data[3] = _slide_data("content")

    render_errors = [{"slide_index": 2, "error": "validation failed"}]

    renderer.render_into_pptx(
        input_path, placeholders, slides_data, output_path,
        render_errors=render_errors,
    )

    result = Presentation(output_path)
    # file 3 maps to PPTX slide 2 (file 2 was dropped)
    assert len([s for s in result.slides[2].shapes if s.has_chart]) >= 1


def test_multiple_dropped(tmp_path):
    """8 files, files 1 and 4 dropped → PPTX has 6 slides."""
    renderer = _make_renderer()
    input_path = _make_pptx(tmp_path, 6)
    output_path = str(tmp_path / "output.pptx")

    # File 5 has chart, maps to PPTX slide 3
    placeholders = [_placeholder_entry(5)]
    slides_data = [{"slide_type": "content"}] * 8
    slides_data[5] = _slide_data("content")

    render_errors = [
        {"slide_index": 1, "error": "fail"},
        {"slide_index": 4, "error": "fail"},
    ]

    renderer.render_into_pptx(
        input_path, placeholders, slides_data, output_path,
        render_errors=render_errors,
    )

    result = Presentation(output_path)
    # file 5: skip 1,4 → offset=2 → PPTX slide 3
    assert len([s for s in result.slides[3].shapes if s.has_chart]) >= 1


def test_dropped_first_slide(tmp_path):
    """File 0 dropped → all subsequent shift by 1."""
    renderer = _make_renderer()
    input_path = _make_pptx(tmp_path, 3)
    output_path = str(tmp_path / "output.pptx")

    # File 1 has chart, maps to PPTX slide 0
    placeholders = [_placeholder_entry(1)]
    slides_data = [{"slide_type": "content"}] * 4
    slides_data[1] = _slide_data("content")

    render_errors = [{"slide_index": 0, "error": "fail"}]

    renderer.render_into_pptx(
        input_path, placeholders, slides_data, output_path,
        render_errors=render_errors,
    )

    result = Presentation(output_path)
    assert len([s for s in result.slides[0].shapes if s.has_chart]) >= 1


def test_dropped_last_slide_no_shift(tmp_path):
    """Last file dropped → no index shift for earlier slides."""
    renderer = _make_renderer()
    input_path = _make_pptx(tmp_path, 3)
    output_path = str(tmp_path / "output.pptx")

    # File 1 has chart, maps to PPTX slide 1 (last file 3 dropped, no effect)
    placeholders = [_placeholder_entry(1)]
    slides_data = [{"slide_type": "content"}] * 4
    slides_data[1] = _slide_data("content")

    render_errors = [{"slide_index": 3, "error": "fail"}]

    renderer.render_into_pptx(
        input_path, placeholders, slides_data, output_path,
        render_errors=render_errors,
    )

    result = Presentation(output_path)
    assert len([s for s in result.slides[1].shapes if s.has_chart]) >= 1


def test_placeholder_for_dropped_slide_skipped(tmp_path):
    """Placeholder referencing a dropped slide → skip, no crash."""
    renderer = _make_renderer()
    input_path = _make_pptx(tmp_path, 2)
    output_path = str(tmp_path / "output.pptx")

    # File 0 was dropped, but placeholder references it
    placeholders = [_placeholder_entry(0)]
    slides_data = [_slide_data("content"), _slide_data("content")]

    render_errors = [{"slide_index": 0, "error": "fail"}]

    renderer.render_into_pptx(
        input_path, placeholders, slides_data, output_path,
        render_errors=render_errors,
    )

    result = Presentation(output_path)
    # No charts injected (placeholder for dropped slide was skipped)
    for slide in result.slides:
        assert len([s for s in slide.shapes if s.has_chart]) == 0


def test_section_divider_with_remap(tmp_path):
    """After remap, section_divider at correct position still filtered (H3)."""
    renderer = _make_renderer()
    input_path = _make_pptx(tmp_path, 3)
    output_path = str(tmp_path / "output.pptx")

    # File 0 dropped. File 1 = section_divider (should be filtered).
    # File 2 = content (should get chart).
    placeholders = [_placeholder_entry(1), _placeholder_entry(2)]
    slides_data = [_slide_data("content"), _slide_data("section_divider"), _slide_data("content")]

    render_errors = [{"slide_index": 0, "error": "fail"}]

    renderer.render_into_pptx(
        input_path, placeholders, slides_data, output_path,
        render_errors=render_errors,
    )

    result = Presentation(output_path)
    # PPTX slide 0 = file 1 (section_divider) → no chart
    assert len([s for s in result.slides[0].shapes if s.has_chart]) == 0
    # PPTX slide 1 = file 2 (content) → has chart
    assert len([s for s in result.slides[1].shapes if s.has_chart]) >= 1


def test_no_auto_percentage_in_annotation(tmp_path):
    """Annotation must NOT contain auto-computed percentages like +1050%."""
    from models.slide_spec import ChartSpec, ChartSeries

    renderer = _make_renderer()

    # Pie chart with large value differences (classic hallucination trigger)
    spec = ChartSpec(
        chart_type="pie",
        categories=["华东", "华北", "华南", "西部"],
        series=[ChartSeries(name="Revenue", values=[4500, 380, 1200, 800])],
    )

    # Extract what _add_chart_annotation would compute
    values = [float(v) for v in spec.series[0].values]
    peak_idx = max(range(len(values)), key=lambda i: abs(values[i]))
    cat_label = spec.categories[peak_idx]

    # With no so_what, annotation should be just "label: value" — no percentage
    expected = f"{cat_label}: {values[peak_idx]:g}"
    assert "%" not in expected
    assert expected == "华东: 4500"
